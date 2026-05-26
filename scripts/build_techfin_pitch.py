#!/usr/bin/env python3
"""生成 科技金融 Pitch Deck PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Color scheme: PSBC green + dark blue
DARK_BLUE = RGBColor(0x0B, 0x2A, 0x3F)
PSBC_GREEN = RGBColor(0x00, 0x7A, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_GOLD = RGBColor(0xD4, 0x9B, 0x2C)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MID_TEXT = RGBColor(0x55, 0x55, 0x55)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xBB)
        p2.font.name = 'Arial'
    # Green accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.6), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = PSBC_GREEN
    line.line.fill.background()
    return slide

def add_section_header(section_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{section_num}"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PSBC_GREEN
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.2), Inches(1.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = PSBC_GREEN
    line.line.fill.background()
    return slide

def add_content_slide(title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Green underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = PSBC_GREEN
    line.line.fill.background()
    # Content bullets
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(5.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        # Check for sub-bullets (indented)
        if bullet.startswith("  -"):
            p.text = bullet.strip()
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = MID_TEXT
            p.space_after = Pt(4)
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.font.bold = bullet.startswith("★")
            p.space_after = Pt(10)
        p.font.name = 'Arial'
    # Footnote
    if note:
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(10)
        p3.font.color.rgb = MID_TEXT
        p3.font.italic = True
    return slide

def add_two_col_slide(title, left_title, left_bullets, right_title, right_bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(1.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = PSBC_GREEN
    line.line.fill.background()
    # Left column
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6), Inches(5.3))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.fill.background()
    txL = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.5))
    pL = txL.text_frame.paragraphs[0]
    pL.text = left_title
    pL.font.size = Pt(20)
    pL.font.bold = True
    pL.font.color.rgb = DARK_BLUE
    txL2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.4), Inches(4.2))
    tfL = txL2.text_frame
    tfL.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tfL.paragraphs[0] if i == 0 else tfL.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    # Right column
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(6), Inches(5.3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GRAY
    right_box.line.fill.background()
    txR = slide.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.5))
    pR = txR.text_frame.paragraphs[0]
    pR.text = right_title
    pR.font.size = Pt(20)
    pR.font.bold = True
    pR.font.color.rgb = DARK_BLUE
    txR2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.4), Inches(4.2))
    tfR = txR2.text_frame
    tfR.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    if note:
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(10)
        p3.font.color.rgb = MID_TEXT
        p3.font.italic = True
    return slide

# ===== SLIDES =====

# Slide 1: Title
add_title_slide(
    "科技金融方法论与战略定位",
    "邮储银行上海分行 科技金融事业部 | 陈颖芳 | 2026年5月"
)

# Slide 2: Agenda
add_content_slide("议程", [
    "1. 科技金融的宏观背景与市场机遇",
    "2. 邮储银行科技金融的战略定位",
    "3. 三原则 + 方法论框架",
    "4. 并购贷款 — 科创金融生态的枢纽",
    "5. 政策建议：从五个工具到一个生态",
    "6. 竞品对标与差异化优势",
    "7. 实施路线图与关键里程碑",
    "8. 下一步行动",
])

# Slide 3: Macro Context
add_content_slide("宏观背景：为什么现在是科技金融的拐点", [
    "★ 国家战略驱动：十五五规划明确\"加快建设金融强国\"，中央经济工作会议提\"科技创新和产业创新深度融合\"",
    "★ 制度框架落位：2025年并购贷款新规+科技企业试点政策，首次允许参股型并购贷款（≥20%即可）",
    "★ 产业端爆炸式增长：全国AI核心产业规模超1.2万亿元，2030年预计破10万亿",
    "★ 银行转型压力：传统对公信贷利差收窄，科技金融是唯一既能做大规模又能做高收益的增长极",
    "★ 竞争窗口期：18城试点+并购贷款新规，先发银行已在跑马圈地，后发者的制度红利窗口有限",
    "",
    "→ 科技金融不是\"要不要做\"，而是\"谁能先做出来\"",
], "Sources: 国家金融监督管理总局, 发改委, 国务院, 2025-2026")

# Slide 4: Market Opportunity
add_two_col_slide(
    "市场机遇：科技金融的四大增量空间",
    "需求侧",
    [
        "• 科创并购融资：百亿级大型并购需要银团，传统银行服务缺位",
        "• 知识产权质押融资：AI算法/芯片专利估值缺乏标准，融资缺口巨大",
        "• 科技中小企业信贷：轻资产、高研发投入、无传统抵押物",
        "• 科技保险+担保联动：银行+保险+政府三方风险分担模式尚未普及",
    ],
    "供给侧",
    [
        "• 邮储银行科技金融事业部：全行唯一专业化科技金融团队（上海）",
        "• 差异化优势：CISA/CIA双证+20年全链路经验+IT风控专长",
        "• 民营经济占比持续提升，传统国企信贷需求增速放缓",
        "• 银行间竞争从\"价格战\"转向\"专业能力战\"",
    ],
    "Sources: 国家金融监督管理总局, 邮储银行内部数据"
)

# Slide 5: Strategy
add_content_slide("战略定位：邮储科技金融的\"三原则\"", [
    "★ 原则一：技术信用替代抵押信用",
    "  传统银行风控看房产+现金流，科技金融风控看技术成熟度(TRL)+知识产权价值+研发团队稳定性",
    "",
    "★ 原则二：生态枢纽而非资金通道",
    "  不只做贷款发放，要做科创金融生态的\"路由器\"——连接VC退出、产业并购、技术交易、ABS流转",
    "",
    "★ 原则三：专业判断替代规模竞争",
    "  不与四大行拼资金成本，拼技术尽调能力、行业认知深度、审批效率",
    "",
    "→ 定位：做\"最懂技术的银行\"，以并购贷款为尖刀产品切入科创金融生态",
])

# Slide 6: Methodology
add_content_slide("方法论框架：从技术评估到风险定价", [
    "★ STAGE 1 — 技术成熟度分层 (TRL 1-9)",
    "  TRL 1-3 实验室阶段 → 专家评议法",
    "  TRL 4-6 小批量产阶段 → 市场可比交易法",
    "  TRL 7-9 规模商业化阶段 → 收益法（DCF）",
    "",
    "★ STAGE 2 — 双签技术评估",
    "  技术专家 → 对技术可行性和先进性负责",
    "  注册资产评估师 → 对价值合理性和可比性负责",
    "",
    "★ STAGE 3 — 风险定价引擎",
    "  输入：技术评级 + 团队评级 + 市场赛道评级 + 知识产权强度",
    "  输出：贷款比例、利率、期限、增信要求",
    "",
    "★ STAGE 4 — 贷后技术监控",
    "  季度技术里程碑审查（研发进度 vs 计划）、知识产权状态变更、核心团队稳定性",
])

# Slide 7: M&A Hub
add_content_slide("并购贷款：科创金融生态的枢纽型基础设施", [
    "★ 生态位定义",
    "  并购贷款连接三个生态圈：VC退出通道 ↔ 科技企业规模化 ↔ 银行技术信用转型",
    "",
    "★ 2025年制度突破",
    "  • 首次允许参股型并购贷款（≥20%即可申请）",
    "  • 控制型贷款比例 60%→70%，期限 7年→10年",
    "  • 科技企业试点：比例80%、期限10年（18城市）",
    "",
    "★ 四大配套短板 → 五个政策建议（详见下页）",
    "  技术价值评估空白 | 担保机制缺位 | 风险分担缺失 | 银团协同不足",
    "",
    "★ 邮储先发优势：已作为18城试点范围银行，率先积累参股型并购贷款案例",
], "Sources: 《商业银行并购贷款管理办法》(金规〔2025〕27号), 2025年12月31日")

# Slide 8: Policy Recommendations
add_content_slide("五条政策建议：构建科创金融生态的五个支柱", [
    "★ 建议一：技术价值评估国家标准 → 生态的\"通用语言\"",
    "  国家知识产权局牵头，TRL分层+双签制度+白名单+银行风控对接",
    "",
    "★ 建议二：科创并购融资增信基金 → 生态的\"风险消化器\"",
    "  财政部+科技部出资，担保基金(20%)+科技保险(30%)+银行(50%)三重分担",
    "",
    "★ 建议三：银团协调机制 → 生态的\"协同网络\"",
    "  标准化合约+信息共享平台+牵头行MPA激励+二级市场转让",
    "",
    "★ 建议四：科创并购贷款ABS试点 → 生态的\"代谢循环\"",
    "  循环购买型ABS，发放→证券化→回收→再发放",
    "",
    "★ 建议五：试点→评估→立法路径 → 生态的\"进化机制\"",
    "  2025-2027试点期→2027-2028评估期→2028纳入《办法》正式修订",
])

# Slide 9: Competitive analysis
add_two_col_slide(
    "竞品对标：邮储vs四大行 科技金融差异化定位",
    "四大行",
    [
        "工商银行：全国布局，客户基数大，但审批链条长，\"总行一刀切\"难以适配科创灵活性",
        "建设银行：科技子公司+建信金科技术积累，但事业部制改革尚未完成",
        "中国银行：跨境并购经验丰富，但境内科创覆盖不足",
        "农业银行：三农+普惠专注，科技金融非核心战略",
    ],
    "邮储差异化优势",
    [
        "上海分行科技金融事业部：全行唯一专业化团队，决策链条短",
        "双证团队(CISA/CIA)：IT审计+风控双背景，看懂技术的能力四大行没有",
        "并购贷款试点先发：18城试点+参股型贷款经验积累",
        "民建+监管沟通管道：政策建议直达决策层",
        "→ 不拼规模拼专业，不拼价格拼认知",
    ],
)

# Slide 10: Roadmap
add_content_slide("实施路线图：2026-2028三阶段推进", [
    "★ Phase 1 — 夯实基础（2026 H2）",
    "  • 完成中级经济师(金融)职称考试",
    "  • 科技金融方法论V2.0：加入并购贷款生态框架",
    "  • 积累3-5笔参股型并购贷款案例",
    "  • 向民建提交社情民意稿件（并购贷款5条建议）",
    "",
    "★ Phase 2 — 扩大影响（2027）",
    "  • 在18城试点内形成可复制的\"上海模式\"",
    "  • 推动技术价值评估国标立项",
    "  • 发起首单科创并购贷款ABS",
    "  • CFA一级",
    "",
    "★ Phase 3 — 生态成型（2028）",
    "  • 试点→评估→立法：成熟经验纳入《办法》正式修订",
    "  • 科技金融方法论V3.0：全行业推广",
    "  • 竞聘支行行长：\"做最懂科技金融的行长\"",
])

# Slide 11: KPIs
add_two_col_slide(
    "关键指标与里程碑",
    "定量指标",
    [
        "• 科创并购贷款累计发放额：目标 ≥5亿元 (Phase 1)",
        "• 参股型并购贷款占比：目标 ≥30%",
        "• 技术评估双签覆盖率：100%",
        "• 知识产权质押作为主要担保方式的比例：从0→20%",
        "• 不良率控制：<2%（优于行业平均）",
    ],
    "定性里程碑",
    [
        "• 科技金融方法论V2.0完成（2026 Q3）",
        "• 民建社情民意稿件提交（2026 Q2）",
        "• 技术价值评估国标起草小组成立（2026 Q4）",
        "• 增信基金上海试点方案报批（2027 Q1）",
        "• 首单科创并购贷款ABS发行（2027 Q3）",
    ],
)

# Slide 12: Next Steps
add_content_slide("下一步行动", [
    "★ 短期（1-3个月）",
    "  • 提交民建社情民意稿件 → 争取纳入年度课题成果",
    "  • 完成中级经济师备考 → 7月考试",
    "  • 内部推动：将并购贷款生态框架纳入分行科技金融展业指引",
    "",
    "★ 中期（3-12个月）",
    "  • 对接上海金融办：推动增信基金上海试点方案",
    "  • 联合律所+评估所：起草技术价值评估标准化流程SOP",
    "  • 推动首笔\"纯知识产权质押+增信基金\"科创并购贷款落地",
    "",
    "★ 长期（1-3年）",
    "  • 从\"上海模式\"走向\"长三角模式\"",
    "  • 推动至少一项政策建议被采纳（评估国标 / 增信基金 / ABS试点）",
    "  • 竞聘支行行长：以科技金融为核心竞争力",
])

# Slide 13: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "构建科创金融生态"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "从\"做业务\"到\"建生态\""
p2.font.size = Pt(24)
p2.font.color.rgb = PSBC_GREEN
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = ""
p3.font.size = Pt(14)
p4 = tf.add_paragraph()
p4.text = "陈颖芳 | 邮储银行上海分行科技金融事业部 | 2026年5月"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0xAA, 0xCC, 0xBB)
p4.alignment = PP_ALIGN.CENTER
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4), Inches(2.3), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = PSBC_GREEN
line.line.fill.background()

# Save
output_path = '/Users/cyingfang/claude/deliverables/career/科技金融Pitch_20260527.pptx'
prs.save(output_path)
size_kb = os.path.getsize(output_path) / 1024
print(f"已生成: {output_path} ({size_kb:.1f} KB)")
