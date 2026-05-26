#!/usr/bin/env python3
"""优化 第一部分.pptx —— 保留原图 + 改进设计"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_NAVY = RGBColor(0x0B, 0x2A, 0x45)
MID_BLUE = RGBColor(0x1B, 0x5E, 0x8A)
LIGHT_BLUE = RGBColor(0xE2, 0xED, 0xF5)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
GREY_TEXT = RGBColor(0x77, 0x77, 0x77)
SOFT_BG = RGBColor(0xF8, 0xFA, 0xFC)

IMG_DIR = '/Users/cyingfang/WorkBuddy/20260429082054/pptx_images'

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_rrect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_textbox(slide, left, top, width, height, text, font='Calibri',
                size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, lsp=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = font; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(0)
    if lsp: p.line_spacing = Pt(size * lsp)
    return tb

def make_tf(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def add_p(tf, text, font='Calibri', size=14, bold=False, color=DARK_TEXT,
          align=PP_ALIGN.LEFT, sa=0, sb=0, lsp=None, first=False):
    if first and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.name = font; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(sa); p.space_before = Pt(sb)
    if lsp: p.line_spacing = Pt(size * lsp)
    return p

def add_img(slide, name, left, top, width, height):
    path = os.path.join(IMG_DIR, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False

def slide_header(slide, label, title):
    add_rect(slide, Cm(0), Cm(0), W, Cm(0.12), MID_BLUE)
    add_textbox(slide, Cm(2), Cm(0.8), Cm(10), Cm(1.2), label,
                size=10, bold=True, color=ACCENT_ORANGE)
    add_textbox(slide, Cm(2), Cm(1.8), Cm(28), Cm(1.8), title,
                size=24, bold=True, color=DARK_NAVY)
    add_rect(slide, Cm(2), Cm(3.6), Cm(3.5), Cm(0.05), ACCENT_ORANGE)
    add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
    add_textbox(slide, Cm(1.5), Cm(18.65), Cm(25), Cm(0.35),
                'Network Community Detection via Neural Embeddings', size=7, color=GREY_TEXT)
    add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35),
                '', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.2), ACCENT_ORANGE)
add_rect(slide, Cm(24), Cm(-2), Cm(14), Cm(10), RGBColor(0x0E, 0x35, 0x55))
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(28), Cm(12), Cm(10), Cm(10))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x0E, 0x35, 0x55); shape.line.fill.background()

add_textbox(slide, Cm(3), Cm(3.5), Cm(22), Cm(5),
            'Network Community\nDetection via\nNeural Embeddings',
            size=36, bold=True, color=WHITE, lsp=1.2)
add_rect(slide, Cm(3), Cm(9), Cm(5), Cm(0.08), ACCENT_ORANGE)

tf = make_tf(slide, Cm(3), Cm(10), Cm(20), Cm(3))
add_p(tf, '基于 node2vec / DeepWalk / LINE 的图嵌入方法', size=14, color=RGBColor(0xBB, 0xCC, 0xDD), first=True)
add_p(tf, '在社区检测中的信息论极限证明', size=14, color=RGBColor(0xBB, 0xCC, 0xDD), sa=8)
add_p(tf, '文献报告  |  第一部分  |  陈于东  |  2026年5月', size=10, color=GREY_TEXT)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)

# ================================================================
# SLIDE 2: OUTLINE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '','')

add_textbox(slide, Cm(2), Cm(1.2), Cm(20), Cm(2),
            '报告提纲', size=28, bold=True, color=DARK_NAVY)
add_rect(slide, Cm(2), Cm(3.2), Cm(4), Cm(0.05), ACCENT_ORANGE)

items = [
    ('1', '研究背景', '网络、社区与图嵌入'),
    ('2', '图嵌入方法', 'DeepWalk / LINE / node2vec'),
    ('3', '核心问题', '神经图嵌入能否达到信息论极限？'),
    ('4', '主要结论', '证明方法与关键结果'),
    ('5', '问题设定', 'SBM / PPM 模型与混合参数'),
    ('6', '信息论极限', '可检测性条件与 μ*'),
    ('7', '数值验证', '实验支撑与讨论'),
]

for i, (num, title, desc) in enumerate(items):
    y = Cm(4.5 + i * 1.8)
    bg = LIGHT_BLUE if int(num) <= 3 else RGBColor(0xFD, 0xF2, 0xE4)
    tc = MID_BLUE if int(num) <= 3 else ACCENT_ORANGE
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2.5), y+Cm(0.1), Cm(1), Cm(1))
    c.fill.solid(); c.fill.fore_color.rgb = bg; c.line.fill.background()
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(12)
    p.font.bold = True; p.font.color.rgb = tc; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_textbox(slide, Cm(4.5), y, Cm(15), Cm(1.2), title, size=16, bold=True, color=DARK_TEXT)
    if i < len(items)-1:
        add_rect(slide, Cm(4.5), y+Cm(1.4), Cm(28), Cm(0.015), LIGHT_BLUE)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '2/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 3: 机器学习
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '基础概念', '名词解释：机器学习')
add_img(slide, 'slide03_img01.jpg', Cm(2), Cm(4.5), Cm(30), Cm(13.5))
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '3/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 4: 神经网络
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '基础概念', '神经网络')
add_img(slide, 'slide04_img02.jpg', Cm(2), Cm(4.5), Cm(30), Cm(13.5))
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '4/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 5: 网络无处不在
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '研究背景', '网络与图表示')

add_img(slide, 'slide05_img02.jpg', Cm(2), Cm(4.5), Cm(14), Cm(9))
add_img(slide, 'slide05_img03.jpg', Cm(17.5), Cm(4.5), Cm(14.5), Cm(9))

tf = make_tf(slide, Cm(2), Cm(14), Cm(30), Cm(3.5))
add_p(tf, '网络以节点和连接边的形式表示复杂系统的结构，在社会科学、交通、金融、', size=11, color=DARK_TEXT, first=True)
add_p(tf, '科学学、神经科学和生物学中无所不在。网络的复杂、高维、离散特性使其', size=11, color=DARK_TEXT)
add_p(tf, '有效表示变得高度非平凡。图嵌入将节点映射为向量，使得机器学习方法', size=11, color=DARK_TEXT)
add_p(tf, '可直接应用于图任务。', size=11, color=DARK_TEXT)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '5/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 6: 图嵌入
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '方法概述', '图嵌入（Graph Embedding）')

# Three images: top-left, top-right, bottom-right
add_img(slide, 'slide06_img03.jpg', Cm(2), Cm(4.5), Cm(14), Cm(6))
add_img(slide, 'slide06_img02.jpg', Cm(17.5), Cm(4.5), Cm(14.5), Cm(6))
add_img(slide, 'slide06_img04.jpg', Cm(17.5), Cm(11), Cm(14.5), Cm(6))

tf = make_tf(slide, Cm(2), Cm(11), Cm(14), Cm(6))
add_p(tf, '图嵌入将每个节点表示为紧凑连续向量空间中的一个点，使得可以直接应用强大的机器学习方法解决各类任务。', size=11, color=DARK_TEXT, first=True)
add_p(tf, '', size=6)
add_p(tf, '然而，由于神经网络的"黑箱"性质，', size=11, color=DARK_TEXT)
add_p(tf, '这些方法的工作原理及其原因在很大', size=11, color=DARK_TEXT)
add_p(tf, '程度上仍属未知。', size=11, color=DARK_TEXT)
add_p(tf, '', size=6)
add_p(tf, '本文致力于建立清晰的理论理解。', size=11, bold=True, color=MID_BLUE)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '6/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 7: SBM
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '问题设定', '随机分块模型（SBM）')

add_img(slide, 'slide07_img03.jpg', Cm(2), Cm(4.5), Cm(30), Cm(7))

tf = make_tf(slide, Cm(2), Cm(12), Cm(18), Cm(5.5))
add_p(tf, 'SBM 是带有社区结构网络的基本生成模型，', size=11, color=DARK_TEXT, first=True)
add_p(tf, '常被用作社区检测算法的基准。', size=11, color=DARK_TEXT)
add_p(tf, '节点被分成若干"块(社区)"，同一社区内', size=11, color=DARK_TEXT)
add_p(tf, '节点连接概率高，不同社区间连接概率低。', size=11, color=DARK_TEXT)

add_rrect(slide, Cm(21), Cm(12), Cm(11), Cm(5.5), LIGHT_BLUE)
tf = make_tf(slide, Cm(22), Cm(12.3), Cm(9), Cm(5))
add_p(tf, '注', size=12, bold=True, color=MID_BLUE, first=True)
add_p(tf, '本文关注稀疏网络（平均度远小于节点数），这是实际应用中最常见也最具挑战的场景。', size=9, color=DARK_TEXT, lsp=1.3)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '7/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 8: 主要结论
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '核心结论', '主要证明结果')

add_img(slide, 'slide08_img02.jpg', Cm(2), Cm(4.5), Cm(30), Cm(7))

tf = make_tf(slide, Cm(2), Cm(12), Cm(30), Cm(5))
add_p(tf, '作者证明了：基于没有非线性激活函数的浅层神经网络图嵌入方法', size=11, color=DARK_TEXT, first=True)
add_p(tf, '（DeepWalk、LINE、node2vec）能够达到 SBM 生成图上社区检测的', size=11, color=DARK_TEXT)
add_p(tf, '信息论极限。这意味着深度学习的两大常见组件——多层和非线性激活——', size=11, color=DARK_TEXT)
add_p(tf, '并非实现最优社区检测所必需的。', size=11, color=DARK_TEXT)
add_p(tf, '', size=4)
add_p(tf, '数值实验：node2vec 嵌入 + K-means 聚类精度接近最优置信传播方法。', size=11, bold=True, color=MID_BLUE)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '8/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 9: PPM
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '问题设定', '种植分块模型（PPM）与混合参数')

add_rrect(slide, Cm(2), Cm(4.5), Cm(15), Cm(9), LIGHT_BLUE)
tf = make_tf(slide, Cm(3.5), Cm(4.8), Cm(12), Cm(8.5))
add_p(tf, 'PPM 是 SBM 的特例', size=14, bold=True, color=DARK_NAVY, first=True)
add_p(tf, '', size=4)
add_p(tf, '• 节点划分为 q 个规模相等的社区', size=11, color=DARK_TEXT)
add_p(tf, '• 同社区连接概率 = p_in', size=11, color=DARK_TEXT)
add_p(tf, '• 不同社区连接概率 = p_out', size=11, color=DARK_TEXT)
add_p(tf, '• 稀疏条件：p ∝ 1/n', size=11, color=DARK_TEXT)
add_p(tf, '', size=6)
add_p(tf, '混合参数：μ = np_out/⟨k⟩', size=13, bold=True, color=ACCENT_ORANGE)
add_p(tf, 'μ → 0：社区分离良好', size=11, color=DARK_TEXT)
add_p(tf, 'μ → 1：无社区结构(ER图)', size=11, color=DARK_TEXT)
add_p(tf, 'μ ∈ [0,1)：社区始终存在', size=11, color=DARK_TEXT)

add_img(slide, 'slide09_img02.jpg', Cm(18), Cm(4.5), Cm(14), Cm(9))

tf = make_tf(slide, Cm(2), Cm(14), Cm(30), Cm(3.5))
add_p(tf, '信息论可检测性极限：当 μ < μ* 时，算法可以优于随机猜测恢复社区；', size=11, color=DARK_TEXT, first=True)
add_p(tf, '当 μ* ≤ μ < 1 时，没有任何算法能优于随机猜测。', size=11, color=DARK_TEXT)
add_p(tf, '本文证明 node2vec 可以达到这个极限 μ*。', size=11, bold=True, color=MID_BLUE)
add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '9/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 10: 总结
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '总结', '数值验证与核心启示')

# Numerical results
add_rrect(slide, Cm(2), Cm(4.5), Cm(14.5), Cm(6), LIGHT_BLUE)
add_textbox(slide, Cm(3), Cm(4.8), Cm(12), Cm(1), '数值验证结果', size=15, bold=True, color=DARK_NAVY)
tf = make_tf(slide, Cm(3), Cm(6), Cm(12), Cm(4))
add_p(tf, '• node2vec + K-means 达到接近', size=11, color=DARK_TEXT, first=True)
add_p(tf, '  最优 BP 方法的精度', size=11, color=DARK_TEXT)
add_p(tf, '• 在度异质性下仍能学习社区', size=11, color=DARK_TEXT)
add_p(tf, '• 社区规模异质性下表现受限', size=11, color=DARK_TEXT)
add_p(tf, '  → 可能源于 K-means 的局限', size=11, color=MID_BLUE)

# Open questions
add_rrect(slide, Cm(17.5), Cm(4.5), Cm(14.5), Cm(6), SOFT_BG)
add_textbox(slide, Cm(18.5), Cm(4.8), Cm(12), Cm(1), '开放问题', size=15, bold=True, color=DARK_NAVY)
tf = make_tf(slide, Cm(18.5), Cm(6), Cm(12), Cm(4))
add_p(tf, '• 线性层是否可以被多层+', size=11, color=DARK_TEXT, first=True)
add_p(tf, '  非线性激活超越？', size=11, color=DARK_TEXT)
add_p(tf, '• 何时"深度"是必要的？', size=11, color=DARK_TEXT)
add_p(tf, '• 更复杂生成模型下', size=11, color=DARK_TEXT)
add_p(tf, '  的理论保证？', size=11, color=DARK_TEXT)

# Takeaway
add_rrect(slide, Cm(2), Cm(11.5), Cm(30), Cm(5.5), DARK_NAVY)
add_textbox(slide, Cm(3.5), Cm(11.8), Cm(27), Cm(1.5), 'Takeaway', size=16, bold=True, color=ACCENT_ORANGE)
tf = make_tf(slide, Cm(3.5), Cm(13.5), Cm(27), Cm(3))
add_p(tf, 'node2vec/DeepWalk/LINE 等浅层神经图嵌入方法', size=13, color=WHITE, lsp=1.4, first=True)
add_p(tf, '在理论上可以达到社区检测的信息论极限。', size=13, color=WHITE, lsp=1.4)
add_p(tf, '', size=4)
add_p(tf, '多一层"深"不如多一分"懂"——理解方法的理论基础，', size=12, color=RGBColor(0xCC, 0xDD, 0xEE))
add_p(tf, '才能在实践中做出最佳选择。', size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.5), SOFT_BG)
add_textbox(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '10/10', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SAVE
# ================================================================
output_path = '/Users/cyingfang/Documents/陈于东/第一部分_优化版.pptx'
prs.save(output_path)
print(f'Saved: {output_path} | Slides: {len(prs.slides)}')
