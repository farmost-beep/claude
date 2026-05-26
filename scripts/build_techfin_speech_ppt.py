#!/usr/bin/env python3
"""科技金融公开演讲 PPT —— 如果10年后回看今天"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x0D, 0x23, 0x4B)
MID_BLUE = RGBColor(0x00, 0x6D, 0xBA)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
GREY_TEXT = RGBColor(0x88, 0x99, 0xAA)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x00, 0x7A, 0x3D)

prs = Presentation()
prs.slide_width = Cm(33.867)   # 16:9 widescreen
prs.slide_height = Cm(19.05)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_name='Microsoft YaHei',
                 font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_section_number(slide, number, label):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.5), Cm(1), Cm(1.5), Cm(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MID_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(slide, Cm(3.5), Cm(1.2), Cm(15), Cm(1.5), label,
                 font_size=14, color=MID_BLUE, bold=False)


def make_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox

def add_para(tf, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
             space_after=0, space_before=0, line_spacing=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return p

def add_rich_text_box(slide, left, top, width, height, paragraphs_data, line_spacing=1.5):
    """paragraphs_data: list of dicts with keys: text, font_size, bold, color, alignment, space_after"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pd['text']
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(pd.get('font_size', 14))
        p.font.bold = pd.get('bold', False)
        p.font.color.rgb = pd.get('color', DARK_TEXT)
        p.alignment = pd.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(pd.get('space_after', 0))
        p.space_before = Pt(pd.get('space_before', 0))
        if pd.get('line_spacing'):
            p.line_spacing = Pt(pd['font_size'] * pd['line_spacing'])
        elif line_spacing:
            p.line_spacing = Pt(pd.get('font_size', 14) * line_spacing)
    return txBox


# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(26), Cm(-3), Cm(12), Cm(12))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x15, 0x2D, 0x55)
shape.line.fill.background()

add_text_box(slide, Cm(4), Cm(4), Cm(24), Cm(4),
             '如果10年后回看今天',
             font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

txBox = make_text_box(slide, Cm(4), Cm(8.5), Cm(24), Cm(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '哪些今天看起来“不合理”的决定，\n会被证明是明智的？'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(20)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.LEFT
p.line_spacing = Pt(32)

add_rect(slide, Cm(4), Cm(12), Cm(6), Cm(0.08), GOLD)
add_text_box(slide, Cm(4), Cm(13.5), Cm(20), Cm(2),
             '邮储银行上海分行科技金融事业部\n2026年5月',
             font_size=11, bold=False, color=GREY_TEXT, alignment=PP_ALIGN.LEFT)


# ================================================================
# SLIDE 2: 开场三问
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_number(slide, 1, '开场')

add_text_box(slide, Cm(1.5), Cm(3), Cm(24), Cm(2),
             '先问三个问题',
             font_size=30, bold=True, color=DARK_BLUE)
add_rect(slide, Cm(1.5), Cm(5.2), Cm(4), Cm(0.06), GOLD)

questions = [
    (1, '为什么中国最赚钱的银行，和中国最具创新力的科技企业之间，仍然隔着一堵看不见的墙？'),
    (2, '为什么我们愿意为一家还在亏损的科技企业提供数十亿授信，却很难用一句话说清楚它的核心价值？'),
    (3, '如果2036年回看今天，哪些让我们犹豫不决的决定，会被证明是理所当然的？'),
]

for i, (num, q_body) in enumerate(questions):
    y = Cm(6.5 + i * 3.5)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.8), y, Cm(1), Cm(1))
    c.fill.solid()
    c.fill.fore_color.rgb = MID_BLUE
    c.line.fill.background()
    tf = c.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(slide, Cm(3.5), y - Cm(0.1), Cm(28), Cm(3.2), q_body,
                 font_size=16, color=DARK_TEXT, line_spacing=1.5)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), LIGHT_BLUE)
add_text_box(slide, Cm(1.5), Cm(18.7), Cm(30), Cm(0.5),
             '这三个问题没有标准答案。但试图回答它们的过程，本身就是科技金融认知的进阶之路。',
             font_size=9, color=MID_BLUE)


# ================================================================
# SLIDE 3: 时间困境
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_number(slide, 2, '时间困境')

add_text_box(slide, Cm(1.5), Cm(3), Cm(30), Cm(2),
             '科技金融最大的困境：时间的错配',
             font_size=28, bold=True, color=DARK_BLUE)
add_rect(slide, Cm(1.5), Cm(5.2), Cm(4), Cm(0.06), GOLD)

# Left column
add_rect(slide, Cm(1.5), Cm(6.5), Cm(14), Cm(4.5), LIGHT_BLUE)
add_text_box(slide, Cm(2.5), Cm(6.8), Cm(12), Cm(1.5),
             '技术的逻辑', font_size=20, bold=True, color=DARK_BLUE)
txBox = make_text_box(slide, Cm(2.5), Cm(8.5), Cm(12), Cm(2.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '技术从研发到商业化的周期\n5-15年\n需要的是“耐心资本”'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = DARK_TEXT
p.line_spacing = Pt(22)

# Right column
add_rect(slide, Cm(17.5), Cm(6.5), Cm(14.5), Cm(4.5), LIGHT_BLUE)
add_text_box(slide, Cm(18.5), Cm(6.8), Cm(12), Cm(1.5),
             '金融的逻辑', font_size=20, bold=True, color=DARK_BLUE)
txBox = make_text_box(slide, Cm(18.5), Cm(8.5), Cm(12), Cm(2.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '金融产品的期限结构\n1-5年\n追求的是“确定性”'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = DARK_TEXT
p.line_spacing = Pt(22)

# Bottom insight
add_rect(slide, Cm(1.5), Cm(12), Cm(30.5), Cm(2.5), DARK_BLUE)
txBox = make_text_box(slide, Cm(2.5), Cm(12.3), Cm(28), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '“合理的决定”是顺应这个错配——用短钱做短事。\n“不合理的决定”是试图弥合这个错配——用长钱做长事。\n而10年后，后者会被证明是明智的。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(13)
p.font.color.rgb = WHITE
p.line_spacing = Pt(20)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), LIGHT_BLUE)
add_text_box(slide, Cm(1.5), Cm(18.7), Cm(30), Cm(0.5),
             '核心洞见：科技金融的本质不是消除技术的不确定性，而是用更长的期限和更精确的认知去拥抱它。',
             font_size=9, color=MID_BLUE)


# ================================================================
# SLIDE 4: 第一个不合理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

add_text_box(slide, Cm(2), Cm(2), Cm(5), Cm(3),
             '01', font_size=60, bold=True, color=GOLD)
add_text_box(slide, Cm(2), Cm(5), Cm(28), Cm(3),
             '在“寒冬”时期，逆势加大科技投资',
             font_size=28, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(8.5), Cm(5), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(2), Cm(9.5), Cm(28), Cm(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '当市场对某一技术领域极度悲观时——比如AI在2018-2022年的“AI寒冬”，或者半导体在2023年的周期低谷——大多数人选择收缩。\n\n但历史反复证明：最佳的布局时机，恰恰是“所有人都在逃离”的时候。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(24)

add_rect(slide, Cm(2), Cm(14), Cm(30), Cm(1.8), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(3), Cm(14.2), Cm(28), Cm(1.5),
             '案例：红杉中国在2023年半导体低谷期密集布局GPU赛道，2025-2026年摩尔线程/沐曦科创板上市回报超百倍。',
             font_size=11, color=GOLD)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.7), Cm(28), Cm(0.5),
             '“别人恐惧时我贪婪”——巴菲特的格言在科技金融中同样适用，只是需要更长的等待周期。',
             font_size=9, color=GREY_TEXT)


# ================================================================
# SLIDE 5: 第二个不合理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

add_text_box(slide, Cm(2), Cm(2), Cm(5), Cm(3),
             '02', font_size=60, bold=True, color=GOLD)
add_text_box(slide, Cm(2), Cm(5), Cm(28), Cm(3),
             '为尚未盈利的科技企业，提供大额授信',
             font_size=28, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(8.5), Cm(5), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(2), Cm(9.5), Cm(28), Cm(4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '传统银行信贷逻辑看的是“过去”——历史财务数据、抵押物、现金流。\n\n科技企业的价值存在“未来”——技术潜力、市场空间、团队能力。\n\n用“过去”的尺子量“未来”的价值，得到的一定是“不合理”的判断。\n\n但如果把评估维度从“历史表现”转向“未来可能”，看似不合理的授信\n就变成了精准的认知投资。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(22)

add_rect(slide, Cm(2), Cm(14.5), Cm(30), Cm(1.8), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(3), Cm(14.7), Cm(28), Cm(1.5),
             '本质：不是让银行变得更冒险，而是让银行更精确地理解技术的“不确定性分布”，从而做出更理性的决策。',
             font_size=11, color=GOLD)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.7), Cm(28), Cm(0.5),
             'TBP评估框架：从技术(Technology)、商业(Business)、人(People)三个维度重新定义“信用”。',
             font_size=9, color=GREY_TEXT)


# ================================================================
# SLIDE 6: 第三个不合理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

add_text_box(slide, Cm(2), Cm(2), Cm(5), Cm(3),
             '03', font_size=60, bold=True, color=GOLD)
add_text_box(slide, Cm(2), Cm(5), Cm(28), Cm(3),
             '放弃短期高回报，追求长期确定性',
             font_size=28, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(8.5), Cm(5), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(2), Cm(9.5), Cm(28), Cm(4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '在科技投资中，“快钱”的诱惑无处不在——追逐热点、参与炒作、短期套利。\n\n但历史证明：真正穿越周期的回报，来自那些愿意等待的人。\n\n“耐心资本”不是一种美德，而是一种策略——因为技术的演化有自己的时间节奏，金融必须尊重这个节奏，而不是试图改变它。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(22)

add_rect(slide, Cm(2), Cm(14.5), Cm(30), Cm(1.8), RGBColor(0x15, 0x2D, 0x55))
txBox = make_text_box(slide, Cm(3), Cm(14.7), Cm(28), Cm(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '《道德经》：“道法自然”——尊重创新的自然时间节奏，不拔苗助长。耐心不是被动等待，而是主动布局后的从容。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(11)
p.font.color.rgb = GOLD

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.7), Cm(28), Cm(0.5),
             '核心追问：你的投资期限是3年还是10年？这个问题的答案，决定了你看待“不合理”的方式。',
             font_size=9, color=GREY_TEXT)


# ================================================================
# SLIDE 7: 第四个不合理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

add_text_box(slide, Cm(2), Cm(2), Cm(5), Cm(3),
             '04', font_size=60, bold=True, color=GOLD)
add_text_box(slide, Cm(2), Cm(5), Cm(28), Cm(3),
             '与科技企业共建生态，而非单纯放贷',
             font_size=28, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(8.5), Cm(5), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(2), Cm(9.5), Cm(28), Cm(4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '传统银企关系：银行提供资金，企业支付利息。\n\n但科技企业需要的不仅是资金——它们需要的是理解技术逻辑的金融合作伙伴。\n\n与科技企业共建生态意味着：\n• 从“资金提供方”进化为“战略合伙人”\n• 从“风险监控”进化为“共同成长”\n• 从“产品推销”进化为“方案共创”'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(22)

add_rect(slide, Cm(2), Cm(14.5), Cm(30), Cm(1.8), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(3), Cm(14.7), Cm(28), Cm(1.5),
             '案例：邮储×恒生电子——AI联合创新实验室模式，从甲乙方关系变为联合创新伙伴。',
             font_size=11, color=GOLD)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.7), Cm(28), Cm(0.5),
             '核心追问：当科技企业不再“需要”你的贷款时，你们还剩下什么关系？——答案是生态。',
             font_size=9, color=GREY_TEXT)


# ================================================================
# SLIDE 8: 第五个不合理
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

add_text_box(slide, Cm(2), Cm(2), Cm(5), Cm(3),
             '05', font_size=60, bold=True, color=GOLD)
add_text_box(slide, Cm(2), Cm(5), Cm(28), Cm(3),
             '建立系统化的科技金融投资框架',
             font_size=28, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(8.5), Cm(5), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(2), Cm(9.5), Cm(28), Cm(4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '当科技金融只是少数人的“直觉判断”时——它是不可复制的。\n\n当科技金融变成一套可持续的方法论时——它才真正成为一门生意。\n\n建立框架的意义：\n• 将个人经验转化为组织能力\n• 将直觉判断转化为可复制的决策流程\n• 将偶然成功转化为持续胜利'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(22)

add_rect(slide, Cm(2), Cm(14.5), Cm(30), Cm(1.8), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(3), Cm(14.7), Cm(28), Cm(1.5),
             '道·法·术·势框架：道(哲学根基) → 法(方法论) → 术(实践工具) → 势(趋势判断)',
             font_size=11, color=GOLD)

add_rect(slide, Cm(0), Cm(18.5), W, Cm(0.55), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.7), Cm(28), Cm(0.5),
             '核心追问：10年后，当科技金融成为每家银行的标配时，率先建立系统化框架的人会站在哪里？',
             font_size=9, color=GREY_TEXT)


# ================================================================
# SLIDE 9: 共性逻辑
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_number(slide, 3, '共性逻辑')

add_text_box(slide, Cm(1.5), Cm(3), Cm(30), Cm(2),
             '五个“不合理”的共性逻辑',
             font_size=28, bold=True, color=DARK_BLUE)
add_rect(slide, Cm(1.5), Cm(5.2), Cm(4), Cm(0.06), GOLD)

cards = [
    ('01 逆周期', '在寒冬播种', '认知深度'),
    ('02 新尺度', '以未来定价', '改变范式'),
    ('03 长期主义', '用时间换空间', '战略定力'),
    ('04 生态思维', '从交易到共生', '关系升级'),
    ('05 系统化', '从直觉到框架', '能力固化'),
]

for i, (num_label, action, quality) in enumerate(cards):
    x = Cm(1.5 + i * 6.2)
    y = Cm(6.5)
    add_rect(slide, x, y, Cm(5.8), Cm(7.5), LIGHT_BLUE)
    add_text_box(slide, x + Cm(0.5), y + Cm(0.5), Cm(4), Cm(1),
                 num_label, font_size=16, bold=True, color=MID_BLUE)
    add_text_box(slide, x + Cm(0.5), y + Cm(2), Cm(4.8), Cm(2),
                 action, font_size=24, bold=True, color=DARK_BLUE, line_spacing=1.2)
    add_text_box(slide, x + Cm(0.5), y + Cm(4.2), Cm(4), Cm(1),
                 '↓', font_size=20, color=GOLD)
    add_text_box(slide, x + Cm(0.5), y + Cm(5.2), Cm(4.8), Cm(1.5),
                 quality, font_size=14, bold=False, color=MID_BLUE)

add_rect(slide, Cm(1.5), Cm(15), Cm(30.5), Cm(2.5), DARK_BLUE)
txBox = make_text_box(slide, Cm(2.5), Cm(15.3), Cm(28), Cm(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '所有“不合理”的决定都有一个共性：\n它们挑战的是现有的思维范式，而不是基本的商业逻辑。\n10年后的“明智”，源自今天的“敢不同”——但“不同”必须有框架支撑。'
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.line_spacing = Pt(21)


# ================================================================
# SLIDE 10: 结语
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Cm(0), Cm(0), W, Cm(0.3), GOLD)

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(24), Cm(10), Cm(14), Cm(14))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x15, 0x2D, 0x55)
shape.line.fill.background()

add_text_box(slide, Cm(3), Cm(3), Cm(26), Cm(4),
             '回到最初的问题',
             font_size=36, bold=True, color=WHITE)
add_rect(slide, Cm(3), Cm(7.5), Cm(6), Cm(0.06), GOLD)

txBox = make_text_box(slide, Cm(3), Cm(8.5), Cm(26), Cm(6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = ('“如果10年后回看今天，哪些今天看起来“不合理”的决定，\n'
          '会被证明是明智的？”\n\n'
          '我的回答是：\n'
          '那些在所有人犹疑时敢于下注的人，\n'
          '那些用10年的尺度思考问题的人，\n'
          '那些不仅提供资金、更提供认知的人。\n\n'
          '不是因为他们是冒险家。\n'
          '而是因为他们理解：\n\n'
          '科技金融的第一性原理不是风险，而是信任。\n'
          '科技的第一性原理不是效率，而是可能。')
p.font.name = 'Microsoft YaHei'
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.line_spacing = Pt(26)

add_rect(slide, Cm(0), Cm(18), W, Cm(1), RGBColor(0x15, 0x2D, 0x55))
add_text_box(slide, Cm(2), Cm(18.2), Cm(28), Cm(0.8),
             '邮储银行上海分行科技金融事业部  |  2026年5月  |  感谢聆听  |  欢迎交流',
             font_size=10, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)


# ================================================================
# SAVE
# ================================================================
output_dir = '/Users/cyingfang/WorkBuddy/20260429082054'
output_path = os.path.join(output_dir, '科技金融公开演讲_如果10年后回看今天.pptx')
prs.save(output_path)
print(f'Saved: {output_path} | Slides: {len(prs.slides)}')
