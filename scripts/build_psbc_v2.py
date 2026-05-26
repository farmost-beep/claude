#!/usr/bin/env python3
"""
邮储银行上海分行科技金融业务完整介绍 PPT (优化版)
Style reference: 小微企业主办行综合金融服务手册.pdf
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

OUTPUT = '/Users/cyingfang/Documents/科技金融/邮储银行上海分行科技金融业务完整介绍_优化版.pptx'

# ===== COLOR PALETTE (from reference manual) =====
GREEN = RGBColor(0x00, 0x96, 0x53)          # 邮储绿 primary
GREEN_DARK = RGBColor(0x00, 0x6B, 0x3B)     # dark green
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xEE)    # light green bg
GREEN_BRIGHT = RGBColor(0x00, 0xB3, 0x64)   # bright accent green
NAVY = RGBColor(0x17, 0x2B, 0x4F)           # deep navy blue
NAVY_LIGHT = RGBColor(0x2C, 0x3E, 0x6B)     # lighter navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x17)

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
SW = prs.slide_width
SH = prs.slide_height

# ===== HELPERS =====
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = 'Microsoft YaHei'; p.alignment = align
    return tb

def mtext(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.alignment = align
    return tb

def bar(slide, l, t, w, color=GREEN, h=4):
    return rect(slide, l, t, w, h, color)

def pn(slide, n):
    txt(slide, SW-Emu(1500000), SH-Emu(450000), Emu(1300000), Emu(350000),
        str(n), size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def footer(slide):
    rect(slide, 0, SH-Emu(350000), SW, Emu(350000), LIGHT_GRAY)
    txt(slide, Emu(500000), SH-Emu(300000), Emu(8000000), Emu(250000),
        '邮储银行上海分行 · 科技金融业务介绍 · 内部资料 · 2026年5月', size=8, color=MID_GRAY)

def section_header(slide, num, title, sub=''):
    """Full-slide section divider in navy"""
    bg(slide, NAVY)
    # Triangle decoration
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Emu(SW-Emu(3600000)), 0, Emu(3600000), Emu(3600000))
    tri.fill.solid(); tri.fill.fore_color.rgb = GREEN; tri.line.fill.background()
    tri.rotation = 180.0
    txt(slide, Emu(1000000), Emu(1200000), Emu(8000000), Emu(500000),
        f'PART {num}', size=14, color=GREEN_BRIGHT, bold=True)
    txt(slide, Emu(1000000), Emu(1800000), Emu(9000000), Emu(1000000),
        title, size=42, color=WHITE, bold=True)
    bar(slide, Emu(1000000), Emu(3000000), Emu(2500000), GREEN, 3)
    if sub:
        txt(slide, Emu(1000000), Emu(3200000), Emu(9000000), Emu(500000),
            sub, size=18, color=RGBColor(0xAA, 0xBB, 0xCC))

def page_header(slide, title, sub=''):
    """Standard page header with navy bar"""
    rect(slide, 0, 0, SW, Emu(1100000), NAVY)
    # Green accent line at bottom of header
    bar(slide, 0, Emu(1100000), SW, GREEN, 4)
    txt(slide, Emu(600000), Emu(250000), Emu(10000000), Emu(550000),
        title, size=26, color=WHITE, bold=True)
    if sub:
        txt(slide, Emu(600000), Emu(750000), Emu(10000000), Emu(300000),
            sub, size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

def card(slide, l, t, w, h, title, lines, title_color=NAVY):
    """Shadow card with top accent"""
    rect(slide, l, t, w, h, WHITE, LIGHT_GRAY)
    bar(slide, l, t, w, title_color, 5)
    txt(slide, l+Emu(180000), t+Emu(180000), w-Emu(360000), Emu(350000),
        title, size=14, color=title_color, bold=True)
    for i, line in enumerate(lines):
        txt(slide, l+Emu(180000), t+Emu(600000)+i*Emu(260000), w-Emu(360000), Emu(240000),
            f'▸ {line}', size=11, color=DARK_GRAY)

def green_bullet_text(slide, l, t, w, lines, size=12, max_lines=6):
    """Text with green triangle bullets"""
    for i, line in enumerate(lines[:max_lines]):
        txt(slide, l+Emu(50000), t+i*Emu(280000), Emu(150000), Emu(260000),
            '▶', size=8, color=GREEN)
        txt(slide, l+Emu(250000), t+i*Emu(280000), w-Emu(250000), Emu(260000),
            line, size=size, color=DARK_GRAY)

def metric_card(slide, l, t, w, h, value, unit, label, color=GREEN):
    """Big number metric card"""
    rect(slide, l, t, w, h, WHITE, LIGHT_GRAY)
    bar(slide, l, t, w, color, 4)
    txt(slide, l+Emu(100000), t+Emu(200000), w-Emu(200000), Emu(500000),
        value, size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l+Emu(100000), t+Emu(700000), w-Emu(200000), Emu(250000),
        unit, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, l+Emu(100000), t+Emu(950000), w-Emu(200000), Emu(350000),
        label, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

def flow_arrow(slide, l, t, w, color=GREEN):
    """Right-pointing arrow for flow diagrams"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(l), Emu(t), Emu(w), Emu(350000))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color; arrow.line.fill.background()

def info_banner(slide, l, t, w, h, text, bg_color=GREEN_LIGHT, text_color=GREEN_DARK):
    """Highlighted info banner"""
    rect(slide, l, t, w, h, bg_color)
    bar(slide, l, t, Emu(50000), color=GREEN, h=h)
    txt(slide, l+Emu(180000), t+Emu(150000), w-Emu(200000), h-Emu(300000),
        text, size=12, color=text_color, bold=True)

# ================================================================
# SLIDE 1: COVER (Reference style)
# ================================================================
print("Building slides...")
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, NAVY)

# Large green accent block on right side
rect(slide, Emu(7000000), 0, Emu(5200000), SH, GREEN_DARK)
# Green overlay ribbon
rect(slide, Emu(6000000), Emu(500000), Emu(6200000), Emu(1500000), GREEN)

# Left side title
txt(slide, Emu(800000), Emu(1200000), Emu(5500000), Emu(700000),
    '中国邮政储蓄银行', size=26, color=WHITE)
txt(slide, Emu(800000), Emu(2000000), Emu(5500000), Emu(1200000),
    '上海分行', size=48, color=WHITE, bold=True)
bar(slide, Emu(800000), Emu(3300000), Emu(2500000), GREEN, 4)
txt(slide, Emu(800000), Emu(3500000), Emu(5500000), Emu(600000),
    '科技金融业务完整介绍', size=28, color=WHITE)

# Right side - PSBC branding
txt(slide, Emu(7500000), Emu(200000), Emu(4000000), Emu(400000),
    'POSTAL SAVINGS BANK OF CHINA', size=10, color=RGBColor(0x88, 0xAA, 0x88))

# Bottom info
txt(slide, Emu(800000), Emu(4800000), Emu(5000000), Emu(400000),
    '科技金融事业部  |  内部资料  |  2026年5月', size=12, color=RGBColor(0x99, 0xAA, 0xBB))
txt(slide, Emu(800000), Emu(5200000), Emu(5000000), Emu(400000),
    '"投早、投小、投硬核、投长期"  ·  打造科技金融生力军', size=14, color=GREEN_BRIGHT)

bar(slide, 0, SH-Emu(60000), SW, GREEN, 3)
print("Slide 1: Cover")

# ================================================================
# SLIDE 2: TABLE OF CONTENTS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '目  录', 'CONTENTS')

toc = [
    ('01', '战略定位与时代机遇', '政策东风 · 市场空间 · 全行战略 · 上海定位'),
    ('02', '产品体系与核心能力', '"U益创"五大矩阵 · 全生命周期产品 · "看未来"评价 · 数字化平台'),
    ('03', '差异化竞争优势', '九大优势 · 万亿俱乐部对标 · 综合化经营生态'),
    ('04', '营销实战与拓客路径', '园区深耕 · 产业链批量 · 投贷联动 · 政府合作 · 数字化触达 · 五步营销法'),
    ('05', '客户服务与典型案例', '三大标杆案例 · 长三角服务成果 · 客户服务流程'),
    ('06', '风险管理与合规保障', '四大风险防控 · 三维评价体系 · 尽职免责 · 合规红线'),
    ('07', '合作模式与发展蓝图', '合作路径 · 重点领域 · "早小硬"蓝海 · 三步走路线图'),
]

for i, (num, title, desc) in enumerate(toc):
    y = Emu(1500000) + i * Emu(720000)
    rect(slide, Emu(600000), y, Emu(11000000), Emu(650000),
         GREEN_LIGHT if i % 2 == 0 else WHITE, LIGHT_GRAY)
    txt(slide, Emu(800000), y+Emu(80000), Emu(800000), Emu(450000),
        num, size=30, color=GREEN, bold=True)
    txt(slide, Emu(1800000), y+Emu(80000), Emu(4000000), Emu(300000),
        title, size=17, color=NAVY, bold=True)
    txt(slide, Emu(1800000), y+Emu(380000), Emu(9000000), Emu(250000),
        desc, size=10, color=MID_GRAY)

footer(slide); pn(slide, 2)
print("Slide 2: Agenda")

# ================================================================
# SLIDE 3: SECTION DIVIDER - Strategy
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 1, '战略定位与时代机遇',
               '科技金融上升为国家战略，长三角科创浪潮全面来临')
pn(slide, 3)
print("Slide 3: Sec-Strategy")

# ================================================================
# SLIDE 4: Policy Landscape
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '宏观政策体系——科技金融全面提速', '2026年政府工作报告、"十五五"规划纲要、七部委方案、三部门联合发文层层加码')

policies = [
    ('2026年\n政府工作报告', '引导金融机构加力支持科技创新；打造集成电路、航空航天、生物医药、低空经济等新兴支柱产业；大力发展创业投资、天使投资；加强科技创新全链条全生命周期金融服务'),
    ('十五五\n规划纲要', '大力发展科技金融、绿色金融、普惠金融、养老金融、数字金融"五篇大文章"；加快高水平科技自立自强，引领发展新质生产力；促进创新链产业链资金链人才链深度融合'),
    ('七部委\n工作方案', '全面加强金融服务专业能力建设；打造"天使投资—创业投资—私募股权投资—银行贷款—资本市场"多元化接力式金融服务；完善科技信贷管理体制'),
    ('2026.4.30\n三部门发文', '科技创新再贷款扩大至14个领域；民营中小企业首次获得再贷款纳入资格；AI设备和软件服务获专项金融服务支持；电子信息、人工智能等新领域投放空间全面打开'),
]

for i, (title, desc) in enumerate(policies):
    x = Emu(400000) + i * Emu(2900000)
    card(slide, x, Emu(1400000), Emu(2700000), Emu(4800000), title, desc.split('；'), NAVY if i < 2 else GREEN)

# Bottom banner
info_banner(slide, Emu(400000), Emu(6350000), Emu(11400000), Emu(380000),
    '💡 政策窗口期判断：2026年是科技金融政策红利集中释放年，三部门联合扩大再贷款+七部委多元化接力式服务+政府工作报告首次明确科技金融首位度——邮储银行科技金融事业部面临三重政策叠加的历史性机遇',
    GREEN_LIGHT, GREEN_DARK)

footer(slide); pn(slide, 4)
print("Slide 4: Policy")

# ================================================================
# SLIDE 5: Market Opportunity
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '长三角科技金融市场机遇', '中国科技创新最活跃区域，TAM/SAM/SOM三层市场空间测算')

# Market sizing
cols = [
    ('TAM 总可用市场', '长三角全年新增科技金融\n需求超 5,000亿元', GREEN),
    ('SAM 可服务市场', '邮储银行年均可服务\n需求约 800-1,000亿元', NAVY),
    ('SOM 可获市场', '3年目标市场份额3-5%\n年新增500-800亿元', ACCENT_BLUE),
]
for i, (title, desc, color) in enumerate(cols):
    x = Emu(400000) + i * Emu(3900000)
    rect(slide, x, Emu(1400000), Emu(3700000), Emu(2000000), WHITE, LIGHT_GRAY)
    bar(slide, x, Emu(1400000), Emu(3700000), color, 5)
    txt(slide, x+Emu(150000), Emu(1500000), Emu(3400000), Emu(500000),
        title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(150000), Emu(2100000), Emu(3400000), Emu(1000000),
        desc, size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Key indicators
indicators = [
    ('长三角高新企业', '10万+家', '科创板上市全国领先'),
    ('"早小硬"覆盖率', '<20%', '理论市场空间1200亿'),
    ('六大行科技贷款', '>23万亿', '信贷加速向长三角集聚'),
    ('邮储未来五年承诺', '≥3万亿', '科技金融融资总规模'),
]
for i, (label, value, note) in enumerate(indicators):
    x = Emu(400000) + i * Emu(2900000)
    rect(slide, x, Emu(3700000), Emu(2700000), Emu(1200000), GREEN_LIGHT)
    txt(slide, x+Emu(150000), Emu(3800000), Emu(2400000), Emu(300000),
        label, size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(150000), Emu(4100000), Emu(2400000), Emu(450000),
        value, size=26, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(150000), Emu(4550000), Emu(2400000), Emu(250000),
        note, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Shanghai industry focus
txt(slide, Emu(400000), Emu(5200000), Emu(11000000), Emu(300000),
    '上海四大重点科创产业与拓客机遇', size=16, color=NAVY, bold=True)
industries = [
    ('集成电路', '张江/临港\n芯片设计制造封测', GREEN),
    ('人工智能', '张江/西岸\n大模型/算力/应用', NAVY),
    ('生物医药', '张江药谷\n创新药/医疗器械', ACCENT_BLUE),
    ('低空经济', '临港新片区\n低空/火箭/卫星', ACCENT_ORANGE),
]
for i, (name, region, color) in enumerate(industries):
    x = Emu(400000) + i * Emu(2900000)
    rect(slide, x, Emu(5600000), Emu(2700000), Emu(1100000), color)
    txt(slide, x+Emu(100000), Emu(5650000), Emu(2500000), Emu(450000),
        name, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(100000), Emu(6100000), Emu(2500000), Emu(500000),
        region, size=11, color=WHITE, align=PP_ALIGN.CENTER)

footer(slide); pn(slide, 5)
print("Slide 5: Market")

# ================================================================
# SLIDE 6: PSBC National Data
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '邮储银行科技金融全行发展成果', '截至2026年Q1——科技金融"生力军"实现里程碑式突破')

# Top metrics row
metrics = [
    ('11万+', '户', '服务科技企业', GREEN),
    ('9,500亿', '元', '科技贷款余额', GREEN_DARK),
    ('17.75%', '', '科技贷款增速(国有大行第一)', NAVY),
    ('1,064亿', '元', '科创贷产品结余', ACCENT_BLUE),
    ('1.65%', '', '净息差(国有大行最优)', GREEN),
    ('0.99%', '', '不良率(行业领先)', NAVY),
]
for i, (val, unit, label, color) in enumerate(metrics):
    x = Emu(300000) + i * Emu(2000000)
    metric_card(slide, x, Emu(1400000), Emu(1800000), Emu(1600000), val, unit, label, color)

# Bottom: 2026 targets
rect(slide, Emu(300000), Emu(3200000), Emu(11600000), Emu(500000), NAVY)
txt(slide, Emu(600000), Emu(3250000), Emu(3000000), Emu(400000),
    '2026年全行目标', size=16, color=WHITE, bold=True)

targets = [
    ('科技贷款净增', '1,500-1,800亿'),
    ('科技客户净增', '3万户'),
    ('FTP优惠', '科创贷20BP激励'),
    ('经济资本优惠', '90%系数'),
]
for i, (label, value) in enumerate(targets):
    x = Emu(3800000) + i * Emu(2100000)
    txt(slide, x, Emu(3250000), Emu(1900000), Emu(400000),
        f'{label}：{value}', size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Non-interest income highlight
info_banner(slide, Emu(300000), Emu(3900000), Emu(11600000), Emu(380000),
    '📈 第二增长曲线：Q1手续费及佣金净收入 +16.83%（近三年新高）| 投行收入 +38.5% | 交易银行收入 +32.7% | 票据交易非息 10.69亿（+50%）——为科技金融投贷联动提供强力支撑')

# 万亿俱乐部 positioning
rect(slide, Emu(300000), Emu(4500000), Emu(11600000), Emu(2200000), GREEN_LIGHT)
txt(slide, Emu(600000), Emu(4600000), Emu(11000000), Emu(400000),
    '🏆 行业地位：距"万亿俱乐部"仅差500亿元', size=18, color=GREEN_DARK, bold=True)

club_data = [
    ('工商银行', '6万亿', '≈20%'), ('建设银行', '5.25万亿', '18.9%'),
    ('农业银行', '4.7万亿', '20.1%'), ('中国银行', '4.82万亿', '18.8%'),
    ('交通银行', '1.58万亿', '10.7%'), ('兴业银行', '1.12万亿', '18.5%'),
    ('中信银行', '1.07万亿', '14.8%'), ('浦发银行', '1.05万亿', '—'),
    ('招商银行', '1.04万亿', '8.1%'), ('邮储银行 ★', '0.95万亿', '>13%'),
]
for i, (name, amt, growth) in enumerate(club_data):
    col = i % 5
    row = i // 5
    x = Emu(400000) + col * Emu(2300000)
    y = Emu(5100000) + row * Emu(700000)
    is_psbc = '邮储' in name
    c = GREEN if is_psbc else WHITE
    rect(slide, x, y, Emu(2200000), Emu(600000), c, GREEN if is_psbc else LIGHT_GRAY)
    txt(slide, x+Emu(80000), y+Emu(50000), Emu(1500000), Emu(250000),
        name, size=10 if is_psbc else 10, color=GREEN if is_psbc else DARK_GRAY, bold=is_psbc)
    txt(slide, x+Emu(80000), y+Emu(280000), Emu(1500000), Emu(250000),
        f'{amt} ({growth})', size=10, color=DARK_GRAY)

footer(slide); pn(slide, 6)
print("Slide 6: National Data")

# ================================================================
# SLIDE 7: PSBC Strategy
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '邮储银行科技金融战略架构', '"六化"升级总纲 × "五大行动"战略嵌合 × "1+N"组织体系')

# Left: strategy architecture
rect(slide, Emu(400000), Emu(1400000), Emu(5600000), Emu(5200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(600000), Emu(1500000), Emu(5000000), Emu(400000),
    '战略总纲："六化"升级', size=18, color=NAVY, bold=True)
strategies = [
    '特色化——打造具有邮储特色的科技金融差异化路径',
    '轻型化——轻资本运营，中间业务收入占比持续提升',
    '综合化——"商行+投行+保险+投资"一站式服务',
    '生态化——"五子联动"邮政金融千亿级基金生态',
    '精细化——29家科技金融事业部+70家特色网点专营',
    '数智化——"技术流"评价+鸿蒙生态银行+260+AI应用场景',
]
for i, s in enumerate(strategies):
    txt(slide, Emu(700000), Emu(2000000)+i*Emu(500000), Emu(5000000), Emu(450000),
        f'▸ {s}', size=12, color=DARK_GRAY)

# Right: five actions
txt(slide, Emu(6500000), Emu(1400000), Emu(5000000), Emu(400000),
    '"五大行动"与科技金融嵌合', size=18, color=GREEN_DARK, bold=True)

actions = [
    ('公司业务提升', '1+N综合服务，新增30万户公司客户\n产业链链式拓客，聚焦新动能领域', GREEN),
    ('城市业务攻坚', '重点城市行优先发展，私行客户三年翻番\n聚焦城市科创客群，40+私行中心', NAVY),
    ('网点效能跃升', '近4万网点转型"获客服务价值阵地"\nAI应用场景260+，数字人服务', ACCENT_BLUE),
    ('手机银行突破', '月活超9000万户，AI赋能智能触达\n线上获客→运营→线下服务全链路闭环', ACCENT_ORANGE),
    ('服务强县富镇', '普惠小微贷款1.90万亿(Q1)\n下沉市场科创企业全覆盖', GREEN_DARK),
]
for i, (title, desc, color) in enumerate(actions):
    y = Emu(1900000) + i * Emu(840000)
    rect(slide, Emu(6500000), y, Emu(5200000), Emu(750000), color)
    txt(slide, Emu(6700000), y+Emu(80000), Emu(4800000), Emu(300000),
        title, size=14, color=WHITE, bold=True)
    txt(slide, Emu(6700000), y+Emu(380000), Emu(4800000), Emu(350000),
        desc, size=10, color=WHITE)

footer(slide); pn(slide, 7)
print("Slide 7: Strategy")

# ================================================================
# SLIDE 8: Section - Products
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, '产品体系与核心能力',
               '"U益创"五大产品矩阵 · 全生命周期覆盖 · "看未来"六维评价 · 数字化赋能')
pn(slide, 8)
print("Slide 8: Sec-Products")

# ================================================================
# SLIDE 9: U益创 Overview
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '"U益创"全周期科技金融服务体系', '核心品牌：U益创（优益创）——五大一体化产品矩阵，覆盖科技企业完整生命周期')

# Five product pillars
colors = [GREEN, NAVY, ACCENT_ORANGE, ACCENT_BLUE, GREEN_DARK]
pillars = [
    ('U创融', '融', '资金融通', ['初创贷——无营收初创期专属', '科创e贷——线上化随借随还', '科创贷——成长期标准化信贷', '上市贷——拟上市企业长期资金', '并购贷——最长10年中长期贷款']),
    ('U创投', '投', '资本运作', ['认股权贷——以未来股权价值融资', '银投连贷——创投+银行协同', '债转股——AIC首笔20亿已落地', '股权投资——中邮资本千亿级基金']),
    ('U创通', '通', '支付结算', ['跨境结算——本外币一体化', '资金管理——全球资金归集', '供应链金融——产业链全流程', '票据服务——灵活贴现+开票']),
    ('U创富', '富', '财富管理', ['企业财富增值——理财+存款', '"邮银财管+"产品体系', '高管个人金融服务', '员工薪酬福利方案']),
    ('U创慧', '慧', '智库咨询', ['行业研究——产业趋势分析', '政策解读——补贴申报指导', '产业对接——上下游撮合', '路演展示——"科融通"平台']),
]
for i, (name, icon, subtitle, items) in enumerate(pillars):
    x = Emu(200000) + i * Emu(2360000)
    rect(slide, x, Emu(1400000), Emu(2200000), Emu(4800000), WHITE, LIGHT_GRAY)
    bar(slide, x, Emu(1400000), Emu(2200000), colors[i], 6)
    txt(slide, x+Emu(80000), Emu(1500000), Emu(2000000), Emu(350000),
        icon, size=22, color=colors[i], bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(80000), Emu(1850000), Emu(2000000), Emu(350000),
        name, size=15, color=colors[i], bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(80000), Emu(2200000), Emu(2000000), Emu(250000),
        subtitle, size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(slide, x+Emu(80000), Emu(2600000)+j*Emu(420000), Emu(2000000), Emu(380000),
            f'▸ {item}', size=10, color=DARK_GRAY)

# Philosophy banner
info_banner(slide, Emu(200000), Emu(6350000), Emu(11800000), Emu(380000),
    '核心理念："投早、投小、投硬核、投长期"  |  不看抵押看未来，不看现在看发展，用"成长性"替代"盈利性"作为核心评价维度')

footer(slide); pn(slide, 9)
print("Slide 9: U益创")

# ================================================================
# SLIDE 10: Lifecycle Matrix
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '全生命周期产品匹配矩阵', '从种子期到成熟期，每个发展阶段均有精准产品覆盖——FPA综合价值最大化')

# Lifecycle stages with products
stages = [
    ('种子期/天使期\n(初创0-3年)', [
        '初创贷 ≤2000万',
        '人才贷（以人定贷）',
        '科创e贷（线上化）',
        '银投连贷 ≤3000万',
        '创新积分贷',
    ], GREEN),
    ('成长期A/B轮\n(3-5年)', [
        '科创贷（三种模式）',
        '投贷联动组合',
        '知识产权质押融资',
        '供应链金融',
        '科技履约贷',
    ], NAVY),
    ('扩张期C轮+\n(5-8年)', [
        '大额科创贷增额',
        '项目贷款',
        '并购贷款 ≤3亿',
        '跨境金融服务',
        '"邮银财管+"',
    ], ACCENT_BLUE),
    ('成熟期Pre-IPO\n(上市前及上市后)', [
        '上市贷（中长期）',
        '并购贷（最长10年）',
        '投融资顾问',
        '高管私行服务',
        '中邮证券IPO辅导',
    ], GREEN_DARK),
]
for i, (stage, products, color) in enumerate(stages):
    x = Emu(200000) + i * Emu(2950000)
    rect(slide, x, Emu(1400000), Emu(2750000), Emu(4800000), WHITE, LIGHT_GRAY)
    rect(slide, x, Emu(1400000), Emu(2750000), Emu(800000), color)
    txt(slide, x+Emu(100000), Emu(1450000), Emu(2550000), Emu(700000),
        stage, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, prod in enumerate(products):
        txt(slide, x+Emu(100000), Emu(2400000)+j*Emu(480000), Emu(2550000), Emu(420000),
            f'▸ {prod}', size=11, color=DARK_GRAY)

# FPA decomposition
rect(slide, Emu(200000), Emu(6400000), Emu(11800000), Emu(350000), GREEN_LIGHT)
txt(slide, Emu(400000), Emu(6400000), Emu(11400000), Emu(350000),
    'FPA综合价值 = 存款余额 + 贷款余额(60-70%) + 理财余额 + 保险保费(10-15%) + 票据/债券(5-10%) + 国际业务折算(5-10%)  |  目标：从单一信贷提供者升级为综合金融服务主办行',
    size=10, color=GREEN_DARK, bold=True)

footer(slide); pn(slide, 10)
print("Slide 10: Lifecycle")

# ================================================================
# SLIDE 11: 科创贷详解
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科创贷——成长期科技企业核心融资工具', '截至2026年4月末：结余1,064亿元 | 不良率1.04% | 占小企业贷款比重超10% | 持续增长态势')

# Three modes
modes = [
    ('模式一：一般模式', '适用7类科技企业资质\n（高新/专精特新/科小等）\n以政府认定资质为准入基础\n综合评估企业成长性\n担保：信用/抵押/保证灵活组合\n额度：根据评级和规模确定', GREEN),
    ('模式二：创新积分模式', '适用已纳入各地"创新积分"\n评价体系并获得积分的企业\n基于创新能力数据授信\n突破传统财务指标限制\n纯信用为主，无需抵押\n额度与积分等级挂钩', NAVY),
    ('模式三：银投联动模式', '适用已获优质创投机构\n股权投资的企业\n"创投先投、银行跟进"\n跟贷额度最高3,000万/户\n合作投资机构分L2-L5等级\n投贷联动综合服务', ACCENT_BLUE),
]
for i, (title, desc, color) in enumerate(modes):
    x = Emu(300000) + i * Emu(3900000)
    rect(slide, x, Emu(1400000), Emu(3700000), Emu(3200000), WHITE, LIGHT_GRAY)
    rect(slide, x, Emu(1400000), Emu(3700000), Emu(550000), color)
    txt(slide, x+Emu(150000), Emu(1450000), Emu(3400000), Emu(500000),
        title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(150000), Emu(2100000), Emu(3400000), Emu(2300000),
        desc, size=12, color=DARK_GRAY)

# Bottom notes
rect(slide, Emu(300000), Emu(4800000), Emu(3800000), Emu(700000), GREEN_LIGHT)
txt(slide, Emu(500000), Emu(4900000), Emu(3400000), Emu(500000),
    '三大风控底线', size=14, color=GREEN_DARK, bold=True)
txt(slide, Emu(500000), Emu(5300000), Emu(3400000), Emu(200000),
    '① 经营正常（非停产/半停产）  ② 资信良好（无不良征信/欠息/违约）  ③ 用途合规（严禁流贷固用/转借/挪用）',
    size=10, color=DARK_GRAY)

rect(slide, Emu(4300000), Emu(4800000), Emu(7600000), Emu(700000), GREEN_LIGHT)
txt(slide, Emu(4500000), Emu(4900000), Emu(7200000), Emu(500000),
    '科创贷专属优惠政策', size=14, color=GREEN_DARK, bold=True)
txt(slide, Emu(4500000), Emu(5300000), Emu(7200000), Emu(200000),
    'FTP内部计价：额外给予20BP优惠  |  经济资本占用：90%优惠系数  |  科技创新再贷款政策贴息叠加  |  绿色审批通道：最快7天完成审批',
    size=10, color=DARK_GRAY)

# Qualification requirements
rect(slide, Emu(300000), Emu(5700000), Emu(11600000), Emu(1000000), WHITE, LIGHT_GRAY)
txt(slide, Emu(500000), Emu(5750000), Emu(5000000), Emu(250000),
    '7类科技企业准入资质', size=13, color=NAVY, bold=True)
quals = ['高新技术企业', '科技型中小企业', '专精特新企业', '国家技术创新示范企业', '制造业单项冠军', '科创板/创业板上市企业', '纳入创新积分评价企业']
for i, q in enumerate(quals):
    x = Emu(500000) + i * Emu(1580000)
    txt(slide, x, Emu(6050000), Emu(1480000), Emu(350000),
        q, size=10, color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

footer(slide); pn(slide, 11)
print("Slide 11: 科创贷")

# ================================================================
# SLIDE 12: 初创贷详解
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '初创贷——"投早投小"的差异化拳头产品', '专为无营收/低营收初创期科技企业设计，破解科技成果转化"最先一公里"融资难题')

# Product specification
specs = [
    ('目标客群', [
        '成立≤5年（技术研发周期长可至10年）',
        '营收<1亿元，研发费用占比≥5%',
        '已有订单但尚未规模盈利',
        '拥有Ⅰ/Ⅱ类知识产权',
        '属于新质生产力范畴企业',
    ]),
    ('两大业务模式', [
        '模式一：人才贷——以人定贷',
        '· 凭核心团队学术/产业背景',
        '· 高层次人才创业企业专享',
        '· 评级BB-及以上即可准入',
        '模式二：银投连贷——投贷联动',
        '· 创投机构已投企业',
        '· 银行跟贷最高3000万',
        '· 合作GP超300家',
    ]),
    ('授信要素', [
        '最高额度：2,000万元',
        '流贷期限：最长2年',
        '固贷期限：最长10年',
        '担保方式：信用为主，灵活组合',
        '审批时效：绿色通道7-10天',
        '还款方式：按季/月付息，到期还本',
        '利率优惠：科创专属定价+FTP 20BP优惠',
    ]),
    ('创新风控逻辑', [
        '• 用未来订单判断现金流',
        '• 关注知识产权+人才+财政支持',
        '• 评估内驱力（技术）与外源动力（资本）',
        '• 三维评价：科学家+投资人+企业家',
        '• 与政府风险补偿池联动缓释风险',
    ]),
]
for i, (title, items) in enumerate(specs):
    col = i % 2
    row = i // 2
    x = Emu(300000) + col * Emu(6000000)
    y = Emu(1400000) + row * Emu(2600000)
    rect(slide, x, y, Emu(5800000), Emu(2400000), WHITE, LIGHT_GRAY)
    rect(slide, x, y, Emu(5800000), Emu(420000), GREEN if row == 0 else NAVY)
    txt(slide, x+Emu(150000), y+Emu(50000), Emu(5500000), Emu(350000),
        title, size=15, color=WHITE, bold=True)
    for j, item in enumerate(items):
        txt(slide, x+Emu(150000), y+Emu(550000)+j*Emu(320000), Emu(5500000), Emu(300000),
            f'{item}', size=11, color=DARK_GRAY)

footer(slide); pn(slide, 12)
print("Slide 12: 初创贷")

# ================================================================
# SLIDE 13: 投贷联动 + 技术流
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '投贷联动 & "看未来"技术流评价体系', '商行+投行协同新模式 | 重新定义科技企业信用评价标准')

# Left: 投贷联动
rect(slide, Emu(300000), Emu(1400000), Emu(5800000), Emu(5200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(500000), Emu(1450000), Emu(5400000), Emu(350000),
    '投贷联动——生态金融新模式', size=17, color=NAVY, bold=True)

# Flow
flow_steps = ['创投发现\n项目', '股权\n投资', '银行\n跟贷', '综合\n服务', '上市/\n退出']
for i, step in enumerate(flow_steps):
    x = Emu(400000) + i * Emu(1060000)
    rect(slide, x, Emu(2000000), Emu(900000), Emu(700000), GREEN if i < 3 else NAVY)
    txt(slide, x+Emu(50000), Emu(2050000), Emu(800000), Emu(600000),
        step, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        flow_arrow(slide, x+Emu(920000), Emu(2150000), Emu(150000), GREEN)

txt(slide, Emu(500000), Emu(2900000), Emu(5400000), Emu(1000000),
    '▶ 服务企业：300+户 | 跟贷规模：200亿+\n▶ 银投联贷额度上限：3000万/户\n▶ 合作投资机构等级：L2-L5(超300家)\n▶ AIC首笔债转股：20亿落地(首批8亿)\n▶ 中邮资本：千亿级基金生态联动',
    size=12, color=DARK_GRAY)

# Right: 技术流
rect(slide, Emu(6300000), Emu(1400000), Emu(5600000), Emu(5200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(6500000), Emu(1450000), Emu(5200000), Emu(350000),
    '"看未来"六维评估体系', size=17, color=GREEN_DARK, bold=True)

dims = [
    ('技术创新力', '专利质量/研发投入/技术壁垒'),
    ('团队竞争力', '核心团队/行业经验/人才储备'),
    ('市场成长力', '市场空间/客户质量/收入趋势'),
    ('资本认可度', '创投背书/估值/融资轮次'),
    ('政策适配度', '高新认定/专精特新/政策支持'),
    ('产业链价值', '产业链地位/上下游稳定/集中度'),
]
for i, (name, desc) in enumerate(dims):
    y = Emu(1950000) + i * Emu(680000)
    rect(slide, Emu(6500000), y, Emu(5200000), Emu(600000), WHITE, LIGHT_GRAY)
    bar(slide, Emu(6500000), y, Emu(60000), color=GREEN, h=Emu(600000))
    txt(slide, Emu(6650000), y+Emu(80000), Emu(2000000), Emu(400000),
        name, size=13, color=NAVY, bold=True)
    txt(slide, Emu(8700000), y+Emu(80000), Emu(2800000), Emu(400000),
        desc, size=11, color=DARK_GRAY)

# Application stats
txt(slide, Emu(6500000), Emu(6100000), Emu(5200000), Emu(400000),
    '📊 覆盖企业1.6万+户 | Q1批复1.17万亿(+12.66%) | 新增11个行业模型',
    size=11, color=GREEN_DARK, bold=True)

footer(slide); pn(slide, 13)
print("Slide 13: 投贷联动+技术流")

# ================================================================
# SLIDE 14: 易企营
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '"易企营"数字化SaaS平台 & 知识产权融资',
            '面向中小微科创企业免费开放，覆盖"人财物产供销"全链路数字化管理')

# 易企营 modules
modules = ['易办公', '易财税', '易薪酬', '易经营', '易发票', '易招采', '易CRM', '易库存', '易报表']
for i, mod in enumerate(modules):
    col = i % 3; row = i // 3
    x = Emu(300000) + col * Emu(4000000)
    y = Emu(1400000) + row * Emu(1500000)
    rect(slide, x, y, Emu(3800000), Emu(1300000), WHITE, LIGHT_GRAY)
    bar(slide, x, y, Emu(3800000), NAVY if row == 0 else GREEN, 5)
    txt(slide, x+Emu(150000), y+Emu(300000), Emu(3500000), Emu(700000),
        mod, size=22, color=NAVY if row == 0 else GREEN, bold=True, align=PP_ALIGN.CENTER)

# Bottom: IP financing
rect(slide, Emu(300000), Emu(6000000), Emu(11600000), Emu(700000), GREEN_LIGHT)
txt(slide, Emu(500000), Emu(6050000), Emu(11200000), Emu(500000),
    '知识产权融资专题：专利质押贷款 · 商标质押贷款 · 著作权质押贷款 | 评估方式：收益法/成本法/市场法 | 质押率：发明专利≤40%、实用新型≤30% | 配套：知识产权保险+维权基金',
    size=12, color=GREEN_DARK)

footer(slide); pn(slide, 14)
print("Slide 14: 易企营")

# ================================================================
# SLIDE 15: Section - Advantages
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, '差异化竞争优势',
               '九大核心优势构筑竞争壁垒 · 万亿俱乐部精准对标 · 综合化经营生态闭环')
pn(slide, 15)
print("Slide 15: Sec-Advantages")

# ================================================================
# SLIDE 16: 9 Advantages
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '邮储银行科技金融九大差异化竞争优势', '依托国有大行信用背书与近4万网点覆盖，形成不可替代的竞争壁垒')

advantages = [
    ('普惠基因', '近4万个网点\n覆盖最广下沉最深', GREEN),
    ('客户基础', '个人客户超6.8亿\n企业服务转化入口', GREEN),
    ('资金成本', '净息差1.65%\n国有大行最优\n定价空间最大', GREEN),
    ('资产质量', '不良率0.99%\n行业领先\n风控体系成熟', NAVY),
    ('技术流', '"看未来"评价体系\n覆盖1.6万户+\n批复3.6万亿+', NAVY),
    ('综合化', '证券+保险+理财+AIC\n千亿级基金联动\n一站式服务', NAVY),
    ('非息增长', '手续费+16.83%\n投行收入+38.5%\n第二增长曲线', ACCENT_BLUE),
    ('品牌信任', '国有大行信用背书\n安全稳健\n长期合作伙伴', GREEN_DARK),
    ('政策响应', '"五大行动"战略\n"六化"升级总纲\n组织敏捷高效', GREEN_DARK),
]
for i, (title, desc, color) in enumerate(advantages):
    col = i % 3; row = i // 3
    x = Emu(300000) + col * Emu(3900000)
    y = Emu(1400000) + row * Emu(1700000)
    rect(slide, x, y, Emu(3700000), Emu(1500000), WHITE, LIGHT_GRAY)
    bar(slide, x, y, Emu(3700000), color, 5)
    txt(slide, x+Emu(150000), y+Emu(180000), Emu(3400000), Emu(400000),
        title, size=16, color=color, bold=True)
    txt(slide, x+Emu(150000), y+Emu(650000), Emu(3400000), Emu(750000),
        desc, size=12, color=DARK_GRAY)

footer(slide); pn(slide, 16)
print("Slide 16: Advantages")

# ================================================================
# SLIDE 17: Competition Deep-Dive (2-page spread feel)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融竞争格局深度对标', '竞争已从"产品同质化"进入"模式差异化"新阶段——精准对标，找准差异化定位')

# Competition table
txt(slide, Emu(400000), Emu(1300000), Emu(2000000), Emu(300000),
    '万亿俱乐部九强', size=14, color=NAVY, bold=True)

banks = [
    ('1', '工商银行', '6万亿', '≈20%', '科技贷款全行业第一'),
    ('2', '建设银行', '5.25万亿', '18.9%', '制造业贷款超5万亿'),
    ('3', '中国银行', '4.82万亿', '18.8%', '长三角近1.3万亿'),
    ('4', '农业银行', '4.7万亿', '20.1%', '300+科技专业支行'),
    ('5', '交通银行', '1.58万亿', '10.7%', '长三角营收占比38%'),
    ('6', '兴业银行', '1.12万亿', '18.5%', '"第四张名片"+AIC'),
    ('7', '中信银行', '1.07万亿', '14.8%', '数字经济贷款2625亿'),
    ('8', '浦发银行', '1.05万亿', '—', '科创板覆盖率>70%'),
    ('9', '招商银行', '1.04万亿', '8.1%', '"1+20+100"六专机制'),
]
for i, (rank, name, amt, growth, position) in enumerate(banks):
    y = Emu(1650000) + i * Emu(340000)
    bg_c = GREEN_LIGHT if i % 2 == 0 else WHITE
    rect(slide, Emu(400000), y, Emu(5600000), Emu(310000), bg_c)
    txt(slide, Emu(450000), y+Emu(20000), Emu(300000), Emu(250000),
        rank, size=10, color=GREEN, bold=True)
    txt(slide, Emu(750000), y+Emu(20000), Emu(1300000), Emu(250000),
        name, size=11, color=DARK_GRAY, bold=True)
    txt(slide, Emu(2100000), y+Emu(20000), Emu(1500000), Emu(250000),
        amt, size=11, color=NAVY, bold=True)
    txt(slide, Emu(3600000), y+Emu(20000), Emu(800000), Emu(250000),
        f'增速:{growth}', size=10, color=MID_GRAY)
    txt(slide, Emu(4400000), y+Emu(20000), Emu(1500000), Emu(250000),
        position, size=10, color=DARK_GRAY)

# Right: PSBC position + key competitors analysis
rect(slide, Emu(6200000), Emu(1300000), Emu(5700000), Emu(1200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(6400000), Emu(1400000), Emu(5300000), Emu(400000),
    '邮储银行 第10位 → 0.95万亿 → 距万亿仅差500亿 ★', size=14, color=GREEN, bold=True)
txt(slide, Emu(6400000), Emu(1800000), Emu(5300000), Emu(500000),
    '按现有增速（>13%）预计2026年内突破万亿，将成为邮储科技金融战略里程碑事件，对上海分行拓客增户信心提振和品牌背书意义重大。',
    size=11, color=DARK_GRAY)

# Key competitor modes
txt(slide, Emu(6400000), Emu(2700000), Emu(5300000), Emu(300000),
    '关键竞品差异化模式解析', size=13, color=NAVY, bold=True)
competitors = [
    ('兴业银行', '"1+20+150"体系+AIC股债贷联动，客户增速18.47%行业领先'),
    ('浦发银行', '"科技五力模型"服务7500家融资超1000亿，510家科技支行'),
    ('中信银行', '"望远镜"评估(技术/团队/未来)+"7+8+N"产品矩阵'),
    ('招商银行', '"1+20+100"阵型+"六专"机制+四链融合'),
    ('民生银行', '投联贷"股债贷保"一站式，重构风险收益逻辑'),
    ('华夏银行', '增速53.74%领跑全行业(中小银行聚焦细分赛道快速迭代)'),
]
for i, (name, desc) in enumerate(competitors):
    y = Emu(3100000) + i * Emu(480000)
    txt(slide, Emu(6500000), y, Emu(1500000), Emu(400000),
        name, size=11, color=NAVY, bold=True)
    txt(slide, Emu(8000000), y, Emu(3800000), Emu(400000),
        desc, size=10, color=DARK_GRAY)

footer(slide); pn(slide, 17)
print("Slide 17: Competition")

# ================================================================
# SLIDE 18: Section - Marketing & Channels
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, '营销实战与拓客路径',
               '六大拓客渠道 · 五步营销法 · 五大重点场景 · 长三角精准布局')
pn(slide, 18)
print("Slide 18: Sec-Marketing")

# ================================================================
# SLIDE 19: Six Channels
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融六大拓客渠道', '园区深耕 × 产业链批量 × 政府平台 × 投贷联动 × 数字化触达 × 生态圈协同')

channels = [
    ('科创园区深耕', '入驻园区金融服务站\n建立企业动态项目库\n定期举办银企对接会\n批量授信方案', GREEN),
    ('产业链批量', '"一行一链多企"对接\n"链主"企业上下游延伸\n产业贷批量授信\n产业链图谱营销', NAVY),
    ('政府平台合作', '对接"创新积分"平台\n参与风险补偿资金池\n科创助力贷/履约贷\n政银产品承接', ACCENT_BLUE),
    ('投贷联动协同', '头部GP战略合作\n共享项目源+尽调\n300+户投贷联动\nAIC债转股项目', GREEN_DARK),
    ('数字化触达', '"技术流"自动授信\n科创e贷线上产品\n鸿蒙生态银行\n易企营SaaS引流', ACCENT_ORANGE),
    ('生态圈协同', '"商行+投行+保险"\n深交所"科融通"路演\n工商联合作机制\n中邮投资一体化', NAVY),
]
for i, (title, desc, color) in enumerate(channels):
    col = i % 3; row = i // 3
    x = Emu(300000) + col * Emu(3900000)
    y = Emu(1400000) + row * Emu(2500000)
    rect(slide, x, y, Emu(3700000), Emu(2300000), WHITE, LIGHT_GRAY)
    rect(slide, x, y, Emu(3700000), Emu(550000), color)
    txt(slide, x+Emu(150000), y+Emu(80000), Emu(3400000), Emu(430000),
        title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(150000), y+Emu(700000), Emu(3400000), Emu(1400000),
        desc, size=12, color=DARK_GRAY)

footer(slide); pn(slide, 19)
print("Slide 19: Channels")

# ================================================================
# SLIDE 20: Marketing 5-Step Method
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融对客营销五步法', '精准识别 → 深度需求挖掘 → 定制方案推介 → 高效签约落地 → 长期深度绑定')

steps_data = [
    ('STEP 1\n精准识别', GREEN, [
        '七类科技资质名单批量触达',
        '人行下发高新/专精特新/科小名单',
        '"泛科技企业"标签预营销',
        '投资机构转介绍客户',
        '园区/科委/火炬中心渠道',
    ]),
    ('STEP 2\n需求挖掘', NAVY, [
        '"三问"挖掘法：行业→阶段→痛点',
        '技术壁垒/市场空间/团队实力评估',
        '现有融资结构分析（股+债）',
        '发展规划与资金缺口测算',
        'FPA综合价值预评估',
    ]),
    ('STEP 3\n方案推介', ACCENT_BLUE, [
        '基于企业生命周期匹配产品组合',
        '"不看抵押看未来"差异化切入',
        '利率优势+FTP优惠展示',
        '成功案例(同业/同阶段)佐证',
        'FPA综合方案（存贷理保票汇）',
    ]),
    ('STEP 4\n签约落地', GREEN_DARK, [
        '绿色通道：最快7天审批',
        '专属客户经理1对1服务',
        '线上+线下多渠道签约',
        '首笔放款全流程跟进',
        '开户/代发/网银同步办理',
    ]),
    ('STEP 5\n深度绑定', ACCENT_ORANGE, [
        '主办行关系确立（结算主账户）',
        'FPA季度复盘与交叉销售',
        '企业成长阶段升级产品匹配',
        '转介上下游产业链客户',
        '长期陪伴→IPO/并购全流程服务',
    ]),
]
for i, (title, color, items) in enumerate(steps_data):
    x = Emu(150000) + i * Emu(2400000)
    rect(slide, x, Emu(1400000), Emu(2250000), Emu(3800000), WHITE, LIGHT_GRAY)
    rect(slide, x, Emu(1400000), Emu(2250000), Emu(1200000), color)
    txt(slide, x+Emu(80000), Emu(1450000), Emu(2100000), Emu(1100000),
        title, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(slide, x+Emu(80000), Emu(2750000)+j*Emu(420000), Emu(2100000), Emu(380000),
            f'{j+1}. {item}', size=10, color=DARK_GRAY)

# Scene-based talk tracks
info_banner(slide, Emu(150000), Emu(5400000), Emu(11900000), Emu(1300000),
    '💬 场景化话术参考\n'
    '初次拜访："您好！我是邮储银行科技金融专属客户经理。我们有一款专为科技企业设计的产品，不看抵押看技术和成长潜力——只要您有科技资质或已获股权投资，'
    '科创贷最高可到3000万，审批最快7天。"\n'
    '投贷联动："我看到贵公司已获得XX基金投资，这非常好！我们与300+知名投资机构建立了合作，可以直接走\'银投联贷\'，额度最高3000万，比普通利率更优惠。"\n'
    '暂时不需要："完全理解！我建议先把授信额度做好，相当于预备一个\'弹药库\'，用的时候随借随还。同时把基本账户和薪酬代发先做好，让邮储成为您的\'第一家银行\'。"',
    GREEN_LIGHT, GREEN_DARK)

footer(slide); pn(slide, 20)
print("Slide 20: 5-Step")

# ================================================================
# SLIDE 21: Section - Cases
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 5, '客户服务与典型案例',
               '三大标杆案例 · 全流程客户服务路径 · 长三角区域服务成果')
pn(slide, 21)
print("Slide 21: Sec-Cases")

# ================================================================
# SLIDE 22: Case 1 - 园区深耕
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '典型案例一：溧阳高新区"园区服务专班"', '园区深耕——批量获客最高效渠道，授信超80亿元')

# Left: case details
txt(slide, Emu(400000), Emu(1400000), Emu(5500000), Emu(350000),
    '背景与挑战', size=16, color=NAVY, bold=True)
txt(slide, Emu(400000), Emu(1800000), Emu(5500000), Emu(800000),
    '江苏溧阳高新区聚集大量动力电池、新能源企业，企业融资需求旺盛但缺乏合适抵质押物，传统信贷模式难以覆盖园区内中小微科技企业。',
    size=12, color=DARK_GRAY)

txt(slide, Emu(400000), Emu(2700000), Emu(5500000), Emu(350000),
    '邮储解决方案', size=16, color=GREEN_DARK, bold=True)
green_bullet_text(slide, Emu(400000), Emu(3150000), Emu(5500000), [
    '成立园区服务专班，派驻专职客户经理驻点',
    '建立园区企业动态项目库，分层分类精细化管理',
    '推出"园区贷"批量授信方案，大幅提高服务效率',
    '与园区运营方建立战略合作+利益共享机制',
    '累计开立银票近27亿元，贴现超100亿元',
], size=12)

# Right: results
rect(slide, Emu(6300000), Emu(1400000), Emu(5500000), Emu(3000000), GREEN_LIGHT)
txt(slide, Emu(6600000), Emu(1500000), Emu(4900000), Emu(400000),
    '服务成效', size=18, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
results = [('80亿+', '园区企业累计授信'), ('27亿', '累计开立银票'),
           ('100亿+', '累计贴现'), ('首选银行', '成为园区首选合作银行')]
for i, (num, label) in enumerate(results):
    col = i % 2; row = i // 2
    x = Emu(6400000) + col * Emu(2700000)
    y = Emu(2100000) + row * Emu(1000000)
    txt(slide, x, y, Emu(2500000), Emu(400000), num, size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x, y+Emu(400000), Emu(2500000), Emu(250000), label, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom: replication
rect(slide, Emu(400000), Emu(4600000), Emu(7800000), Emu(2100000), WHITE, LIGHT_GRAY)
txt(slide, Emu(600000), Emu(4700000), Emu(7400000), Emu(350000),
    '可复制经验 & 对上海分行的启示', size=16, color=NAVY, bold=True)
green_bullet_text(slide, Emu(600000), Emu(5150000), Emu(7400000), [
    '① 专班制是实现园区深耕的最有效组织方式——上海张江/临港/漕河泾可直接复制',
    '② 与园区运营方建立利益共享机制是持续合作关键——存量介绍+增量分成',
    '③ 批量授信方案可大幅提高服务效率——从"一户一策"到"一类一策"',
    '④ 上海目标园区：张江科学城、临港新片区、漕河泾开发区、市北高新',
    '⑤ 每个园区目标：新增授信企业50-100家 → 批量授信5-10亿元',
], size=11)

# Right: 上海落地计划
rect(slide, Emu(8400000), Emu(4600000), Emu(3500000), Emu(2100000), GREEN_LIGHT)
txt(slide, Emu(8600000), Emu(4700000), Emu(3100000), Emu(350000),
    '上海分行行动', size=16, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Emu(8600000), Emu(5200000), Emu(3100000), Emu(1300000),
    '▶ 张江专班已组建\n▶ 临港特色网点已选址\n▶ 漕河泾合作洽谈中\n▶ 2026年目标：\n  深耕3-5个重点园区\n  新增科创客户200家\n  批量授信20-30亿元', size=11, color=DARK_GRAY)

footer(slide); pn(slide, 22)
print("Slide 22: Case1")

# ================================================================
# SLIDE 23: Case 2 - 产业链 + Case 3 - 信用贷
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '典型案例二&三：产业链金融 & 科技信用贷', '集群整体授信+科技信用贷组合——两大模式可批量复制至上海市场')

# Left: Case 2
rect(slide, Emu(300000), Emu(1400000), Emu(5800000), Emu(5200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(500000), Emu(1450000), Emu(5400000), Emu(350000),
    '案例二：新化电子陶瓷"产业链金融"', size=16, color=NAVY, bold=True)
green_bullet_text(slide, Emu(500000), Emu(1950000), Emu(5400000), [
    '• 产业集群有数百家企业，单家难以达到银行授信标准',
    '• 以产业集群整体授信，开发"产业贷"产品',
    '• 从1000万元授信逐步增加到5.2亿元',
    '• 助力某陶瓷企业成长为"制造业单项冠军"',
    '• 其温控器陶瓷产品全球市场份额达80%',
    '• 上海可复制：集成电路/生物医药/AI产业链',
], size=11)

# Right: Case 3
rect(slide, Emu(6300000), Emu(1400000), Emu(5600000), Emu(5200000), WHITE, LIGHT_GRAY)
txt(slide, Emu(6500000), Emu(1450000), Emu(5200000), Emu(350000),
    '案例三：冠甲电子"科技信用贷+贴息"', size=16, color=GREEN_DARK, bold=True)
green_bullet_text(slide, Emu(6500000), Emu(1950000), Emu(5200000), [
    '• 墨盒回收与智能再制造（轻资产科创企业）',
    '• 痛点：跨境电商备货3个月+回收现金结算→资金链拉长至7个月',
    '• 痛点：每增加1美元销售需垫付6-7元流动资金',
    '• 方案："科技信用贷+财政贴息"合计1000万元',
    '• 配套：美元存款+结售汇+美元掉期跨境服务',
    '• 成为"科技金融生力军"标杆案例',
    '• 上海复制：大量跨境电商/轻资产科创企业可批量复制',
], size=11)

# Bottom: Key takeaways
rect(slide, Emu(300000), Emu(5900000), Emu(11600000), Emu(800000), GREEN_LIGHT)
txt(slide, Emu(500000), Emu(6000000), Emu(8000000), Emu(500000),
    '🔑 三大案例核心启示', size=15, color=GREEN_DARK, bold=True)
txt(slide, Emu(500000), Emu(6450000), Emu(11200000), Emu(200000),
    '① 园区专班制+批量授信是最高效拓客模式  ② "核心企业+上下游"产业链批量化可复制至长三角集成电路/生物医药/AI产业链  ③ "科技信用贷+贴息"组合是轻资产科创企业融资最优解',
    size=10, color=DARK_GRAY)

footer(slide); pn(slide, 23)
print("Slide 23: Case2+3")

# ================================================================
# SLIDE 24: Customer Journey
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融客户全流程服务路径', '从初次触达到长期深度绑定——"主办行"关系升级路径')

# Customer journey
journey = [
    ('初次触达', '名单触达\n园区路演\n投资机构转介\n老客户介绍', GREEN),
    ('需求诊断', '"三问"诊断\n行业→阶段→痛点\nFPA预评估\n产品匹配', NAVY),
    ('方案定制', '产品组合设计\n利率/额度测算\n综合服务方案\n同业比较优势', ACCENT_BLUE),
    ('签约落地', '绿色通道审批\n专属经理1对1\n线上+线下签约\n首笔快速放款', GREEN_DARK),
    ('深度经营', 'FPA季度复盘\n交叉销售升级\n产品周期匹配\n上下游转介', ACCENT_ORANGE),
    ('长期绑定', '主办行关系确立\nIPO/并购全流程\n跨境金融扩展\n私行+保险+理财', NAVY),
]
for i, (title, desc, color) in enumerate(journey):
    x = Emu(200000) + i * Emu(1980000)
    rect(slide, x, Emu(1400000), Emu(1800000), Emu(3800000), WHITE, LIGHT_GRAY)
    rect(slide, x, Emu(1400000), Emu(1800000), Emu(1000000), color)
    txt(slide, x+Emu(60000), Emu(1450000), Emu(1700000), Emu(900000),
        title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(60000), Emu(2600000), Emu(1700000), Emu(2400000),
        desc, size=11, color=DARK_GRAY)
    if i < 5:
        txt(slide, x+Emu(1820000), Emu(3300000), Emu(180000), Emu(180000),
            '→', size=18, color=GREEN, bold=True)

# Bottom: metrics
rect(slide, Emu(200000), Emu(5400000), Emu(11800000), Emu(1300000), NAVY)
txt(slide, Emu(500000), Emu(5500000), Emu(11200000), Emu(400000),
    '从单一信贷提供者 → 综合金融服务"主办行"', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Emu(500000), Emu(6000000), Emu(11200000), Emu(500000),
    '目标：让邮储银行成为科技企业的"第一家银行、主要结算行、核心融资行、财富管理行"——FPA综合价值最大化',
    size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

footer(slide); pn(slide, 24)
print("Slide 24: Customer Journey")

# ================================================================
# SLIDE 25: Section - Risk
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 6, '风险管理与合规保障',
               '四大风险防控 · 三维协同评价 · 尽职免责机制 · 合规操作红线')
pn(slide, 25)
print("Slide 25: Sec-Risk")

# ================================================================
# SLIDE 26: Risk Management
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融四大风险与防控体系', '建立"三维协同评价体系"，实现风险可控前提下的业务高质量发展')

risks = [
    ('⚠ 伪科创风险', '企业虚构科创属性\n套取科创贷优惠利率', [
        '交叉验证多项科技资质',
        '实地调研研发团队实力',
        '核实专利质量与有效性',
        '定期复核资质有效期',
        '一票否决制',
    ], ACCENT_RED),
    ('⚠ 技术路线风险', '技术路线失败\n产品研发不及预期', [
        '科学家评技术可行性',
        '投资人评商业化价值',
        '企业家评运营模式',
        '"三维协同评价体系"',
        '投贷联动分散风险',
    ], ACCENT_ORANGE),
    ('⚠ 股权结构风险', 'PE/VC退出后股权不稳\n实控人频繁变更', [
        '核查股权穿透结构',
        '识别一致行动人',
        '关注PE投后管理协议',
        '股权分散企业增加担保',
        '设置实控人变更触发条款',
    ], ACCENT_BLUE),
    ('⚠ 信息不对称风险', '银行不了解企业真实\n技术含量和商业价值', [
        '"技术流"评价体系',
        '创新积分数字化工具',
        '行业专家智库支持',
        '行业研究所深度研究',
        '投资人交叉验证',
    ], NAVY),
]
for i, (title, desc, controls, color) in enumerate(risks):
    col = i % 2; row = i // 2
    x = Emu(300000) + col * Emu(6000000)
    y = Emu(1400000) + row * Emu(2600000)
    rect(slide, x, y, Emu(5800000), Emu(2400000), WHITE, LIGHT_GRAY)
    rect(slide, x, y, Emu(5800000), Emu(550000), color)
    txt(slide, x+Emu(150000), y+Emu(80000), Emu(5500000), Emu(420000),
        title, size=16, color=WHITE, bold=True)
    green_bullet_text(slide, x+Emu(150000), y+Emu(700000), Emu(5500000), [
        f'• 风险：{desc.split(chr(10))[0]}',
        f'• {desc.split(chr(10))[1] if chr(10) in desc else ""}',
    ] + [f'• 防控：{c}' for c in controls], size=10, max_lines=10)

# Bottom: 尽职免责
rect(slide, Emu(300000), Emu(6300000), Emu(11600000), Emu(450000), GREEN_LIGHT)
txt(slide, Emu(500000), Emu(6300000), Emu(11200000), Emu(450000),
    '🛡️ 科技金融尽职免责机制：建议不良容忍度3-5% | 科技金融专营机构差异化考核 | "负面清单+应免尽免" | 参照徽商银行3个百分点不良容忍度标杆实践',
    size=12, color=GREEN_DARK)

footer(slide); pn(slide, 26)
print("Slide 26: Risk")

# ================================================================
# SLIDE 27: Compliance
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '科技金融合规操作四原则', '贷款独立 · 投资透明 · 资金专用 · 尽职勤勉——监管红线不可逾越')

principles = [
    ('贷款独立性', '贷款决策与投资决策严格分离\n银行信贷审批不受投资方影响\n独立风险评估和授信决策\n严禁"以投定贷"或变相兜底', GREEN),
    ('投资透明度', '投贷联动项目信息披露充分\n关联方交易必须如实申报\n投资方与银行信息共享机制\n利益冲突申报与回避制度', NAVY),
    ('资金专用性', '信贷资金必须用于约定用途\n严禁信贷资金违规流入股市/房市\n资金流向全程监控（受托支付）\n定期贷后检查资金用途', ACCENT_BLUE),
    ('尽职勤勉', '贷前调查充分（双人实地调查）\n贷中审查严格（独立审批人）\n贷后管理持续（每季现场+非现场）\n严格按照尽职免责办法执行', GREEN_DARK),
]
for i, (title, desc, color) in enumerate(principles):
    col = i % 2; row = i // 2
    x = Emu(300000) + col * Emu(6000000)
    y = Emu(1400000) + row * Emu(2600000)
    rect(slide, x, y, Emu(5800000), Emu(2400000), WHITE, LIGHT_GRAY)
    bar(slide, x, y, Emu(60000), color=color, h=Emu(2400000))
    txt(slide, x+Emu(250000), y+Emu(200000), Emu(5300000), Emu(400000),
        title, size=18, color=color, bold=True)
    txt(slide, x+Emu(250000), y+Emu(800000), Emu(5300000), Emu(1400000),
        desc, size=12, color=DARK_GRAY)

# Regulatory basis
rect(slide, Emu(300000), Emu(6300000), Emu(11600000), Emu(450000), GREEN_LIGHT)
txt(slide, Emu(500000), Emu(6300000), Emu(11200000), Emu(450000),
    '📋 监管制度依据：科技金融授信业务尽职免责实施细则(2025年版) | 银投联动合作投资机构评价与管理指引 | 中小微企业科创贷业务操作规程(2025年修订版) | 科技企业授信业务尽职认定标准',
    size=11, color=GREEN_DARK)

footer(slide); pn(slide, 27)
print("Slide 27: Compliance")

# ================================================================
# SLIDE 28: Section - Cooperation
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 7, '合作模式与发展蓝图',
               '多元化合作路径 · 四大重点领域 · "早小硬"蓝海市场 · 三步走战略路线图')
pn(slide, 28)
print("Slide 28: Sec-Cooperation")

# ================================================================
# SLIDE 29: Cooperation Models
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '多元化合作模式', '政银企投四方联动，共建科技金融生态圈')

coops = [
    ('政府合作', GREEN, [
        '对接上海市科委，获取高新企业白名单',
        '接入"创新积分"平台，开发"积分贷"',
        '承接科创助力贷/科技履约贷',
        '参与政府风险补偿资金池',
        '联合举办"走进张江"品牌活动',
    ]),
    ('创投机构', NAVY, [
        '与头部GP建立投贷联动战略合作',
        '共享项目源+协同尽职调查',
        'AIC债转股项目推荐与合作',
        '联合举办"科融通"路演',
        '中邮投资一体化平行作业',
    ]),
    ('园区/孵化器', ACCENT_BLUE, [
        '设立园区金融服务站/特色网点',
        '与运营方战略合作+利益共享',
        '推出"园区贷"批量授信方案',
        '孵化器推荐+初创贷快速通道',
        '定期举办银企对接+融资路演',
    ]),
    ('产业链/协会', GREEN_DARK, [
        '围绕"链主"企业拓展上下游',
        '开发"产业贷"批量产品',
        '对接行业协会/产业联盟',
        '"一行一链多企"精准对接',
        '全国工商联常态化合作机制',
    ]),
]
for i, (title, color, items) in enumerate(coops):
    col = i % 2; row = i // 2
    x = Emu(300000) + col * Emu(6000000)
    y = Emu(1400000) + row * Emu(2600000)
    rect(slide, x, y, Emu(5800000), Emu(2400000), WHITE, LIGHT_GRAY)
    rect(slide, x, y, Emu(5800000), Emu(500000), color)
    txt(slide, x+Emu(150000), y+Emu(70000), Emu(5500000), Emu(400000),
        title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    green_bullet_text(slide, x+Emu(150000), y+Emu(700000), Emu(5500000), items, size=12)

footer(slide); pn(slide, 29)
print("Slide 29: Cooperation")

# ================================================================
# SLIDE 30: Focus Areas + Early-Stage
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '重点拓展领域 & "早小硬"蓝海市场', '四大产业赛道 + 早期硬科技企业1200亿市场缺口')

# Four focus areas
areas = [
    ('集成电路与半导体', '芯片设计/制造/封测 · 张江/无锡/合肥', GREEN),
    ('新一代信息技术与AI', '大模型/算力/应用 · 上海/杭州', NAVY),
    ('新能源与新型储能', '储能/锂电/光伏 · 上海/宁波/合肥', ACCENT_ORANGE),
    ('低空经济与商业航天', '低空/火箭/卫星 · 上海临港', GREEN_DARK),
]
for i, (area, region, color) in enumerate(areas):
    x = Emu(300000) + i * Emu(2950000)
    rect(slide, x, Emu(1400000), Emu(2750000), Emu(1500000), color)
    txt(slide, x+Emu(100000), Emu(1450000), Emu(2550000), Emu(700000),
        area, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(100000), Emu(2200000), Emu(2550000), Emu(500000),
        region, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# "早小硬" analysis
rect(slide, Emu(300000), Emu(3200000), Emu(11600000), Emu(3400000), GREEN_LIGHT)
txt(slide, Emu(600000), Emu(3300000), Emu(8000000), Emu(400000),
    '"早小硬"科技企业——1200亿蓝海市场待开拓', size=20, color=GREEN_DARK, bold=True)

txt(slide, Emu(600000), Emu(3850000), Emu(11000000), Emu(250000),
    '长三角"早小硬"科技企业信贷覆盖率不足20%——这是邮储银行相对城商行的核心机会点', size=14, color=NAVY, bold=True)

txt(slide, Emu(600000), Emu(4250000), Emu(3500000), Emu(200000),
    '市场空间测算', size=14, color=GREEN_DARK, bold=True)
txt(slide, Emu(600000), Emu(4550000), Emu(3500000), Emu(200000),
    '10万家×(1-20%)×30%转化率≈2.4万家\n若平均单户授信500万元\n理论市场空间≈1,200亿元', size=11, color=DARK_GRAY)

txt(slide, Emu(4300000), Emu(4250000), Emu(3500000), Emu(200000),
    '邮储银行切入路径', size=14, color=GREEN_DARK, bold=True)
txt(slide, Emu(4300000), Emu(4550000), Emu(3500000), Emu(200000),
    '① 依托"技术流"评价体系\n② 国有大行资金成本优势\n③ 初创信用贷/人才贷/认股权贷\n④ 与国投创合等GP建立合作通道', size=11, color=DARK_GRAY)

txt(slide, Emu(8000000), Emu(4250000), Emu(3500000), Emu(200000),
    '优先行动', size=14, color=GREEN_DARK, bold=True)
txt(slide, Emu(8000000), Emu(4550000), Emu(3500000), Emu(200000),
    'P0: 锁定张江/临港首批50家\nP1: 复制初创贷/人才贷模式\nP2: 与5-10家GP签署合作协议', size=11, color=DARK_GRAY)

# Product comparison
txt(slide, Emu(600000), Emu(5200000), Emu(11000000), Emu(300000),
    '针对"早小硬"企业的五款核心产品', size=14, color=GREEN_DARK, bold=True)
products_early = [
    ('初创信用贷', '无抵押/无流水', 'P0'), ('人才贷2.0', '以人定贷', 'P0'),
    ('认股权贷', '股权稀释顾虑', 'P1'), ('创新积分贷', '轻资产/无抵押', 'P1'),
    ('远期共赢贷', '短期收益难实现', 'P2'),
]
for i, (prod, pain, priority) in enumerate(products_early):
    x = Emu(600000) + i * Emu(2200000)
    rect(slide, x, Emu(5600000), Emu(2000000), Emu(800000), WHITE, LIGHT_GRAY)
    p_color = GREEN if priority == 'P0' else (NAVY if priority == 'P1' else MID_GRAY)
    txt(slide, x+Emu(80000), Emu(5650000), Emu(1850000), Emu(350000),
        prod, size=12, color=p_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(80000), Emu(5950000), Emu(1850000), Emu(200000),
        pain, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)
    txt(slide, x+Emu(80000), Emu(6150000), Emu(1850000), Emu(200000),
        priority, size=9, color=p_color, bold=True, align=PP_ALIGN.CENTER)

footer(slide); pn(slide, 30)
print("Slide 30: Focus Areas")

# ================================================================
# SLIDE 31: Roadmap
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '邮储银行上海分行科技金融发展路线图', '三步走战略——从基础能力建设到长三角科技金融标杆银行')

phases = [
    ('第一阶段：基础能力建设\n（0-6个月）', GREEN, [
        '设立科技金融专营机构',
        '张江/临港增设特色网点',
        '对接科委获取企业白名单',
        '培训"科技金融专员"队伍',
        '上线本地化"技术流"系统',
        '对接"创新积分"平台',
        '与5家头部GP签合作协议',
    ]),
    ('第二阶段：渠道拓展与批量获客\n（6-18个月）', NAVY, [
        '深耕3-5个重点园区',
        '每个园区新增授信50-100家',
        '开发2-3条重点产业链',
        '产业链批量授信10-20亿',
        '新增政策性科创贷款200家',
        '投贷联动项目20-30个',
        '举办首场科创金融品牌峰会',
    ]),
    ('第三阶段：生态圈建设与品牌塑造\n（18-36个月）', GREEN_DARK, [
        '科技金融事业部升级',
        '专营机构达到30家',
        '办年度科创金融峰会',
        '服务科技企业突破2万家',
        '科技金融融资余额>1000亿',
        '构建科创企业数字化平台',
        '成为长三角科技金融标杆银行',
    ]),
]
for i, (phase, color, actions) in enumerate(phases):
    x = Emu(200000) + i * Emu(3950000)
    rect(slide, x, Emu(1400000), Emu(3750000), Emu(3800000), WHITE, LIGHT_GRAY)
    rect(slide, x, Emu(1400000), Emu(3750000), Emu(1100000), color)
    txt(slide, x+Emu(100000), Emu(1450000), Emu(3550000), Emu(1000000),
        phase, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, action in enumerate(actions):
        txt(slide, x+Emu(100000), Emu(2700000)+j*Emu(470000), Emu(3550000), Emu(420000),
            f'✓  {action}', size=11, color=DARK_GRAY)

# Bottom KPIs
rect(slide, Emu(200000), Emu(5400000), Emu(11800000), Emu(1300000), NAVY)
txt(slide, Emu(500000), Emu(5500000), Emu(11200000), Emu(500000),
    '三年战略目标', size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Emu(500000), Emu(6100000), Emu(11200000), Emu(400000),
    '服务科技企业突破 2万家  |  科技金融融资余额突破 1,000亿元  |  科技金融专营机构达到 30家  |  成为长三角科技金融领军银行之一',
    size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

footer(slide); pn(slide, 31)
print("Slide 31: Roadmap")

# ================================================================
# SLIDE 32: Immediate Actions
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
page_header(slide, '近期优先行动清单（1-3个月）', 'P0级行动——抢占政策窗口期与万亿突破品牌红利')

# P0 actions
txt(slide, Emu(400000), Emu(1400000), Emu(11000000), Emu(350000),
    'P0 优先级——立即执行', size=18, color=ACCENT_RED, bold=True)

p0_actions = [
    ('P0-1', '借势万亿突破品牌红利', '邮储科技贷款9500亿距万亿仅差500亿——制作"邮储银行科技金融里程碑"营销素材，强化上海市场品牌认知度', GREEN),
    ('P0-2', '对接中邮资本AIC', 'AIC首笔20亿债转股已落地——梳理长三角优质科技企业债转股目标名单(新能源/AI/集成电路)，争取1-2个债转股项目落地上海', NAVY),
    ('P0-3', '梳理"早小硬"目标客户', '行业覆盖率不足20%——基于"技术流"评价体系，锁定张江/临港首批"早小硬"目标客户50家，制定精准触达方案', ACCENT_BLUE),
    ('P0-4', '研究浦发/兴业对标方案', '浦发"科技五力"服务7500家融资>1000亿，兴业客户增速18.5%——评估邮储"技术流"升级路径，形成对标优化报告', GREEN_DARK),
    ('P0-5', '园区批量拓客启动', '依托"易企营"SaaS免费优势——在张江/临港开展科创企业数字化赋能活动，实现批量获客新增30-50家', ACCENT_ORANGE),
]
for i, (code, title, desc, color) in enumerate(p0_actions):
    y = Emu(1900000) + i * Emu(880000)
    rect(slide, Emu(400000), y, Emu(11400000), Emu(800000), WHITE, LIGHT_GRAY)
    bar(slide, Emu(400000), y, Emu(60000), color=color, h=Emu(800000))
    txt(slide, Emu(700000), y+Emu(80000), Emu(800000), Emu(250000),
        code, size=12, color=color, bold=True)
    txt(slide, Emu(1500000), y+Emu(80000), Emu(3800000), Emu(300000),
        title, size=14, color=NAVY, bold=True)
    txt(slide, Emu(1500000), y+Emu(400000), Emu(9800000), Emu(350000),
        desc, size=10, color=DARK_GRAY)

# P1/P2
rect(slide, Emu(400000), Emu(6400000), Emu(11400000), Emu(350000), GREEN_LIGHT)
txt(slide, Emu(600000), Emu(6400000), Emu(11000000), Emu(350000),
    'P1 跟进：华夏银行增速53.7%策略研究 | 北京分行"1+2+5+3+X"组织架构对标 | "U亿算"算力贷上海复制 | 绿色"降碳贷"上海落地  |  P2 储备：鸿蒙生态银行接入 | 差异化考核机制推动',
    size=10, color=GREEN_DARK, bold=True)

footer(slide); pn(slide, 32)
print("Slide 32: Actions")

# ================================================================
# SLIDE 33: Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, NAVY)

# Decorative triangles
for i in range(3):
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(SW-Emu(6000000)+i*Emu(2000000)), 0, Emu(2500000), Emu(2500000))
    tri.fill.solid(); tri.fill.fore_color.rgb = GREEN_DARK if i % 2 == 0 else GREEN
    tri.line.fill.background(); tri.rotation = 180.0

# Main content
txt(slide, SW//6, Emu(800000), SW*2//3, Emu(800000),
    '感谢聆听', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
bar(slide, SW//3, Emu(1800000), SW//3, GREEN, 3)

txt(slide, SW//6, Emu(2100000), SW*2//3, Emu(600000),
    '携手共建科技金融生态，赋能科创企业全生命周期成长', size=22, color=WHITE, align=PP_ALIGN.CENTER)

# Contact
txt(slide, SW//6, Emu(3100000), SW*2//3, Emu(400000),
    '中国邮政储蓄银行 · 上海分行 · 科技金融事业部', size=16, color=GREEN_BRIGHT, align=PP_ALIGN.CENTER)

contact_items = [
    '"投早、投小、投硬核、投长期"',
    '打造具有邮储特色的科技金融生力军',
]
for i, item in enumerate(contact_items):
    txt(slide, SW//6, Emu(3700000)+i*Emu(400000), SW*2//3, Emu(380000),
        item, size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

# Bottom
bar(slide, 0, SH-Emu(60000), SW, GREEN, 3)
txt(slide, Emu(2000000), SH-Emu(450000), Emu(8000000), Emu(350000),
    '内部资料 · 请勿外传 · 2026年5月 · 邮储银行上海分行科技金融事业部 © 版权所有',
    size=9, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)
pn(slide, 33)

print("Slide 33: Thank You")

# ================================================================
# SAVE
# ================================================================
print(f"\nTotal: 33 slides")
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Size: {os.path.getsize(OUTPUT)/1024/1024:.1f} MB")
