#!/usr/bin/env python3
"""第一部分 PPT —— 扩展优化版（14页，保留原图，增强理解）"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_NAVY = RGBColor(0x0B, 0x2A, 0x45)
MID_BLUE = RGBColor(0x1B, 0x5E, 0x8A)
LIGHT_BLUE = RGBColor(0xE2, 0xED, 0xF5)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
GREY_TEXT = RGBColor(0x77, 0x77, 0x77)
SOFT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xCC, 0x33, 0x33)
CARD_BG = RGBColor(0xFD, 0xF2, 0xE4)

IMG_DIR = '/Users/cyingfang/WorkBuddy/20260429082054/pptx_images'

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s


def add_rrect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s


def add_tb(slide, l, t, w, h, text, font='Calibri', size=18, bold=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT, lsp=1.3):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.name = font; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0)
    if lsp: p.line_spacing = Pt(size * lsp)
    return tb


def mktf(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def addp(tf, text, font='Calibri', size=14, bold=False, color=DARK_TEXT,
         align=PP_ALIGN.LEFT, sa=0, sb=0, lsp=None, first=False):
    if first and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.name = font; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    p.space_after = Pt(sa); p.space_before = Pt(sb)
    if lsp: p.line_spacing = Pt(size * lsp)
    return p


def add_img(slide, name, l, t, w, h=None):
    path = os.path.join(IMG_DIR, name)
    if os.path.exists(path):
        if h is None:
            slide.shapes.add_picture(path, l, t, w)
        else:
            slide.shapes.add_picture(path, l, t, w, h)
        return True
    return False


def slide_header(slide, label, title, add_page=True, pnum=''):
    add_rect(slide, 0, 0, W, Cm(0.1), MID_BLUE)
    add_tb(slide, Cm(2), Cm(0.8), Cm(10), Cm(1), label, size=10, bold=True, color=ACCENT_ORANGE)
    add_tb(slide, Cm(2), Cm(1.6), Cm(30), Cm(1.5), title, size=22, bold=True, color=DARK_NAVY)
    add_rect(slide, Cm(2), Cm(3.2), Cm(3), Cm(0.04), ACCENT_ORANGE)
    if add_page:
        add_rect(slide, 0, Cm(18.5), W, Cm(0.5), SOFT_BG)
        add_tb(slide, 1.5, Cm(18.65), Cm(25), Cm(0.35),
               'Network Community Detection via Neural Embeddings', size=7, color=GREY_TEXT)
        if pnum:
            add_tb(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), pnum, size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


def add_bullet_card(tf, items, size=11, color=DARK_TEXT, first=True):
    for i, item in enumerate(items):
        f = first and i == 0
        addp(tf, f'• {item}', size=size, color=color, first=f)


def add_card(slide, l, t, w, h, color, title, items, title_size=14, item_size=11):
    add_rrect(slide, l, t, w, h, color)
    add_tb(slide, l+Cm(0.8), t+Cm(0.3), w-Cm(1.6), Cm(1), title, size=title_size, bold=True, color=DARK_NAVY)
    tf = mktf(slide, l+Cm(0.8), t+Cm(1.3), w-Cm(1.6), h-Cm(1.6))
    for i, item in enumerate(items):
        addp(tf, f'• {item}', size=item_size, color=DARK_TEXT, first=(i==0))
    return tf


# ================================================================
# SLIDE 1 — TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, 0, 0, W, Cm(0.2), ACCENT_ORANGE)
add_rect(slide, Cm(24), Cm(-2), Cm(14), Cm(10), RGBColor(0x0E, 0x35, 0x55))
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(28), Cm(12), Cm(10), Cm(10))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x0E, 0x35, 0x55); shape.line.fill.background()

add_tb(slide, Cm(3), Cm(3.5), Cm(22), Cm(5),
       'Network Community\nDetection via\nNeural Embeddings',
       size=34, bold=True, color=WHITE, lsp=1.15)
add_rect(slide, Cm(3), Cm(9), Cm(5), Cm(0.08), ACCENT_ORANGE)
tf = mktf(slide, Cm(3), Cm(10), Cm(24), Cm(4))
addp(tf, '证明 node2vec / DeepWalk / LINE 等浅层图嵌入方法', size=14, color=RGBColor(0xBB, 0xCC, 0xDD), first=True)
addp(tf, '在社区检测中可达信息论极限', size=14, color=RGBColor(0xBB, 0xCC, 0xDD), sa=4)
addp(tf, '', size=6)
addp(tf, '文献报告 · 第一部分', size=11, color=GREY_TEXT)
addp(tf, '陈于东  |  2026年5月', size=11, color=GREY_TEXT)
add_rect(slide, 0, Cm(18.5), W, Cm(0.5), SOFT_BG)

# ================================================================
# SLIDE 2 — OUTLINE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '', '', add_page=False)

add_tb(slide, Cm(2), Cm(1.2), Cm(20), Cm(2), '报告提纲', size=28, bold=True, color=DARK_NAVY)
add_rect(slide, Cm(2), Cm(3.2), Cm(4), Cm(0.05), ACCENT_ORANGE)

items = [
    ('背景', [
        ('1', '网络、社区与图表示'),
        ('2', '机器学习与神经网络基础'),
        ('3', '图嵌入方法概述'),
    ]),
    ('问题', [
        ('4', '核心问题：理论空白'),
    ]),
    ('模型', [
        ('5', '随机分块模型（SBM）'),
        ('6', '种植分块模型与混合参数 μ'),
    ]),
    ('结果', [
        ('7', '信息论可检测极限'),
        ('8', '主要结论与证明思路'),
        ('9', '数值验证与启示'),
    ]),
]

x_start = Cm(2)
for si, (section, subs) in enumerate(items):
    x = x_start + Cm(si * 8)
    # Section header
    add_rrect(slide, x, Cm(4.2), Cm(7), Cm(1.2), DARK_NAVY)
    add_tb(slide, x+Cm(0.3), Cm(4.4), Cm(6.4), Cm(1), section, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Sub items
    for i, (num, title) in enumerate(subs):
        y = Cm(5.8 + i * 1.3)
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Cm(0.5), y, Cm(0.9), Cm(0.9))
        c.fill.solid(); c.fill.fore_color.rgb = MID_BLUE; c.line.fill.background()
        tf = c.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(10)
        p.font.bold = True; p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_tb(slide, x+Cm(1.8), y-Cm(0.05), Cm(5), Cm(1.2), title, size=11, color=DARK_TEXT)

# Page info
add_rect(slide, 0, Cm(18.5), W, Cm(0.5), SOFT_BG)
add_tb(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '2/14', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 3 — 网络、社区与图表示
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '背景 1/3', '网络、社区与图表示', pnum='3/14')

# Two images side by side
add_img(slide, 'slide05_img02.jpg', Cm(2), Cm(4), Cm(14.5), Cm(9.5))
add_img(slide, 'slide05_img03.jpg', Cm(17.5), Cm(4), Cm(14.5), Cm(9.5))

# Bottom explanation
tf = mktf(slide, Cm(2), Cm(14), Cm(30), Cm(4))
addp(tf, '网络无处不在', size=13, bold=True, color=DARK_NAVY, first=True)
addp(tf, '从社交网络到交通网络、从金融系统到神经网络——网络以节点和边的形式', size=10, color=DARK_TEXT, sa=2)
addp(tf, '表示复杂系统的结构。社区检测旨在发现网络中"内部紧密、外部稀疏"的功能模块。', size=10, color=DARK_TEXT)
addp(tf, '然而，网络的复杂、高维、离散特性使其有效表示变得高度非平凡。', size=10, color=DARK_TEXT)

# ================================================================
# SLIDE 4 — 机器学习与神经网络基础
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '背景 2/3', '机器学习与神经网络基础', pnum='4/14')

# Left: ML image
add_rrect(slide, Cm(2), Cm(4), Cm(14.5), Cm(9), LIGHT_BLUE)
add_tb(slide, Cm(3), Cm(4.3), Cm(12), Cm(1), '机器学习', size=14, bold=True, color=DARK_NAVY)
add_img(slide, 'slide03_img01.jpg', Cm(2.5), Cm(5.5), Cm(13.5), Cm(7))

# Right: Neural network image
add_rrect(slide, Cm(17.5), Cm(4), Cm(14.5), Cm(9), LIGHT_BLUE)
add_tb(slide, Cm(18.5), Cm(4.3), Cm(12), Cm(1), '神经网络', size=14, bold=True, color=DARK_NAVY)
add_img(slide, 'slide04_img02.jpg', Cm(18), Cm(5.5), Cm(13.5), Cm(7))

# Key insight at bottom
tf = mktf(slide, Cm(2), Cm(13.5), Cm(30), Cm(4.5))
addp(tf, '为什么这里要提 ML 和神经网络？', size=13, bold=True, color=ACCENT_ORANGE, first=True)
addp(tf, '本文要证明的核心命题是关于"神经网络图嵌入方法的社区检测能力"', size=11, color=DARK_TEXT, sa=2)
addp(tf, '——理解 ML 基本概念和神经网络结构，是理解后续证明的前提。', size=11, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '关键：本文讨论的是"浅层线性神经网络"（无多层非线性激活函数），与"深度"学习形成对比。', size=11, bold=True, color=MID_BLUE)

# ================================================================
# SLIDE 5 — 图嵌入方法
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '背景 3/3', '图嵌入方法：DeepWalk / LINE / node2vec', pnum='5/14')

# Three method cards with images
methods = [
    ('DeepWalk', 'slide06_img03.jpg', '随机游走 → Skip-gram'),
    ('LINE', 'slide06_img04.jpg', '一阶+二阶邻近保持'),
    ('node2vec', 'slide06_img02.jpg', '灵活游走策略 (DFS/BFS)'),
]

for i, (name, img, desc) in enumerate(methods):
    x = Cm(2 + i * 10.3)
    add_rrect(slide, x, Cm(4), Cm(9.5), Cm(8.5), LIGHT_BLUE)
    add_tb(slide, x+Cm(0.5), Cm(4.3), Cm(8.5), Cm(1), name, size=16, bold=True, color=DARK_NAVY)
    add_rect(slide, x+Cm(1), Cm(5.5), Cm(7.5), Cm(0.02), MID_BLUE)
    add_img(slide, img, x+Cm(1.5), Cm(5.8), Cm(6.5), Cm(4))
    add_tb(slide, x+Cm(1), Cm(10.2), Cm(7.5), Cm(1.5), desc, size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# How they work
tf = mktf(slide, Cm(2), Cm(13.5), Cm(30), Cm(4.5))
addp(tf, '核心思想', size=13, bold=True, color=DARK_NAVY, first=True)
addp(tf, '将每个节点映射为低维连续向量，使得"相似"的节点在向量空间中距离更近。', size=11, color=DARK_TEXT, sa=2)
addp(tf, '相似性定义因方法而异：DeepWalk 用随机游走共现，LINE 用直接/间接邻居，node2vec 用 biased random walk。', size=11, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '共同特点：都是浅层网络（1层隐藏层），没有非线性激活函数。', size=11, bold=True, color=MID_BLUE)
addp(tf, '差异对比：node2vec 引入了 biased random walk 参数 p 和 q，可灵活控制搜索策略。', size=10, color=GREY_TEXT)

# ================================================================
# SLIDE 6 — 核心问题
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '问题提出', '图嵌入方法的理论空白', pnum='6/14')

# Known vs Unknown
add_rrect(slide, Cm(2), Cm(4), Cm(30), Cm(4.5), LIGHT_BLUE)
add_tb(slide, Cm(3.5), Cm(4.3), Cm(27), Cm(1.2), '已知（实证层面）', size=16, bold=True, color=GREEN)
tf = mktf(slide, Cm(3.5), Cm(5.8), Cm(27), Cm(2.5))
addp(tf, 'node2vec / DeepWalk / LINE 在实践中被广泛用于社区检测任务，取得了非常好的效果。', size=12, color=DARK_TEXT, first=True)
addp(tf, '大量实验表明这些方法能有效发现网络中的功能模块。', size=12, color=DARK_TEXT)

add_rrect(slide, Cm(2), Cm(9), Cm(30), Cm(2), CARD_BG)
add_tb(slide, Cm(3.5), Cm(9.1), Cm(27), Cm(1), '未知（理论层面）', size=16, bold=True, color=ACCENT_ORANGE)
add_tb(slide, Cm(3.5), Cm(10.2), Cm(27), Cm(1), '我们缺乏对这些方法工作原理的清晰理论理解——它们是"黑箱"。', size=12, color=DARK_TEXT)

# Three key questions
add_tb(slide, Cm(2), Cm(11.8), Cm(30), Cm(1.2), '本文回答的三个关键问题', size=14, bold=True, color=DARK_NAVY)

qs = [
    ('Q1', '神经图嵌入在稀疏\n网络上表现如何？'),
    ('Q2', '边稀疏性在多大程度\n上削弱检测能力？'),
    ('Q3', '与传统谱方法\n相比如何？'),
]

for i, (q, text) in enumerate(qs):
    x = Cm(2 + i * 10.3)
    add_rrect(slide, x, Cm(13.2), Cm(9.5), Cm(4.5), SOFT_BG)
    add_tb(slide, x+Cm(0.5), Cm(13.5), Cm(8.5), Cm(1), q, size=14, bold=True, color=MID_BLUE)
    add_tb(slide, x+Cm(0.5), Cm(14.5), Cm(8.5), Cm(2.5), text, size=11, color=DARK_TEXT, lsp=1.4)

# ================================================================
# SLIDE 7 — SBM 模型
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '模型 1/2', '随机分块模型（Stochastic Block Model）', pnum='7/14')

add_img(slide, 'slide07_img03.jpg', Cm(2), Cm(4), Cm(30), Cm(5.5))

# Left: model explanation
add_rrect(slide, Cm(2), Cm(10), Cm(14.5), Cm(7.5), LIGHT_BLUE)
tf = mktf(slide, Cm(3.5), Cm(10.3), Cm(11.5), Cm(7))
addp(tf, '什么是 SBM？', size=14, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=4)
addp(tf, '• 网络社区结构的基本生成模型', size=11, color=DARK_TEXT)
addp(tf, '• 节点被分为 q 个社区（块）', size=11, color=DARK_TEXT)
addp(tf, '• 同社区节点连接概率 = p_in', size=11, color=DARK_TEXT)
addp(tf, '• 不同社区连接概率 = p_out', size=11, color=DARK_TEXT)
addp(tf, '• p_in > p_out → 有社区结构', size=11, color=DARK_TEXT)
addp(tf, '• p_in = p_out → ER随机图', size=11, color=DARK_TEXT)
addp(tf, '', size=6)
addp(tf, '为什么重要？', size=13, bold=True, color=ACCENT_ORANGE)
addp(tf, 'SBM 是社区检测算法的标准基准。', size=11, color=DARK_TEXT)
addp(tf, '评估一个算法能否"检测"社区，', size=11, color=DARK_TEXT)
addp(tf, '就是在 SBM 生成的图上测试。', size=11, color=DARK_TEXT)

# Right: Sparse network note
add_rrect(slide, Cm(17.5), Cm(10), Cm(14.5), Cm(3.5), SOFT_BG)
tf = mktf(slide, Cm(19), Cm(10.3), Cm(11.5), Cm(3))
addp(tf, '稀疏网络（Sparse Regime）', size=13, bold=True, color=DARK_NAVY, first=True)
addp(tf, '本文关注稀疏网络：', size=11, color=DARK_TEXT, sa=2)
addp(tf, 'p_in, p_out 与节点数 n 成反比', size=11, color=DARK_TEXT)
addp(tf, '平均度 ⟨k⟩ 不依赖于 n', size=11, color=DARK_TEXT)
addp(tf, '这是实际中最具挑战的场景。', size=11, bold=True, color=MID_BLUE)

# ================================================================
# SLIDE 8 — PPM 与混合参数 μ
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '模型 2/2', '种植分块模型（PPM）与混合参数 μ', pnum='8/14')

# Left: PPM explanation
add_rrect(slide, Cm(2), Cm(4), Cm(15), Cm(6), LIGHT_BLUE)
tf = mktf(slide, Cm(3.5), Cm(4.3), Cm(12), Cm(5.5))
addp(tf, 'PPM（Planted Partition Model）', size=14, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=2)
addp(tf, '• SBM 的特例', size=11, color=DARK_TEXT)
addp(tf, '• 所有社区规模相等', size=11, color=DARK_TEXT)
addp(tf, '• 只有两种连接概率：', size=11, color=DARK_TEXT)
addp(tf, '  p_in (内部) 和 p_out (外部)', size=11, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '社区检测目标：', size=12, bold=True, color=ACCENT_ORANGE)
addp(tf, '基于生成的网络结构恢复', size=11, color=DARK_TEXT)
addp(tf, '节点的社区归属', size=11, color=DARK_TEXT)

# Right: PPM image
add_img(slide, 'slide09_img02.jpg', Cm(18), Cm(4), Cm(14), Cm(6))

# Mixing parameter section
add_rrect(slide, Cm(2), Cm(10.5), Cm(30), Cm(7), SOFT_BG)
add_tb(slide, Cm(3.5), Cm(10.8), Cm(15), Cm(1.2), '混合参数 μ = np_out / ⟨k⟩', size=16, bold=True, color=DARK_NAVY)

tf = mktf(slide, Cm(3.5), Cm(12.2), Cm(27), Cm(2))
addp(tf, 'μ 衡量社区之间混合的程度：μ → 0 时社区分离良好（p_out ≪ p_in）；', size=11, color=DARK_TEXT, first=True)
addp(tf, 'μ → 1 时无社区结构（p_out = p_in，退化为 ER 随机图）。', size=11, color=DARK_TEXT)
addp(tf, '社区存在于整个 μ ∈ [0,1) 区间，但检测难度随 μ 增加而增加。', size=11, color=DARK_TEXT)

# Detectability spectrum
addp(tf, '', size=4)
# Three regions visually
for i, (label, desc, color) in enumerate([
    ('μ 小 → 容易', '社区分离良好\n算法可完美恢复', GREEN),
    ('μ ≈ μ*', '信息论极限\n检测能力的边界', ACCENT_ORANGE),
    ('μ 大 → 不可检测', '无法优于随机猜测\nμ* ≤ μ < 1', RED),
]):
    x = Cm(3.5 + i * 9.5)
    add_rrect(slide, x, Cm(14.5), Cm(8.5), Cm(2.5), color)
    add_tb(slide, x+Cm(0.3), Cm(14.5), Cm(7.9), Cm(1), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(slide, x+Cm(0.3), Cm(15.3), Cm(7.9), Cm(1.5), desc, size=9, color=WHITE, align=PP_ALIGN.CENTER, lsp=1.2)

# ================================================================
# SLIDE 9 — 信息论可检测极限
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '理论基础', '信息论可检测性极限', pnum='9/14')

# Explanation
tf = mktf(slide, Cm(2), Cm(4), Cm(30), Cm(3))
addp(tf, '核心概念：在 SBM 生成的图中，存在一个"信息论极限" μ*：', size=13, color=DARK_TEXT, first=True)
addp(tf, '当 μ < μ* 时，有算法能以优于随机猜测的方式恢复社区（可检测）。', size=12, color=DARK_TEXT, sa=4)
addp(tf, '当 μ* ≤ μ < 1 时，没有任何算法能优于随机猜测（不可检测）。', size=12, color=DARK_TEXT)
addp(tf, '这个极限是由图的随机结构本身决定的，与具体算法无关。', size=12, color=DARK_TEXT)

# Visual spectrum
add_rrect(slide, Cm(2), Cm(8), Cm(30), Cm(4), LIGHT_BLUE)

# μ axis line
add_rect(slide, Cm(4), Cm(9.5), Cm(26), Cm(0.08), DARK_NAVY)

# μ=0
add_tb(slide, Cm(3.5), Cm(9.8), Cm(1.5), Cm(1), '0', size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)
# μ=μ*
add_rect(slide, Cm(13.5), Cm(9.2), Cm(0.08), Cm(0.6), ACCENT_ORANGE)
add_tb(slide, Cm(10.5), Cm(9.8), Cm(4), Cm(1), 'μ*', size=12, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
add_tb(slide, Cm(9.5), Cm(10.8), Cm(6), Cm(1), '极限', size=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
# μ=1
add_tb(slide, Cm(28), Cm(9.8), Cm(2), Cm(1), '1', size=11, color=GREY_TEXT, align=PP_ALIGN.CENTER)

# Labels
add_tb(slide, Cm(4.5), Cm(8.3), Cm(9), Cm(1.2), '可检测 ✓', size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_tb(slide, Cm(14.5), Cm(8.3), Cm(13.5), Cm(1.2), '不可检测 ✗', size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Arrow labels
add_tb(slide, Cm(4.5), Cm(10.5), Cm(9), Cm(1.5), '算法优于随机猜测', size=9, color=GREEN, align=PP_ALIGN.CENTER)
add_tb(slide, Cm(14.5), Cm(10.5), Cm(13.5), Cm(1.5), '任何算法都不优于随机猜测', size=9, color=RED, align=PP_ALIGN.CENTER)

# Key insight
add_rrect(slide, Cm(2), Cm(12.5), Cm(30), Cm(5), CARD_BG)
tf = mktf(slide, Cm(3.5), Cm(12.8), Cm(27), Cm(4.5))
addp(tf, '为什么这对本文重要？', size=14, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=3)
addp(tf, '本文要证明的正是：node2vec 的嵌入可以达到这个极限 μ*。', size=12, bold=True, color=MID_BLUE)
addp(tf, '', size=3)
addp(tf, '这意味着：', size=11, color=DARK_TEXT)
addp(tf, '1. node2vec 在社区检测任务上是"信息论最优"的——理论上不会比任何其他方法差。', size=11, color=DARK_TEXT)
addp(tf, '2. 深层网络和非线性激活不是达到最优所必需的——浅层线性网络就已经足够。', size=11, color=DARK_TEXT)
addp(tf, '3. 嵌入 + K-means 聚类的两步法，具有坚实的理论基础。', size=11, color=DARK_TEXT)

# ================================================================
# SLIDE 10 — 主要结论
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, 0, 0, W, Cm(0.15), ACCENT_ORANGE)

add_tb(slide, Cm(2), Cm(1.2), Cm(5), Cm(1), '核心结论', size=10, bold=True, color=ACCENT_ORANGE)
add_tb(slide, Cm(2), Cm(2), Cm(30), Cm(2), '本文的主要理论贡献', size=24, bold=True, color=WHITE)
add_rect(slide, Cm(2), Cm(4), Cm(4), Cm(0.04), ACCENT_ORANGE)

# Three findings
findings = [
    ('定理 1', 'node2vec 的嵌入可达到 SBM 生成图上\n社区检测的信息论极限。', '即：理论上最优，不会低于任何其他方法。'),
    ('定理 2', '深层网络和非线性激活\n并非实现最优社区检测所必需的。', '浅层线性神经网络就已足够。\n挑战了"越深越好"的直觉。'),
    ('核心洞见', 'node2vec 嵌入与对称归一化拉普拉斯\n矩阵的谱嵌入之间存在等价性。', '解释了为什么浅层方法也能\n捕捉深层社区结构。'),
]

for i, (title, main, note) in enumerate(findings):
    y = Cm(5 + i * 4)
    add_rect(slide, Cm(2), y, Cm(0.12), Cm(2.5), ACCENT_ORANGE)
    add_tb(slide, Cm(2.8), y, Cm(10), Cm(1), title, size=14, bold=True, color=ACCENT_ORANGE)
    tf = mktf(slide, Cm(2.8), y+Cm(1.2), Cm(28), Cm(2.5))
    addp(tf, main, size=12, color=WHITE, lsp=1.5, first=True)

# Proof image
add_img(slide, 'slide08_img02.jpg', Cm(2), Cm(16.2), Cm(30), Cm(1.8))

add_rect(slide, 0, Cm(18.5), W, Cm(0.5), SOFT_BG)
add_tb(slide, Cm(1.5), Cm(18.65), Cm(25), Cm(0.35),
       '核心洞见：深度学习的"深度"不是必须的', size=7, color=GREY_TEXT)
add_tb(slide, Cm(31), Cm(18.6), Cm(2), Cm(0.35), '10/14', size=7, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ================================================================
# SLIDE 11 — 证明思路
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '证明思路', 'node2vec 嵌入与谱方法的等价性', pnum='11/14')

# Three steps
steps = [
    ('Step 1', '嵌入 = 矩阵分解',
     'node2vec 的损失函数在无限采样\n极限下收敛到确定的矩阵分解问题。\n\n这意味着随机游走采样的随机性\n不会影响嵌入的最终几何结构。'),
    ('Step 2', '矩阵 = 拉普拉斯谱',
     '该矩阵与对称归一化拉普拉斯矩阵\n的谱分解等价。\n\nnode2vec 学习到的嵌入与拉普拉斯\n矩阵的 d 个最大特征向量一致。'),
    ('Step 3', '谱方法 = 信息论最优',
     '谱聚类在 SBM 上已知可以达到\n信息论可检测极限。\n\n因此 node2vec 继承了这个最优性。\n等价性 = 最优性。'),
]

for i, (title, summary, detail) in enumerate(steps):
    y = Cm(4.5 + i * 4.3)
    # Number
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2), y+Cm(0.2), Cm(1.2), Cm(1.2))
    c.fill.solid(); c.fill.fore_color.rgb = MID_BLUE; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.text = str(i+1)
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Summary box
    add_rrect(slide, Cm(4), y, Cm(13), Cm(3.5), LIGHT_BLUE)
    add_tb(slide, Cm(5), y+Cm(0.2), Cm(11), Cm(1), title, size=13, bold=True, color=MID_BLUE)
    add_tb(slide, Cm(5), y+Cm(1.2), Cm(11), Cm(2), summary, size=11, color=DARK_TEXT, lsp=1.4)

    # Detail box
    add_rrect(slide, Cm(17.5), y, Cm(14.5), Cm(3.5), SOFT_BG)
    add_tb(slide, Cm(18.5), y+Cm(0.3), Cm(12.5), Cm(3), detail, size=10, color=GREY_TEXT, lsp=1.4)

    if i < len(steps) - 1:
        add_tb(slide, Cm(2), y+Cm(3.7), Cm(3), Cm(0.8), '⬇', size=14, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 12 — 谱方法详解 (NEW)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '深入理解', '为什么是谱方法？', pnum='12/14')

# Left: what is spectral method
add_rrect(slide, Cm(2), Cm(4), Cm(14.5), Cm(6.5), LIGHT_BLUE)
tf = mktf(slide, Cm(3.5), Cm(4.3), Cm(11.5), Cm(6))
addp(tf, '谱聚类（Spectral Clustering）', size=15, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=3)
addp(tf, '1. 构建图的拉普拉斯矩阵 L', size=12, color=DARK_TEXT)
addp(tf, '   L = D − A (D: 度矩阵, A: 邻接矩阵)', size=11, color=DARK_TEXT)
addp(tf, '2. 计算 L 的 d 个最小特征向量', size=12, color=DARK_TEXT)
addp(tf, '3. 用特征向量作为节点的新表示', size=12, color=DARK_TEXT)
addp(tf, '4. 在新表示上运行 K-means', size=12, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '谱方法是社区检测中', size=11, color=GREY_TEXT)
addp(tf, '理论最成熟的方法之一。', size=11, color=GREY_TEXT)

# Right: equivalence
add_rrect(slide, Cm(17.5), Cm(4), Cm(14.5), Cm(6.5), SOFT_BG)
tf = mktf(slide, Cm(19), Cm(4.3), Cm(11.5), Cm(6))
addp(tf, '等价性为什么成立？', size=15, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=3)
addp(tf, 'node2vec 的 Skip-gram 目标函数', size=12, bold=True, color=MID_BLUE)
addp(tf, '本质上是在分解一个矩阵 M，其中', size=11, color=DARK_TEXT)
addp(tf, 'M_ij = log(邻接矩阵的某种变换)', size=11, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '这个 M 与归一化拉普拉斯矩阵', size=11, color=DARK_TEXT)
addp(tf, 'L_sym = D⁻¹ʹ² A D⁻¹ʹ²', size=12, bold=True, color=DARK_TEXT)
addp(tf, '的特征空间是一致的。', size=11, color=DARK_TEXT)

# Bottom insight
add_rrect(slide, Cm(2), Cm(11), Cm(30), Cm(6.5), CARD_BG)
tf = mktf(slide, Cm(3.5), Cm(11.3), Cm(27), Cm(6))
addp(tf, '直观理解', size=15, bold=True, color=DARK_NAVY, first=True)
addp(tf, '', size=3)
addp(tf, '谱方法的本质：利用"拉普拉斯矩阵的特征向量"来揭示图的聚类结构。', size=12, color=DARK_TEXT)
addp(tf, 'node2vec 的本质：利用"随机游走 + Skip-gram"来学习节点的向量表示。', size=12, color=DARK_TEXT)
addp(tf, '', size=3)
addp(tf, '本文证明：这两条路径最终到达同一个数学结构。', size=13, bold=True, color=ACCENT_ORANGE)
addp(tf, '就像两个不同方向的登山队，最终登上了同一座山峰。', size=11, color=GREY_TEXT)
addp(tf, '', size=6)
addp(tf, '这意味着：node2vec 的"黑箱"其实并不是真的黑——它的行为可以用谱图理论精确解释。', size=12, bold=True, color=MID_BLUE)

# ================================================================
# SLIDE 13 — 数值验证
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '实验验证', '数值测试与观察', pnum='13/14')

# Two cards
add_card(slide, Cm(2), Cm(4), Cm(14.5), Cm(6), LIGHT_BLUE, '实验结果',
         ['node2vec + K-means 达到接近最优置信传播(BP)方法的精度',
          '在度异质性（节点度数差异大）下仍能学习社区',
          '社区规模异质性下表现受限——可能源于 K-means 缺陷',
          'K-means 在聚类规模差异大时存在困难'])

add_card(slide, Cm(17.5), Cm(4), Cm(14.5), Cm(6), SOFT_BG, '开放问题',
         ['线性层是否可以被多层非线性激活超越？',
          '何时"深度"是必要的？',
          '更复杂生成模型（如度校正 SBM）下的理论保证？',
          '扩展到重叠社区、带权网络？'])

# Key findings
add_rrect(slide, Cm(2), Cm(10.5), Cm(30), Cm(7), CARD_BG)
add_tb(slide, Cm(3.5), Cm(10.8), Cm(27), Cm(1.5), '实验设计的精妙之处', size=16, bold=True, color=DARK_NAVY)

tf = mktf(slide, Cm(3.5), Cm(12.5), Cm(27), Cm(4.5))
addp(tf, '论文的数值实验刻意设计了多组对比：', size=12, color=DARK_TEXT, first=True)
addp(tf, '', size=3)
addp(tf, '1. 改变混合参数 μ，观察从"容易"到"不可检测"的完整过渡', size=11, color=DARK_TEXT)
addp(tf, '2. 对比 node2vec + K-means 与最优 BP 方法', size=11, color=DARK_TEXT)
addp(tf, '3. 测试度异质性和社区规模异质性下的鲁棒性', size=11, color=DARK_TEXT)
addp(tf, '', size=4)
addp(tf, '结果清晰地落在理论预测的曲线上——实证与理论完美吻合。', size=12, bold=True, color=MID_BLUE)

# ================================================================
# SLIDE 14 — 总结
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, '总结', '回顾与启示', pnum='14/14')

# Summary cards
add_card(slide, Cm(2), Cm(4), Cm(14.5), Cm(6), LIGHT_BLUE, '理论贡献',
         ['证明了 node2vec 嵌入可达到 SBM 社区检测的信息论极限',
          '揭示了浅层线性网络已经足够——深度和激活函数不是必需的',
          '建立了 node2vec 与谱方法之间的等价性'],
         item_size=10)

add_card(slide, Cm(17.5), Cm(4), Cm(14.5), Cm(6), SOFT_BG, '实践启示',
         ['图嵌入方法有坚实的理论基础，不是经验黑箱',
          '选择方法时，"对"比"深"更重要',
          '嵌入 + 聚类是端到端可解释的社区检测范式'],
         item_size=10)

# Takeaway
add_rrect(slide, Cm(2), Cm(10.5), Cm(30), Cm(7), DARK_NAVY)
add_tb(slide, Cm(3.5), Cm(10.8), Cm(27), Cm(1.5), '核心启迪', size=18, bold=True, color=ACCENT_ORANGE)

tf = mktf(slide, Cm(3.5), Cm(12.8), Cm(27), Cm(4))
addp(tf, '多一层"深"不如多一分"懂"', size=16, bold=True, color=WHITE, lsp=1.5, first=True)
addp(tf, '', size=4)
addp(tf, '本文的核心价值不仅是证明了一个定理，更是打开了一扇窗——', size=12, color=RGBColor(0xCC, 0xDD, 0xEE))
addp(tf, '让我们看清神经网络"黑箱"内部的结构。', size=12, color=RGBColor(0xCC, 0xDD, 0xEE))
addp(tf, '', size=6)
addp(tf, '理解方法的理论基础 → 在实践中做出最佳选择', size=13, bold=True, color=WHITE)
addp(tf, '今天的理论理解，就是明天的实践优势。', size=11, color=RGBColor(0xAA, 0xBB, 0xCC))

# ================================================================
# SAVE
# ================================================================
output_path = '/Users/cyingfang/Documents/陈于东/第一部分_优化版.pptx'
prs.save(output_path)
print(f'Saved: {output_path} | Slides: {len(prs.slides)}')
