#!/usr/bin/env python3
"""陈颖芳个人介绍 PPT — 22 slides"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
GREEN = RGBColor(0x00, 0x96, 0x53)
NAVY = RGBColor(0x17, 0x2B, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT = RGBColor(0xFF, 0x6B, 0x35)
GOLD = RGBColor(0xD4, 0xA0, 0x1F)

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
SW = prs.slide_width
SH = prs.slide_height

def rect(s, l, t, w, h, fill=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    if fill:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    else:
        r.fill.background()
    r.line.fill.background()
    return r

def txt(s, l, t, w, h, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tb

def mtext(s, l, t, w, h, lines):
    tb = s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.LEFT
    return tb

def bar(s, l, t, w, color=GREEN, h=4):
    return rect(s, l, t, w, h, color)

def pn(s, n):
    txt(s, SW - Emu(800000), SH - Emu(400000), Emu(600000), Emu(300000),
        str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def section(s, num, title, sub=''):
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, SH - Emu(50000), SW, Emu(50000), GREEN)
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Emu(400000), Emu(1400000), Emu(350000), Emu(350000))
    tri.fill.solid()
    tri.fill.fore_color.rgb = GREEN
    tri.line.fill.background()
    tri.rotation = -15.0
    txt(s, Emu(900000), Emu(1800000), Emu(10000000), Emu(1000000), title, size=38, color=WHITE, bold=True)
    if sub:
        txt(s, Emu(900000), Emu(3000000), Emu(10000000), Emu(600000), sub, size=15, color=RGBColor(0xBB, 0xBB, 0xBB))
    txt(s, Emu(400000), Emu(5800000), Emu(3000000), Emu(300000), f'SECTION {num}', size=13, color=GREEN, bold=True)

def timeline_item(s, l, t, w, period, org, role, highlight='', color=GREEN):
    bar(s, l, t, Emu(50000), color, Emu(600000))
    txt(s, l + Emu(120000), t + Emu(30000), Emu(2000000), Emu(250000),
        period, size=11, color=color, bold=True)
    txt(s, l + Emu(2200000), t + Emu(30000), Emu(w - Emu(2320000)), Emu(250000),
        org, size=11, color=NAVY, bold=True)
    txt(s, l + Emu(2200000), t + Emu(280000), Emu(w - Emu(2320000)), Emu(280000),
        role, size=10, color=DARK)
    if highlight:
        txt(s, l + Emu(120000), t + Emu(330000), Emu(w - Emu(240000)), Emu(250000),
            highlight, size=9, color=GRAY)

def skill_bar(s, l, t, w, label, level, color=GREEN):
    txt(s, l, t, Emu(1500000), Emu(280000), label, size=11, color=NAVY, bold=True)
    rect(s, l + Emu(1600000), t + Emu(60000), Emu(8000000), Emu(160000), LGRAY)
    rect(s, l + Emu(1600000), t + Emu(60000), Emu(int(8000000 * level)), Emu(160000), color)
    txt(s, l + Emu(9800000), t, Emu(1000000), Emu(280000), f'{int(level * 100)}%', size=10, color=color)

def info_card(s, l, t, w, h, title, items, color=GREEN):
    rect(s, l, t, w, h, WHITE)
    bar(s, l, t, w, color, 4)
    txt(s, l + Emu(80000), t + Emu(120000), w - Emu(160000), Emu(300000),
        title, size=14, color=color, bold=True)
    y = t + Emu(450000)
    for item in items:
        txt(s, l + Emu(100000), y, w - Emu(200000), Emu(260000),
            f'> {item}', size=10, color=DARK)
        y += Emu(280000)


# ================================================================
# SLIDE 1: COVER
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, SH - Emu(70000), SW, Emu(70000), GREEN)
for x, y, sz, rot in [
    (Emu(700000), Emu(1200000), Emu(500000), 0.0),
    (Emu(9000000), Emu(4200000), Emu(400000), 45.0),
    (Emu(1500000), Emu(5000000), Emu(250000), -30.0),
]:
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, sz, sz)
    tri.fill.solid()
    tri.fill.fore_color.rgb = GREEN
    tri.line.fill.background()
    tri.rotation = rot

mtext(s, Emu(1000000), Emu(1800000), Emu(10000000), Emu(2200000), [
    ('陈颖芳', 48, WHITE, True),
    ('CHEN YINGFANG', 18, RGBColor(0x88, 0x88, 0x88), False),
    ('', 12, WHITE, False),
    ('邮储银行上海分行 科技金融事业部 高级副经理', 20, GREEN, False),
    ('资深对公业务专家 | CISA / CIA | 中国科学技术大学本科', 13, RGBColor(0xAA, 0xAA, 0xAA), False),
])
txt(s, Emu(1000000), Emu(5200000), Emu(5000000), Emu(300000),
    '2026年5月 | 个人介绍', size=11, color=RGBColor(0x88, 0x88, 0x88))
pn(s, 1)
print("Slide 1: Cover")

# ================================================================
# SLIDE 2: AGENDA
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(600000), Emu(180000), Emu(10000000), Emu(500000),
    '目录 CONTENTS', size=28, color=NAVY, bold=True)
bar(s, Emu(600000), Emu(800000), Emu(3500000), NAVY, 3)

agenda = [
    ('01', '个人档案', '基本信息与身份定位'),
    ('02', '教育背景', '中国科学技术大学 双学位'),
    ('03', '职业履历', '20+年金融IT与银行业经验'),
    ('04', '现任职务', '邮储银行科技金融事业部'),
    ('05', '核心能力', 'IT治理·审计·科技金融'),
    ('06', '主要成就', '获奖与行业认可'),
    ('07', '同业网络', '科大校友金融人脉'),
    ('08', '阅读与知识体系', '14条阅读主线·475本馆藏'),
]
y = Emu(1100000)
for num, title, desc in agenda:
    txt(s, Emu(600000), y, Emu(600000), Emu(350000), num, size=22, color=GREEN, bold=True)
    txt(s, Emu(1300000), y, Emu(9000000), Emu(350000), title, size=15, color=NAVY, bold=True)
    txt(s, Emu(1300000), y + Emu(300000), Emu(9000000), Emu(250000), desc, size=10, color=GRAY)
    y += Emu(680000)

pn(s, 2)
print("Slide 2: Agenda")

# ================================================================
# SLIDE 3: Personal Profile (section)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 1, '个人档案', 'Personal Profile')
pn(s, 3)
print("Slide 3: Profile section")

# SLIDE 4: Profile details
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '个人档案', size=26, color=NAVY, bold=True)

profile_items = [
    ('姓名', '陈颖芳', '籍贯', '浙江'),
    ('出生年月', '1977年2月', '常住地', '上海'),
    ('政治面貌', '中国民主建国会（民建）', '学历', '中国科学技术大学 本科（双学位）'),
    ('现任职务', '邮储银行上海分行 科技金融事业部 高级副经理', '', ''),
    ('专业认证', 'CISA（国际注册信息系统审计师）/ CIA（国际注册内部审计师）', '', ''),
]
y = Emu(800000)
for left_label, left_val, right_label, right_val in profile_items:
    txt(s, Emu(500000), y, Emu(1500000), Emu(280000), left_label, size=12, color=GRAY)
    txt(s, Emu(2100000), y, Emu(3500000), Emu(280000), left_val, size=13, color=NAVY, bold=True)
    if right_label:
        txt(s, Emu(6200000), y, Emu(1500000), Emu(280000), right_label, size=12, color=GRAY)
        txt(s, Emu(7800000), y, Emu(4000000), Emu(280000), right_val, size=13, color=NAVY, bold=True)
    y += Emu(500000)

y += Emu(300000)
bar(s, Emu(500000), y, Emu(11000000), NAVY, 2)
y += Emu(200000)
txt(s, Emu(500000), y + Emu(100000), Emu(11000000), Emu(350000),
    '个人定位', size=16, color=GREEN, bold=True)
y += Emu(450000)
bullets = [
    '20年银行IT研发、治理、审计全链路经验，从程序员到部门总经理的完整成长路径',
    '科技金融业务专家：深度研究长三角科技金融拓客增户、FPA增长策略',
    '坚持系统阅读：14条阅读主线、475本馆藏学术经典、持续的跨领域知识整合',
    '工作方法论：要事为先、以终为始、化繁为简',
]
for b in bullets:
    txt(s, Emu(600000), y, Emu(10800000), Emu(260000), f'> {b}', size=11, color=DARK)
    y += Emu(300000)

pn(s, 4)
print("Slide 4: Profile")

# ================================================================
# SLIDE 5: Education section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 2, '教育背景', 'Education — 中国科学技术大学')
pn(s, 5)
print("Slide 5: Education section")

# SLIDE 6: Education detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '教育背景', size=26, color=NAVY, bold=True)

rect(s, Emu(500000), Emu(800000), Emu(11000000), Emu(2000000), LGRAY)
bar(s, Emu(500000), Emu(800000), Emu(60000), GREEN, Emu(2000000))
txt(s, Emu(700000), Emu(880000), Emu(5000000), Emu(350000),
    '中国科学技术大学（USTC / 985·211）', size=20, color=NAVY, bold=True)
txt(s, Emu(700000), Emu(1300000), Emu(8000000), Emu(250000),
    '班级编号：9509/959 | 1995年入学', size=12, color=GRAY)
txt(s, Emu(700000), Emu(1600000), Emu(4000000), Emu(250000),
    '测控技术与仪器  工学学士（1995-2000）', size=14, color=GREEN, bold=True)
txt(s, Emu(700000), Emu(1900000), Emu(4000000), Emu(250000),
    '计算机科学与技术  第二学士学位（1998-2000）', size=14, color=GREEN, bold=True)

y = Emu(3100000)
txt(s, Emu(500000), y, Emu(5000000), Emu(300000),
    '持续学习计划', size=18, color=GREEN, bold=True)
y += Emu(450000)
for item in [
    'CFA / CPA 专业资格认证（规划中）',
    'EMBA 高级管理人员工商管理硕士（规划中）',
    '华东政法大学 法学硕士（规划中）',
    '机器学习、数据结构与算法研究生课程',
    '经济学家或会计师职称（规划中）',
]:
    txt(s, Emu(600000), y, Emu(10000000), Emu(260000), f'> {item}', size=11, color=DARK)
    y += Emu(300000)

txt(s, Emu(500000), Emu(5300000), Emu(11000000), Emu(300000),
    '科大上海金融校友会核心成员 | 2025年科大上海金融校友年会志愿者兼参会者（"樱花大道桌"）',
    size=11, color=GRAY)

pn(s, 6)
print("Slide 6: Education")

# ================================================================
# SLIDE 7: Career Timeline section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 3, '职业履历', 'Career — 20+年金融IT与银行业深耕')
pn(s, 7)
print("Slide 7: Career section")

# SLIDE 8-9: Career detail
for slide_num, careers in [
    (8, [
        ('2022-至今', '邮储银行上海分行', '科技金融事业部 高级副经理', '科技金融拓客增户·FPA增长·对公业务审批'),
        ('2020-2022', '上海华瑞银行', '审计部总经理 / 监事会办公室主任', '银保监会核准审计部负责人身份'),
        ('2020.2-7', '中美联泰大都会', '内审高级经理', '业务连续性/IT/数据治理审计'),
        ('2015-2019', '东亚银行（中国）', '内审部 高级内审经理', 'IT审计+业务审计+风险评估+数据分析'),
    ]),
    (9, [
        ('2010-2015', '东亚银行（中国）', '科技部 架构师/IT风控合规经理', '3年IT治理规划；CBRC三等奖'),
        ('2008-2010', '东亚银行（中国）', '科技部 需求分析与项目管理经理', '新一代零售信贷/影像档案/ECIF建设'),
        ('2005-2008', '银联数据服务', '研发部 开发室经理', '银行卡新产品开发管理'),
        ('2000-2005', '神州数码（上海）', '金融事业部 项目经理→产品经理', '银行核心系统开发与实施'),
    ]),
]:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, WHITE)
    bar(s, 0, 0, SW, GREEN, 5)
    txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
        f'职业履历（{"上半部分" if slide_num == 8 else "下半部分"}）', size=26, color=NAVY, bold=True)

    for i, (period, org, role, highlight) in enumerate(careers):
        y = Emu(800000) + i * Emu(1400000)
        colors = [GREEN, NAVY, ACCENT, GOLD]
        timeline_item(s, Emu(500000), y, Emu(11000000), period, org, role, highlight, colors[i])

    pn(s, slide_num)
    print(f"Slide {slide_num}: Career")

# ================================================================
# SLIDE 10: Current Role section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 4, '现任职务', '邮储银行上海分行 科技金融事业部')
pn(s, 10)
print("Slide 10: Role section")

# SLIDE 11: PSBC role detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '现任职务：邮储银行科技金融事业部', size=26, color=NAVY, bold=True)

info_card(s, Emu(500000), Emu(750000), Emu(5500000), Emu(3200000),
    '核心职责', [
        '研究同业科技金融成功实践，制定可复制策略',
        '定期研究上海科技金融业务机会、拓客策略、FPA增长',
        '制作科技金融合作方案与演示材料',
        '管理拓客清单与授信审批',
    ], GREEN)

info_card(s, Emu(6300000), Emu(750000), Emu(5500000), Emu(3200000),
    '战略项目', [
        '"U益创"全生命周期科技金融服务体系',
        '首单债转股（20亿元）落地',
        '长三角科技金融拓客增户市场调研',
        '科技贷款余额突破1万亿元，服务11万+科技企业',
    ], NAVY)

y = Emu(4300000)
txt(s, Emu(500000), y, Emu(5000000), Emu(300000),
    '邮储银行关键同僚网络', size=16, color=GREEN, bold=True)
y += Emu(450000)
colleagues = [
    ('梁世栋', '首席风险官，科大管院博士', GREEN),
    ('刘建军', '执行董事、行长', NAVY),
    ('牛新庄', '首席信息官', ACCENT),
    ('敬宗泉', '上海分行副行长代表', GOLD),
]
for i, (name, role, color) in enumerate(colleagues):
    x = Emu(500000) + (i % 2) * Emu(6000000)
    row_y = y + (i // 2) * Emu(380000)
    txt(s, x, row_y, Emu(1500000), Emu(300000), name, size=13, color=color, bold=True)
    txt(s, x + Emu(1600000), row_y, Emu(4000000), Emu(300000), role, size=11, color=GRAY)

pn(s, 11)
print("Slide 11: PSBC Role")

# ================================================================
# SLIDE 12: Competencies section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 5, '核心能力', 'Core Competencies')
pn(s, 12)
print("Slide 12: Competencies section")

# SLIDE 13: Competencies detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '核心能力矩阵', size=26, color=NAVY, bold=True)

y = Emu(900000)
skills_data = [
    ('IT治理与风险管理', 0.95, GREEN),
    ('内部审计与合规', 0.95, GREEN),
    ('科技金融业务', 0.85, NAVY),
    ('项目与团队管理', 0.90, NAVY),
    ('数据分析与建模', 0.80, ACCENT),
    ('对公业务审批', 0.85, ACCENT),
    ('监管沟通与报告', 0.90, GOLD),
    ('跨部门协同', 0.85, GOLD),
]
for label, level, color in skills_data:
    skill_bar(s, Emu(500000), y, Emu(11000000), label, level, color)
    y += Emu(420000)

y += Emu(100000)
txt(s, Emu(500000), y, Emu(5000000), Emu(300000),
    '核心工作方法论', size=16, color=GREEN, bold=True)
y += Emu(400000)
for m in [
    '"要事为先、以终为始、化繁为简" — 高效执行框架',
    '系统化研究→策略制定→可复制方案→持续迭代',
    '跨领域知识整合：IT+金融+法律+管理 复合视角',
]:
    txt(s, Emu(600000), y, Emu(10500000), Emu(260000), f'> {m}', size=11, color=DARK)
    y += Emu(300000)

pn(s, 13)
print("Slide 13: Competencies")

# ================================================================
# SLIDE 14: Achievements section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 6, '主要成就', 'Achievements & Awards')
pn(s, 14)
print("Slide 14: Achievements section")

# SLIDE 15: Achievements detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '主要成就与获奖', size=26, color=NAVY, bold=True)

achievements = [
    ('2014', 'CBRC信息科技风险管理研究三等奖', '唯一及首家获此奖项的外资银行', GREEN),
    ('2013', '上海银监局"信息科技管理优胜奖"', '协助完成"两地三中心"灾备体系建设', NAVY),
    ('2012-2015', 'IT治理体系构建', '编制和审核60余项IT政策制度，识别80余项IT固有风险', ACCENT),
    ('2024-2026', '邮储科技金融业务突破', '助力上海分行科技贷款余额突破1万亿元，服务11万+科技企业', GOLD),
    ('2020', '银保监会核准', '审计部负责人任职资格核准', RGBColor(0x7B, 0x1F, 0xA2)),
]
for i, (year, title, desc, color) in enumerate(achievements):
    y = Emu(800000) + i * Emu(1100000)
    bar(s, Emu(500000), y, Emu(60000), color, Emu(900000))
    txt(s, Emu(700000), y + Emu(30000), Emu(1800000), Emu(280000), year, size=12, color=color, bold=True)
    txt(s, Emu(2600000), y + Emu(30000), Emu(9000000), Emu(280000), title, size=14, color=NAVY, bold=True)
    txt(s, Emu(2600000), y + Emu(360000), Emu(9000000), Emu(250000), desc, size=10, color=GRAY)

pn(s, 15)
print("Slide 15: Achievements")

# ================================================================
# SLIDE 16: Network section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 7, '同业网络', 'Professional Network — 科大校友金融人脉')
pn(s, 16)
print("Slide 16: Network section")

# SLIDE 17: Network detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '同业网络与人脉', size=26, color=NAVY, bold=True)

info_card(s, Emu(500000), Emu(750000), Emu(5500000), Emu(2400000),
    '科大校友网络', [
        'USTC上海校友会（全球约15万校友）',
        'USTC上海金融校友会核心成员，覆盖金融行业228人',
        '2025年科大上海金融校友年会：志愿者兼参会者',
        '"樱花大道桌"：深度交流平台',
    ], GREEN)

info_card(s, Emu(6300000), Emu(750000), Emu(5500000), Emu(2400000),
    '金融行业人脉', [
        '外资银行同业网络：东亚银行、大都会人寿',
        '民营银行同业网络：华瑞银行',
        '国有银行同业网络：邮储银行总分行体系',
        '央国企网络：银联体系、神州数码',
    ], NAVY)

y = Emu(3500000)
txt(s, Emu(500000), y, Emu(5000000), Emu(300000),
    '外部关系管理', size=16, color=GREEN, bold=True)
y += Emu(400000)
for item in [
    '与上海金融监管局保持良好沟通渠道',
    '与多家科技企业建立深度合作联系',
    '与同业科技金融部门保持对标交流',
]:
    txt(s, Emu(600000), y, Emu(10500000), Emu(260000), f'> {item}', size=11, color=DARK)
    y += Emu(300000)

pn(s, 17)
print("Slide 17: Network")

# ================================================================
# SLIDE 18: Reading System section
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section(s, 8, '阅读与知识体系', 'Reading System — 14条阅读主线·475本馆藏')
pn(s, 18)
print("Slide 18: Reading section")

# SLIDE 19: Reading detail
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '阅读与知识体系', size=26, color=NAVY, bold=True)

# Main lines summary in cards
reading_cards = [
    ('主线A：自由vs秩序', '已完成', '伯林→哈耶克→弗里德曼→罗尔斯→诺齐克', GREEN),
    ('主线F：西方哲学史', '7/12', '柏拉图→亚里士多德→休谟→康德→黑格尔→马克思', NAVY),
    ('主线J：中国金融体系', '进行中', '陈元《开发性金融》+ 多维度金融改革', ACCENT),
    ('主线K：凯恩斯vs哈耶克', '已完成', '哈耶克三书→凯恩斯通论（进行中）', GOLD),
    ('主线I：科技金融', '进行中', '科技金融新政+硅谷银行模式+风险管理', GREEN),
]
for i, (title, status, desc, color) in enumerate(reading_cards):
    y = Emu(750000) + i * Emu(700000)
    bar(s, Emu(500000), y, Emu(60000), color, Emu(550000))
    txt(s, Emu(700000), y + Emu(30000), Emu(6000000), Emu(250000), title, size=13, color=color, bold=True)
    txt(s, Emu(6800000), y + Emu(30000), Emu(1200000), Emu(250000), status, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Emu(8100000), y, Emu(3700000), Emu(250000), '', size=10, color=GRAY)
    txt(s, Emu(700000), y + Emu(300000), Emu(4500000), Emu(250000), desc, size=9, color=GRAY)

# Stats row
y = Emu(4400000)
bar(s, Emu(500000), y, Emu(11000000), NAVY, 2)
y += Emu(200000)
stats = [('14', '阅读主线'), ('475', '馆藏图书'), ('25+', '已完成精读报告'), ('5', '核心类别')]
for i, (num, label) in enumerate(stats):
    x = Emu(500000) + i * Emu(2900000)
    txt(s, x, y + Emu(100000), Emu(2500000), Emu(400000), num, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y + Emu(500000), Emu(2500000), Emu(250000), label, size=10, color=GRAY, align=PP_ALIGN.CENTER)

pn(s, 19)
print("Slide 19: Reading")

# ================================================================
# SLIDE 20: Personal Projects
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)
txt(s, Emu(500000), Emu(150000), Emu(10000000), Emu(500000),
    '个人项目与兴趣', size=26, color=NAVY, bold=True)

info_card(s, Emu(500000), Emu(750000), Emu(5500000), Emu(2400000),
    '个人项目', [
        '家族微信小程序（开发中）：信息技术与家族传承的结合',
        '儿子成长PPT（200页+）: 持续20年+成长记录',
        '保研与考研咨询辅导：帮助亲友子女升学规划',
    ], GREEN)

info_card(s, Emu(6300000), Emu(750000), Emu(5500000), Emu(2400000),
    '创作与产出', [
        '科技金融工作笔记（持续积累）',
        '小说创作：历史小说+影视剧本',
        '系统阅读笔记：25+部经典精读报告',
    ], NAVY)

y = Emu(3500000)
txt(s, Emu(500000), y, Emu(5000000), Emu(300000),
    '兴趣领域', size=16, color=GREEN, bold=True)
y += Emu(400000)
interests = [
    '马拉松、徒步等耐力运动',
    '跨领域知识整合：哲学·历史·政法·经济·语言',
    '家族文化建设与信息技术融合',
]
for item in interests:
    txt(s, Emu(600000), y, Emu(10500000), Emu(260000), f'> {item}', size=11, color=DARK)
    y += Emu(300000)

pn(s, 20)
print("Slide 20: Projects")

# ================================================================
# SLIDE 21: Philosophy
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, SH - Emu(70000), SW, Emu(70000), GREEN)

mtext(s, Emu(1200000), Emu(1500000), Emu(10000000), Emu(4000000), [
    ('"要事为先、以终为始、化繁为简"', 30, WHITE, True),
    ('', 16, WHITE, False),
    ('二十余年来，无论技术如何更迭、岗位如何变换，', 16, RGBColor(0xBB, 0xBB, 0xBB), False),
    ('我始终坚持一个理念：', 16, RGBColor(0xBB, 0xBB, 0xBB), False),
    ('在正确的时间，用正确的方法，做正确的事。', 18, GREEN, True),
    ('', 14, WHITE, False),
    ('系统阅读给予我跨领域的视野，', 14, RGBColor(0xAA, 0xAA, 0xAA), False),
    ('工作实践给予我落地的方法，', 14, RGBColor(0xAA, 0xAA, 0xAA), False),
    ('而持续的自我更新，是我对职业生涯的承诺。', 14, RGBColor(0xAA, 0xAA, 0xAA), False),
])

pn(s, 21)
print("Slide 21: Philosophy")

# ================================================================
# SLIDE 22: THANK YOU
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, WHITE)
bar(s, 0, 0, SW, GREEN, 5)

mtext(s, Emu(1500000), Emu(2000000), Emu(9500000), Emu(2500000), [
    ('谢谢！', 48, NAVY, True),
    ('THANK YOU', 18, GRAY, False),
    ('', 14, WHITE, False),
    ('陈颖芳 CHEN YINGFANG', 18, NAVY, True),
    ('邮储银行上海分行 科技金融事业部 高级副经理', 14, GRAY, False),
])

# Contact info
y = Emu(4600000)
contacts = [
    ('上海 · 中国', GREEN),
]
for text, color in contacts:
    txt(s, Emu(1500000), y, Emu(9000000), Emu(300000), text, size=12, color=color, align=PP_ALIGN.CENTER)
    y += Emu(350000)

pn(s, 22)
print("Slide 22: Thank You")

# ================================================================
# SAVE
# ================================================================
OUTPUT = '/Users/cyingfang/Documents/科技金融/陈颖芳个人介绍.pptx'
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
