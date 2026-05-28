#!/usr/bin/env python3
"""Generate English and Japanese versions of 重庆晨兵机械有限公司简介.pptx

Strategy: Join all runs in each shape into one string, apply all translations,
then write back to the first run (clearing others).
"""
from pptx import Presentation

SRC = "/Users/cyingfang/claude/deliverables/career/客户方案/重庆晨兵机械有限公司简介.pptx"
OUT_DIR = "/Users/cyingfang/claude/deliverables/career/客户方案/"

EN_TRANSLATIONS = {
    # Slide 1 - title
    "重庆晨兵机械有限公司": "Chongqing Chenbing Machinery Co., Ltd.",
    "2026年5月": "May 2026",
    # Slide 2 - TOC
    "目 录": "Contents",
    "01 公司概况": "01 Company Profile",
    "02 行业应用": "02 Industry Applications",
    "03 企业资质": "03 Qualifications",
    "04 联系我们": "04 Contact Us",
    # Section headers
    "公司概况": "Company Profile",
    "组织架构": "Organization Structure",
    "行业应用": "Industry Applications",
    "企业资质": "Qualifications",
    "联系我们": "Contact Us",
    # Slide 3 - Company overview
    "重庆晨兵机械有限公司主营业务：液冷不锈钢配件，汽车轻量化部件，半导体，新能源设备结构件，工装夹具，治具等。服务行业：新能源汽车，电池，半导体，AI服务器，医疗器械等。":
    "Main business: liquid-cooling stainless steel components, automotive lightweight parts, semiconductors, new energy equipment structural parts, tooling fixtures, jigs, etc. Industries served: NEVs, batteries, semiconductors, AI servers, medical devices.",
    "公司拥有钻攻机，立加，龙门加工中心，卧式加工中心，五轴加工中心，车床，磨床，激光机，切管机，折弯机，锯切机等制造设备。海克斯康三坐标检测设备，硬度计，膜厚仪，色差仪，粗糙度仪等各类检测仪器。":
    "Equipment: drilling & tapping machines, VMC, gantry MC, horizontal MC, 5-axis MC, lathes, grinders, laser cutters, pipe cutters, bending machines, sawing machines, etc. Inspection: Hexagon CMM, hardness testers, coating thickness gauges, colorimeters, roughness testers.",
    "公司获得国家高新技术企业、重庆市专精特新中小企业、重庆市铜梁区技术创新中心的认定。公司已通过 ISO9001:2015 质量管理体系及 IATF16949:2016 汽车行业质量管理体系认证。公司拥有":
    "Certified: National High-Tech Enterprise, Chongqing SRDI SME, Tongliang District Technology Innovation Center. ISO9001:2015 & IATF16949:2016 certified. The company has",
    "两百名优秀同仁，始终秉持周六保证不休息，周日休息不保证的实干态度。全员勤恳务实，奋勇拼搏，以匠心扎根生产一线，稳稳扛起制造业的脊梁。":
    "over 200 dedicated employees who uphold a pragmatic work ethic. The entire team is diligent, grounded, and committed to manufacturing excellence with craftsmanship at the core.",

    # Product tables - column headers
    "产品名称": "Product",
    "应用领域": "Applications",
    "国外标杆": "Global Benchmark",
    "国内标杆": "Domestic Benchmark",
    "竞争对手": "Competitors",
    "产品展示": "Product Gallery",

    # Slide 5 - Liquid cooling
    "液冷配件": "Liquid Cooling Components",
    "AI 训练集群 / 超算": "AI Training Clusters / HPC",
    "云计算 / 大型 IDC": "Cloud Computing / Large IDCs",
    "储能电站": "Energy Storage Stations",
    "芯片制造": "Chip Manufacturing",
    "高端电子测试设备": "High-End Electronic Test Equipment",
    "动力电池热管理": "Power Battery Thermal Management",

    # Slide 7 - NEV lightweight
    "新能源汽车轻量化结构件": "NEV Lightweight Structural Parts",
    "车身系统铝型材件": "Body System Aluminum Profiles",
    "底盘系统铝型材件": "Chassis System Aluminum Profiles",
    "电池包专属铝型材": "Battery Pack Aluminum Profiles",
    "电驱电控铝型材": "E-Drive & E-Control Aluminum Profiles",
    "座舱内饰铝型材": "Cabin Interior Aluminum Profiles",
    "其他异型工业铝型材": "Other Custom Industrial Aluminum Profiles",

    # Slide 9 - Tooling
    "工装夹治具": "Tooling, Fixtures & Jigs",
    "新能源汽车": "New Energy Vehicles",
    "消费电子": "Consumer Electronics",
    "半导体/集成电路": "Semiconductors / IC",
    "航空航天": "Aerospace",
    "工业自动化": "Industrial Automation",
    "医疗器械": "Medical Devices",

    # Slide 11 - Nitrogen springs
    "氮气弹簧": "Nitrogen Gas Springs",
    "汽车冲压模具": "Automotive Stamping Dies",
    "注塑/压铸模具": "Injection / Die-Casting Molds",
    "车身焊接/总装": "Body Welding / Assembly",
    "底盘/悬架": "Chassis / Suspension",
    "新能源电池": "New Energy Batteries",
    "机床/自动化": "Machine Tools / Automation",

    # Slide 13 - Automation
    "自动化设备及配件": "Automation Equipment & Parts",
    "半导体": "Semiconductors",
    "新能源/锂电/储能": "New Energy / Li-Battery / Energy Storage",
    "汽车制造": "Automotive Manufacturing",
    "通用制造业/加工": "General Manufacturing / Processing",
    "物流仓储": "Logistics & Warehousing",

    # Company names & benchmarks
    "英伟达": "NVIDIA", "华为云/数据中心": "Huawei Cloud / Data Center",
    "宁德时代": "CATL", "台积电": "TSMC", "长江存储": "YMTC", "比亚迪": "BYD",
    "英维克": "Envicool", "高澜股份": "Goaland", "纳百川": "Nabaichuan",
    "中石科技": "Zhongshi Tech", "三花智控": "Sanhua", "银轮股份": "Yinlun",
    "特斯拉": "Tesla", "大众": "Volkswagen", "丰田": "Toyota", "宝马": "BMW",
    "麦格纳": "Magna", "采埃孚": "ZF", "赛力斯": "Seres", "长安汽车": "Changan Auto",
    "吉利": "Geely", "蔚来": "NIO",
    "忠旺集团": "Zhongwang Group", "广东鸿图": "Guangdong Hongtu",
    "南山铝业": "Nanshan Aluminum", "豪美新材": "Haomei New Materials",
    "友升股份": "Yousheng", "亚太科技": "Asia Pacific Tech",
    "西门子": "Siemens", "高通": "Qualcomm", "波音": "Boeing", "ABB": "ABB",
    "飞利浦": "Philips", "富士康": "Foxconn", "长鑫存储": "CXMT",
    "中航工业": "AVIC", "立讯精密": "Luxshare", "欧菲光": "OFILM",
    "博众精工": "Bozhon", "大族激光": "Han's Laser", "利元亨": "LYH",
    "苏州瀚川": "Suzhou Hanchuan", "东莞捷荣": "Dongguan Jierong",
    "深圳吉阳智能": "Shenzhen Jiyang Intelligent", "昆山超秒精密": "Kunshan Chaomiao Precision",
    "奔驰": "Mercedes-Benz", "日本电装": "Denso", "巴斯夫": "BASF",
    "三星新能源": "Samsung SDI", "库卡": "KUKA",
    "银宝山新": "Yinbaoshanxin", "成飞集成": "Chengfei Integration",
    "一汽模具": "FAW Tooling", "佛山精达": "Foshan Jingda",
    "苏州汇川": "Suzhou Inovance",
    "瑞典卡勒": "KALLER (Sweden)", "美国达科": "DADCO (USA)",
    "日本米思米": "MISUMI (Japan)", "湖南邵阳兴达": "Hunan Shaoyang Xingda",
    "东浙江美力科技": "Zhejiang Meili Technology",
    "广东达科": "Guangdong Dako", "深圳铭昊": "Shenzhen Minghao",
    "ASML": "ASML", "发那科": "FANUC", "博世": "Bosch", "德马吉": "DMG MORI",
    "霍尼韦尔": "Honeywell", "长电科技": "JCET", "国轩高科": "Gotion High-Tech",
    "三一重工": "SANY", "顺丰": "SF Express",
    "汇川技术": "Inovance", "中控技术": "SUPCON",
    "科瑞技术": "Coretronic", "赛腾股份": "Saiteng", "博杰股份": "Bojie",

    # Qualification section
    "公司外景及营业执照": "Company Exterior & Business License",
    "生产设备": "Production Equipment",
    "检测设备": "Inspection Equipment",
    "实验室检测能力": "Laboratory Testing Capabilities",
    "主要客户": "Key Customers",
    "体系证书": "System Certifications",
    "专利证书": "Patent Certificates",

    # Equipment tables
    "序号": "No.",
    "设备名称": "Equipment",
    "设备型号": "Model",
    "设备品牌": "Brand",
    "数量": "Qty",
    "加工零件尺寸 （单位：MM）": "Part Dims (mm)",
    "加工精度": "Accuracy",
    "仕兴鸿": "Shi Xing Hong", "德裕顺": "De Yu Shun", "台群": "Tai Qun",
    "宇辰": "Yu Chen", "云南机床": "Yunnan Machine Tool",
    "哈维": "Harvey", "普创": "Pu Chuang", "建德": "Jian De",
    "CNC龙门加工中心": "CNC Gantry MC",
    "CNC加工中心": "CNC MC",
    "卧式加工中心": "Horizontal MC",
    "五轴加工中心": "5-Axis MC",
    "型材加工中心": "Profile MC",
    "车床": "Lathe",
    "车铣复合": "Turn-Mill Compound",
    "侧孔机": "Side Hole Machine",
    "定梁式单头龙门磨": "Fixed-Beam Gantry Grinder",
    "龙门磨": "Gantry Grinder",
    "大水磨": "Surface Grinder",
    "旺磐小磨床": "Bench Grinder",
    "快走丝线切割机床": "Fast Wire EDM",
    "中走丝线切割机床": "Medium Wire EDM",
    "打孔机": "Drilling Machine",
    "高速圆锯机": "High-Speed Circular Saw",
    "铣床": "Milling Machine",
    "钻床": "Bench Drill",
    "拉丝机": "Wire Drawing Machine",
    "电动攻牙机": "Electric Tapping Machine",
    "激光切割机": "Laser Cutter",
    "激光切管机": "Laser Pipe Cutter",
    "数控板料折弯机": "CNC Press Brake",
    "电火花打孔机": "EDM Drilling Machine",
    "二保焊接机": "Dual-Shield Welder",
    "氩弧焊焊接机": "TIG Welder",
    "激光焊接机": "Laser Welder",
    "气体保护焊机": "Gas-Shielded Welder",
    "焊机铝保护焊机": "Aluminum Welder",
    "焊接平台": "Welding Platform",
    "万齐": "Wan Qi", "铭石": "Ming Shi", "万利": "Wan Li",
    "台正": "Tai Zheng", "永升": "Yong Sheng", "智能": "Zhi Neng",
    "宏山": "Hong Shan", "亚威": "Ya Wei", "锐龙": "Rui Long",
    "伊莎": "Yi Sha", "炬亮": "Ju Liang",
    "旺磐": "Wang Pan",
    "永恒擎天": "Yong Heng Qing Tian", "定制": "Custom",
    "宽150*高150": "W150×H150",

    # Inspection table
    "品牌": "Brand",
    "主要型号规格": "Model / Spec",
    "设备行程": "Travel",
    "设备精度": "Accuracy",
    "光纤激光打标机": "Fiber Laser Marker",
    "高度仪": "Height Gauge",
    "万维": "Wan Wei",
    "三坐标测量机": "CMM",
    "海克斯康": "Hexagon",
    "思瑞": "Serein",
    "自动影像仪": "Auto Vision System",
    "硬度计": "Hardness Tester",
    "里氏": "Leeb",
    "洛氏": "Rockwell",
    "涂层测厚仪": "Coating Thickness Gauge",
    "希玛": "Sima",
    "高度规": "Height Gauge",
    "超声波探伤仪": "Ultrasonic Flaw Detector",
    "三丰": "Mitutoyo",

    # Lab capabilities
    "实验室搭载国内外前沿精密检测设备，面向精密零部件、电子元器件、金属材料实现全流程质量管控。":
    "The laboratory features advanced precision testing equipment for full-process quality control of precision parts, electronic components, and metal materials.",
    "几何量测量：依托高精度三坐标测量机，检测工件三维尺寸、形位公差，测量精度高、适配复杂结构件。":
    "Dimensional Measurement: High-precision CMMs for 3D dimensions and geometric tolerances; suitable for complex structural parts.",
    "性能与表面检测：可完成粗糙度、涂层厚度、硬度等指标测试，分析材料及表面特性。":
    "Performance & Surface Testing: Roughness, coating thickness, and hardness testing for material and surface characterization.",
    "无损 & 可靠性测试：检测材料内部缺陷，开展环境适应性试验，验证产品使用可靠性。":
    "NDT & Reliability Testing: Internal defect detection and environmental testing to verify product reliability.",
    "实验室实现微观表面、宏观尺寸、材料性能、环境耐受全场景覆盖，支持产品质量验证、工艺迭代与合规评价。":
    "Full-spectrum coverage from micro-surface to macro-dimensions, material properties, and environmental tolerance — supporting product validation, process improvement, and compliance assessment.",

    # Contact
    "黄 总": "Mr. Huang",
    "重庆市铜梁区蒲吕街道产业大道55号附39号": "No. 55-39, Chanye Ave, Pulu, Tongliang, Chongqing, China",

    # Closing
    "晨启新局，兵拓未来，晨兵智造，精工志远":
    "Forging New Horizons, Expanding the Future — Chenbing Precision Manufacturing, Engineering Excellence",
}

JA_TRANSLATIONS = {
    "重庆晨兵机械有限公司": "重慶晨兵機械有限公司",
    "2026年5月": "2026年5月",
    "目 录": "目 次",
    "01 公司概况": "01 会社概要",
    "02 行业应用": "02 業界応用",
    "03 企业资质": "03 認証・資格",
    "04 联系我们": "04 お問い合わせ",
    "公司概况": "会社概要",
    "组织架构": "組織構成",
    "行业应用": "業界応用",
    "企业资质": "認証・資格",
    "联系我们": "お問い合わせ",

    "重庆晨兵机械有限公司主营业务：液冷不锈钢配件，汽车轻量化部件，半导体，新能源设备结构件，工装夹具，治具等。服务行业：新能源汽车，电池，半导体，AI服务器，医疗器械等。":
    "主事業：液冷ステンレス部品、自動車軽量化部品、半導体、新エネルギー設備構造部品、治工具、ジグ等。対象業界：新エネルギー車、電池、半導体、AIサーバー、医療機器等。",
    "公司拥有钻攻机，立加，龙门加工中心，卧式加工中心，五轴加工中心，车床，磨床，激光机，切管机，折弯机，锯切机等制造设备。海克斯康三坐标检测设备，硬度计，膜厚仪，色差仪，粗糙度仪等各类检测仪器。":
    "設備：ドリルタップ機、立形MC、門形MC、横形MC、5軸MC、旋盤、研削盤、レーザー加工機、パイプ切断機、曲げ機、鋸切断機等。検査：Hexagon三次元測定機、硬度計、膜厚計、色差計、粗さ計等。",
    "公司获得国家高新技术企业、重庆市专精特新中小企业、重庆市铜梁区技术创新中心的认定。公司已通过 ISO9001:2015 质量管理体系及 IATF16949:2016 汽车行业质量管理体系认证。公司拥有":
    "認定：国家高新技術企業、重慶市専精特新中小企業、重慶市銅梁区技術革新センター。ISO9001:2015及びIATF16949:2016認証取得。",
    "两百名优秀同仁，始终秉持周六保证不休息，周日休息不保证的实干态度。全员勤恳务实，奋勇拼搏，以匠心扎根生产一线，稳稳扛起制造业的脊梁。":
    "約200名の優秀な従業員が実直な勤務姿勢でものづくりに励み、匠の精神で製造業を支えています。",

    "产品名称": "製品名", "应用领域": "応用分野",
    "国外标杆": "海外ベンチマーク", "国内标杆": "国内ベンチマーク",
    "竞争对手": "競合他社", "产品展示": "製品展示",

    "液冷配件": "液冷部品",
    "AI 训练集群 / 超算": "AIトレーニングクラスタ / HPC",
    "云计算 / 大型 IDC": "クラウド / 大型IDC",
    "储能电站": "蓄電所",
    "芯片制造": "チップ製造",
    "高端电子测试设备": "先端電子試験装置",
    "动力电池热管理": "動力電池熱管理",

    "新能源汽车轻量化结构件": "NEV軽量化構造部品",
    "车身系统铝型材件": "ボディ用アルミ形材",
    "底盘系统铝型材件": "シャシー用アルミ形材",
    "电池包专属铝型材": "バッテリーパック用アルミ形材",
    "电驱电控铝型材": "電動駆動・制御用アルミ形材",
    "座舱内饰铝型材": "内装用アルミ形材",
    "其他异型工业铝型材": "その他異形工業用アルミ形材",

    "工装夹治具": "治工具・ジグ",
    "新能源汽车": "新エネルギー車",
    "消费电子": "家電・電子",
    "半导体/集成电路": "半導体 / IC",
    "航空航天": "航空宇宙",
    "工业自动化": "FA（工場自動化）",
    "医疗器械": "医療機器",

    "氮气弹簧": "窒素ガススプリング",
    "汽车冲压模具": "自動車プレス金型",
    "注塑/压铸模具": "射出 / ダイカスト金型",
    "车身焊接/总装": "ボディ溶接 / 組立",
    "底盘/悬架": "シャシー / サスペンション",
    "新能源电池": "新エネルギー電池",
    "机床/自动化": "工作機械 / 自動化",

    "自动化设备及配件": "自動化設備・部品",
    "半导体": "半導体",
    "新能源/锂电/储能": "新エネ / リチウム電池 / 蓄電",
    "汽车制造": "自動車製造",
    "通用制造业/加工": "一般製造業 / 加工",
    "物流仓储": "物流・倉庫",

    "英伟达": "NVIDIA", "华为云/数据中心": "Huawei Cloud",
    "宁德时代": "CATL", "台积电": "TSMC", "长江存储": "YMTC", "比亚迪": "BYD",
    "英维克": "Envicool", "高澜股份": "Goaland", "纳百川": "Nabaichuan",
    "中石科技": "Zhongshi", "三花智控": "Sanhua", "银轮股份": "Yinlun",
    "特斯拉": "Tesla", "大众": "VW", "丰田": "Toyota", "宝马": "BMW",
    "麦格纳": "Magna", "采埃孚": "ZF", "赛力斯": "Seres", "长安汽车": "長安汽車",
    "吉利": "Geely", "蔚来": "NIO",
    "忠旺集团": "忠旺G", "广东鸿图": "広東鴻図",
    "南山铝业": "南山Al", "豪美新材": "豪美新材",
    "友升股份": "友升", "亚太科技": "AP Tech",
    "西门子": "Siemens", "高通": "Qualcomm", "波音": "Boeing", "ABB": "ABB",
    "飞利浦": "Philips", "富士康": "Foxconn", "长鑫存储": "CXMT",
    "中航工业": "AVIC", "立讯精密": "Luxshare", "欧菲光": "OFILM",
    "博众精工": "博衆精工", "大族激光": "Han's Laser", "利元亨": "LYH",
    "苏州瀚川": "蘇州瀚川", "东莞捷荣": "東莞捷栄",
    "深圳吉阳智能": "深セン吉陽智能", "昆山超秒精密": "崑山超秒精密",
    "奔驰": "Mercedes-Benz", "日本电装": "Denso", "巴斯夫": "BASF",
    "三星新能源": "Samsung SDI", "库卡": "KUKA",
    "银宝山新": "銀宝山新", "成飞集成": "成飛集成",
    "一汽模具": "第一汽車金型", "佛山精达": "佛山精達",
    "苏州汇川": "蘇州匯川",
    "瑞典卡勒": "KALLER", "美国达科": "DADCO",
    "日本米思米": "MISUMI", "湖南邵阳兴达": "湖南邵陽興達",
    "东浙江美力科技": "浙江美力科技",
    "广东达科": "広東達科", "深圳铭昊": "深セン銘昊",
    "ASML": "ASML", "发那科": "FANUC", "博世": "Bosch", "德马吉": "DMG MORI",
    "霍尼韦尔": "Honeywell", "长电科技": "JCET", "国轩高科": "Gotion",
    "三一重工": "SANY", "顺丰": "SF Express",
    "汇川技术": "Inovance", "中控技术": "SUPCON",
    "科瑞技术": "Coretronic", "赛腾股份": "賽騰", "博杰股份": "博傑",

    "公司外景及营业执照": "会社外観・営業許可証",
    "生产设备": "生産設備",
    "检测设备": "検査設備",
    "实验室检测能力": "試験ラボ能力",
    "主要客户": "主要取引先",
    "体系证书": "認証書",
    "专利证书": "特許証",

    "序号": "No.",
    "设备名称": "設備名",
    "设备型号": "型式",
    "设备品牌": "ブランド",
    "数量": "数量",
    "加工零件尺寸 （单位：MM）": "加工寸法（mm）",
    "加工精度": "精度",
    "仕兴鸿": "仕興鴻", "德裕顺": "德裕順", "台群": "台群",
    "宇辰": "宇辰", "云南机床": "雲南机床",
    "哈维": "Harvey", "普创": "普創", "建德": "建德",
    "CNC龙门加工中心": "CNC門形MC",
    "CNC加工中心": "CNC MC",
    "卧式加工中心": "横形MC",
    "五轴加工中心": "5軸MC",
    "型材加工中心": "形材MC",
    "车床": "旋盤",
    "车铣复合": "複合加工機",
    "侧孔机": "側面孔機",
    "定梁式单头龙门磨": "固定梁門形研削盤",
    "龙门磨": "門形研削盤",
    "大水磨": "平面研削盤",
    "旺磐小磨床": "小型研削盤",
    "快走丝线切割机床": "高速WEDM",
    "中走丝线切割机床": "中速WEDM",
    "打孔机": "穴あけ機",
    "高速圆锯机": "高速丸鋸盤",
    "铣床": "フライス盤",
    "钻床": "ボール盤",
    "拉丝机": "伸線機",
    "电动攻牙机": "電動タップ機",
    "激光切割机": "レーザー切断機",
    "激光切管机": "レーザーパイプ切断機",
    "数控板料折弯机": "CNCプレスブレーキ",
    "电火花打孔机": "放電穴あけ機",
    "二保焊接机": "CO2溶接機",
    "氩弧焊焊接机": "TIG溶接機",
    "激光焊接机": "レーザー溶接機",
    "气体保护焊机": "ガス溶接機",
    "焊机铝保护焊机": "Al溶接機",
    "焊接平台": "溶接定盤",
    "万齐": "万齐", "铭石": "銘石", "万利": "万利",
    "台正": "台正", "永升": "永升", "智能": "智能",
    "宏山": "宏山", "亚威": "亞威", "锐龙": "銳龍",
    "伊莎": "伊莎", "炬亮": "炬亮",
    "永恒擎天": "永恒擎天", "定制": "特注",

    "品牌": "ブランド",
    "主要型号规格": "型式・仕様",
    "设备行程": "ストローク",
    "设备精度": "精度",
    "光纤激光打标机": "ファイバーレーザーマーカー",
    "高度仪": "ハイトゲージ",
    "万维": "万維",
    "三坐标测量机": "三次元測定機",
    "海克斯康": "Hexagon",
    "思瑞": "Serein",
    "自动影像仪": "自動画像測定機",
    "硬度计": "硬度計",
    "里氏": "リーブ",
    "洛氏": "ロックウェル",
    "涂层测厚仪": "膜厚計",
    "希玛": "Sima",
    "高度规": "ハイトゲージ",
    "超声波探伤仪": "超音波探傷装置",
    "三丰": "Mitutoyo",

    "实验室搭载国内外前沿精密检测设备，面向精密零部件、电子元器件、金属材料实现全流程质量管控。":
    "国内外の先端精密測定機器を備え、精密部品・電子部品・金属材料の全工程品質管理を実現。",
    "几何量测量：依托高精度三坐标测量机，检测工件三维尺寸、形位公差，测量精度高、适配复杂结构件。":
    "寸法測定：高精度三次元測定機により三次元寸法・幾何公差を測定。高精度で複雑形状部品に対応。",
    "性能与表面检测：可完成粗糙度、涂层厚度、硬度等指标测试，分析材料及表面特性。":
    "性能・表面試験：粗さ・膜厚・硬さ等を測定し、材料及び表面特性を分析。",
    "无损 & 可靠性测试：检测材料内部缺陷，开展环境适应性试验，验证产品使用可靠性。":
    "非破壊・信頼性試験：材料内部欠陥検出、環境試験により製品信頼性を検証。",
    "实验室实现微观表面、宏观尺寸、材料性能、环境耐受全场景覆盖，支持产品质量验证、工艺迭代与合规评价。":
    "ミクロ表面からマクロ寸法、材料性能、環境耐性まで全領域カバー。品質検証・工程改善・適合性評価を支援。",

    "黄 总": "黄 総経理",
    "重庆市铜梁区蒲吕街道产业大道55号附39号": "中国重慶市銅梁区蒲呂街道産業大道55号附39号",

    "晨启新局，兵拓未来，晨兵智造，精工志远":
    "晨に新局を開き、兵をもって未来を拓く——晨兵精工製造、卓越への挑戦",
}


def collect_all_runs(shape):
    """Collect all runs from a shape's text frame, preserving paragraph structure."""
    result = []  # list of (run, paragraph_index)
    for pi, para in enumerate(shape.text_frame.paragraphs):
        for run in para.runs:
            result.append((run, pi))
    return result


def apply_translations(shape, trans_map):
    """Join all runs, apply all translations, write back to runs."""
    runs = collect_all_runs(shape)
    if not runs:
        return

    # Build full text
    full = "".join(r.text for r, _ in runs)

    # Apply each translation (sort by key length descending to avoid partial matches)
    keys_sorted = sorted(trans_map.keys(), key=len, reverse=True)
    for key in keys_sorted:
        if key in full:
            full = full.replace(key, trans_map[key])

    if not runs:
        return

    # Write back: first run gets all text, others get ""
    runs[0][0].text = full
    for r, _ in runs[1:]:
        r.text = ""

    # Auto-fit: shrink text to fit within original text box bounds
    from pptx.util import Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    try:
        tf = shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
    except Exception:
        pass


def translate_table_cells(shape, trans_map):
    """Translate text in table cells."""
    from pptx.enum.text import MSO_AUTO_SIZE
    for row in shape.table.rows:
        for cell in row.cells:
            runs = []
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    runs.append(run)
            if not runs:
                continue
            full = "".join(r.text for r in runs)
            # Apply translations
            keys_sorted = sorted(trans_map.keys(), key=len, reverse=True)
            for key in keys_sorted:
                if key in full:
                    full = full.replace(key, trans_map[key])
            runs[0].text = full
            for r in runs[1:]:
                r.text = ""
            # Auto-fit text in cell
            try:
                cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                cell.text_frame.word_wrap = True
            except Exception:
                pass


def shrink_font_in_group(shape, min_pt=8, reduce_by=2):
    """Reduce font size for text inside group shapes to prevent row overlap."""
    from pptx.util import Pt
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > Pt(min_pt + reduce_by):
                    new_size = max(Pt(min_pt), run.font.size - Pt(reduce_by))
                    run.font.size = new_size
    if shape.shape_type == 6:
        for child in shape.shapes:
            shrink_font_in_group(child, min_pt, reduce_by)


def process_shape(shape, trans_map, in_group=False):
    """Process a single shape or recurse into groups."""
    if shape.has_table:
        translate_table_cells(shape, trans_map)
    elif shape.has_text_frame:
        apply_translations(shape, trans_map)
        if in_group:
            # Reduce font size to prevent text overlap between rows in product tables
            from pptx.util import Pt
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size / 12700  # EMU to points
                        if pt >= 10:
                            run.font.size = Pt(max(8, pt - 2))
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            process_shape(child, trans_map, in_group=True)


def fix_overlap_slides(prs, slide_indices):
    """Fix text overlap on product-table slides.

    Text boxes are designed for single-line Chinese. English wraps to 2+ lines.
    Strategy: set readable 10pt font, tight line spacing, and increase text box
    heights modestly so wrapped text stays within each box.
    """
    from pptx.util import Pt, Emu

    for idx in slide_indices:
        slide = prs.slides[idx]

        for shape in slide.shapes:
            _format_for_readability(shape)

        # Also format and resize shapes inside groups
        for shape in slide.shapes:
            if shape.shape_type == 6:
                for child in shape.shapes:
                    _format_for_readability(child)

        # Increase height of data-row text boxes (below header at Y≈1.6")
        for shape in slide.shapes:
            _resize_data_rows(shape)


def _resize_data_rows(shape):
    """Increase height of data-row text boxes to fit 2 lines of English text.
    Also remove fill so overlapping bounding boxes don't hide text below."""
    from pptx.util import Emu

    def resize_if_data(sh):
        if sh.has_text_frame and sh.text_frame.text.strip():
            top_inches = sh.top / 914400
            height_inches = sh.height / 914400
            if top_inches >= 2.0 and height_inches < 0.55 * 914400:
                sh.height = int(0.55 * 914400)
                # Remove fill so taller boxes don't cover text in rows below
                try:
                    sh.fill.background()
                except Exception:
                    pass

    if shape.shape_type == 6:
        for child in shape.shapes:
            resize_if_data(child)
            if child.shape_type == 6:
                _resize_data_rows(child)
    else:
        resize_if_data(shape)


def _format_for_readability(shape):
    """Set readable font size and tight line spacing on a shape."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            # Tight line spacing
            para.line_spacing = Pt(12)
            para.space_before = Pt(0)
            para.space_after = Pt(0)
            for run in para.runs:
                if run.font.size is None:
                    run.font.size = Pt(10)
                else:
                    current_pt = run.font.size / 12700
                    if current_pt > 12:
                        run.font.size = Pt(11)
                    elif current_pt > 10:
                        run.font.size = Pt(10)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    para.line_spacing = Pt(12)
                    para.space_before = Pt(0)
                    para.space_after = Pt(0)
                    for run in para.runs:
                        if run.font.size is None:
                            run.font.size = Pt(10)

    if shape.shape_type == 6:
        for child in shape.shapes:
            _format_for_readability(child)


def translate_pptx(src_path, out_path, trans_map, en_fix_slides=None):
    prs = Presentation(src_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, trans_map)
    if en_fix_slides:
        fix_overlap_slides(prs, en_fix_slides)
    prs.save(out_path)
    print(f"  Saved: {out_path}")


if __name__ == "__main__":
    # Slide indices 4,6,8,10,12 = slides 5,7,9,11,13 (product comparison tables)
    en_overlap_slides = [4, 6, 8, 10, 12]
    print("Generating English version...")
    translate_pptx(SRC, OUT_DIR + "Chongqing_Chenbing_Machinery_EN.pptx", EN_TRANSLATIONS,
                   en_fix_slides=en_overlap_slides)

    print("Generating Japanese version...")
    translate_pptx(SRC, OUT_DIR + "重慶晨兵機械有限公司_JA.pptx", JA_TRANSLATIONS)

    print("Done!")
